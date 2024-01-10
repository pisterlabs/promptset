import numpy as np
import pandas as pd

import openai
from pptx import Presentation

# Replace 'your_api_key_here' with your actual OpenAI API key
openai.api_key = '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'

from openai import OpenAI
client = OpenAI(api_key = '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')
# Function to reword text using GPT-4
#def reword_text_with_gpt4(text):
#    try:
#        # Use the OpenAI API to get a response
#        response = openai.chat.completions.create(
#            model="gpt-3.5-turbo-1106",  # Replace with the appropriate GPT-4 model when available
#            messages=[
#          {"role": "system", "content": "Reword the following text to be clear and concise:\n\n" + text}
#                ]
#            #prompt="Reword the following text to be clear and concise:\n\n" + text,
#            #max_tokens=60  # Adjust max tokens as needed
#        )
#        return response.choices[0].text.strip()
#    except Exception as e:  # Catch a general exception
#        print(f"An error occurred: {e}")
#        return text  # Return the original text if an error occurs

def reword_text_with_gpt4(text, audience_type, shape_type):
    try:
        if shape_type =="Title 1":
            title_prompt = " This is a title of a slide so keep it to less than 8 words."
        else:
            title_prompt =""
        # Use the OpenAI API to get a response
        response = openai.chat.completions.create(
            model="gpt-4-0613",  # Replace with the appropriate model
            messages=[
                {
                    "role": "system", 
                    "content": (
                        "You are an assistant that rewords sentences to be clear and concise. Your output will be no longer than the input in length." +
                        " I am presenting to " + audience_type + ", so make it suitable for this audience." + title_prompt

                                )
                    },
                {"role": "user", "content": text}
            ]
        )
        # Assuming the last message in the list will be the assistant's response
        return str(response.choices[0].message.content)#.text.strip()
    except Exception as e:  # Catch a general exception
        print(f"An error occurred: {e}")
        return text  # Return the original text if an error occurs

def create_exec_summary(text, audience_type):
    try:
        # Use the OpenAI API to get a response
        response = openai.chat.completions.create(
            model="gpt-4-0613",  # Replace with the appropriate model
            messages=[
                {
                    "role": "system", 
                    "content": (
                        "You are an assistant that creates executive summaries." +
                        " I am presenting to " + audience_type + ", so make it suitable for this audience. Produce a simple five bullet point summary"

                                )
                    },
                {"role": "user", "content": text}
            ]
        )
        # Assuming the last message in the list will be the assistant's response
        return str(response.choices[0].message.content)#.text.strip()
    except Exception as e:  # Catch a general exception
        print(f"An error occurred: {e}")
        return text  # Return the original text if an error occurs
# Function to process the PowerPoint file
def process_presentation(
    input_file_path, 
    output_file_path, 
    audience_type,
    executive_summary_slide = False):
    # Load the presentation
    prs = Presentation(input_file_path)
    
    # Iterate through each slide and each text box
    all_text = ""
    for slide_number, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                original_text = shape.text
                # Reword the text using GPT-4
                reworded_text = reword_text_with_gpt4(original_text, audience_type, shape.name)

                # Replace the original text with the reworded text
                shape.text = reworded_text
                all_text = all_text + ". " + original_text
        print(f"Processed slide {slide_number + 1}")
    
    if executive_summary_slide == True:
        slide_layout = prs.slide_layouts[1]
        slide_exec = prs.slides.add_slide(slide_layout)
        slide_exec.placeholders[0].text = "Executive Summary"
        slide_exec.placeholders[1].text = create_exec_summary(all_text, audience_type)
    # Save the presentation
    prs.save(output_file_path)
    print(f"Presentation saved to {output_file_path}")




# Example usage
input_file_path = "C:/Users/Administrator/Documents/GPTB4.pptx"
output_file_path = "C:/Users/Administrator/Documents/GPTA4"
output_file_path_technical = "C:/Users/Administrator/Documents/GPTA4_tech.pptx"
output_file_path_baby = "C:/Users/Administrator/Documents/GPTA4_babies.pptx"
audience_type1 = "a technical audience"
audience_type2 = "a bunch of five year olds who like thomas the tank engine"
process_presentation(input_file_path, output_file_path_technical, audience_type1,executive_summary_slide = True)
#process_presentation(input_file_path, output_file_path_baby, audience_type2)
