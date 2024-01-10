from pptx import Presentation
import openai
from dotenv import load_dotenv

import os
dotenv_path = 'key.env'
load_dotenv(dotenv_path)
print("OpenAI API Key:")
print(f"API Key: {os.getenv('OPENAI_API_KEY')}")
openai.api_key = os.getenv('OPENAI_API_KEY')



def check_grammar_and_spelling(text):
    print(f"Checking grammar and spelling of text: {text}")
    try:
        response = openai.Completion.create(
            model="text-davinci-003",  # This is a model suitable for the completions endpoint.
            prompt=f"Korriegere diesen Text auf Rechtschreib und Grammatikfehler und geben ihn ohne Veränderung der Formatierung zurück:\n\n{text}\n\n",
            max_tokens=60,
            temperature=0.5,
            top_p=1,
            n=1,
            #stop=["\n"],  # This will stop the generation at the first line break after the completion.
            frequency_penalty=0,
            presence_penalty=0
        )
        print(f"OpenAI response: {response}")
        corrected_text = response.choices[0].text.strip()
        print(f"Corrected text: {corrected_text}")
        return corrected_text
    except Exception as e:
        print(f"An error occurred while querying OpenAI: {e}")
        return None

def extract_text_from_pptx(filepath):
    print(f"Extracting text from presentation: {filepath}")
    prs = Presentation(filepath)
    slides_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                print("Skipping shape without text frame")
                continue
            print(f"Shape: {shape}")
              # Nimmt an, dass die erste Shape den Titel enthält
            text = shape.text
            print(f"Title: {text}")
            slides_text.append(text)
            break
    return slides_text

def analyze_presentation(filepath):
    print(f"Analyzing presentation: {filepath}")
    slides_text = extract_text_from_pptx(filepath)
    corrected_texts = []
    for text in slides_text:
        corrected_text = check_grammar_and_spelling(text)
        corrected_texts.append(corrected_text)
    
    return corrected_texts

# Diese Funktion kann in deinem Flask-Handler aufgerufen werden
def perform_analysis(filepath):
    corrected_texts = analyze_presentation(filepath)
    feedback = [f'Slide {i+1}: {text}' for i, text in enumerate(corrected_texts) if text]
    analysis_result = {
        'slide_count': len(corrected_texts),
        'feedback': feedback
    }
    return analysis_result
