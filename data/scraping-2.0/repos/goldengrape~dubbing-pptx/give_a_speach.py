import pptx
import markdown
from typing import List
import time
import openai 
import os


sleep_time=60 
def query_gpt3(prompt,cooldown_time=3):
    global sleep_time
    while True:
        try:
            # start_time=time.time()
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo", 
                messages=[{
                "role": "user", 
                "content": prompt}]
                )
            # print(f"GPT-3 API time: {time.time()-start_time}")
            answer=response.choices[0].message.content.strip()
            time.sleep(cooldown_time)
            # print(f"after sleep 3s, I finished")
            sleep_time = int(sleep_time/2)
            sleep_time = max(sleep_time, 10)
            break
        except:
            print(f"API error, retrying in {sleep_time} seconds...")
            time.sleep(sleep_time)
            sleep_time += 10
            if sleep_time > 120:
                print("API error, aborting...")
                answer=""
                break
    # print(f"Answer: {answer}")
    return answer

def read_pptx(file_path: str) -> List[str]:
    """
    读取pptx文件并返回包含每一页的文本内容的列表
    """
    prs = pptx.Presentation(file_path)
    text_list = []
    for slide in prs.slides:
        text = extract_text_and_notes(slide)
        text_list.append(text)
    return text_list

def extract_text_and_notes(slide: pptx.slide.Slide) -> str:
    """
    从给定的PPT幻灯片中提取正文和备注中的文字
    """
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text
        if hasattr(shape, "notes_text_frame"):
            text += shape.notes_text_frame.text
    return text

def generate_speech(slide_text: str, all_text: str, speech_duration: int, words_per_minute: int = 130) -> str:
    """
    根据给定的幻灯片文本和整个PPT的文本，以及用户输入的每一页期望的演讲时间，使用GPT生成相应字数的演讲词
    """
    
    tone='积极，风趣，幽默'
    content='严谨，科学，有逻辑'

    words_per_slide = speech_duration * words_per_minute // 60
    # TODO: 使用GPT生成演讲词
    prompt=f'''
        请根据整个ppt的内容和当前页面的内容，产生当前页面的演讲词。
        其中当前页面的内容为
        【
        {slide_text}
        】
        整个ppt的内容为
        【
        {all_text}
        】
        注意：
        1. 演讲词大约{words_per_slide}字
        2. 演讲词的语气应当{tone}
        3. 演讲词的内容应当{content}
        '''
    speech=query_gpt3(prompt)
    return speech

def save_to_markdown(speech_list: List[str], file_path: str) -> None:
    """
    将生成的演讲词列表保存为Markdown文件
    """
    with open(file_path, "w", encoding="utf-8") as f:
        for speech in speech_list:
            f.write(markdown.markdown(speech))
            f.write("\n")

def main() -> None:
    """
    主函数，调用其他函数完成任务
    """
    file_path = "test/test.pptx"
    text_list = read_pptx(file_path)
    all_text = "".join(text_list)
    speech_list = []
    for i, slide_text in enumerate(text_list):
        print(f"-----\nGenerating speech for slide {i + 1}...")
        speech = generate_speech(slide_text, all_text, 60)
        print(f"Speech for slide {i + 1} generated: \n{speech}")

        speech_list.append(speech)  # Add the generated speech to the speech_list

    save_to_markdown(speech_list, "output.md")  # Save the speech_list to a markdown file

if __name__ == "__main__":
    openai.api_key=os.environ['OPENAI_API_KEY']
    main()  # Call the main function