import openai, json
from pptx import Presentation
import apikey

openai.api_key = apikey.key   #apikey 


presentation_title = input('Write title of your presentation?? ')

query_json = """
{
    "input_text": "[[QUERY]]",
    "output_format": "json",
    "json_structure": {
        "slides": "{{presentation_slides}}"
    }
}
"""

question = "Generate 10 slide presention on " + presentation_title + " Each slide should have a {{headar}}, {{content}}. Return as JSON."
prompt = query_json.replace("[[QUERY]]",question)
print(prompt)

completion = openai.ChatCompletion.create(model= "gpt-3.5-turbo", messages = [{'role':'user','content':prompt}])

response = completion.choices[0].message.content
print(response)

r = json.loads(response)


slides_data = r["slides"]

presentation = Presentation()


for slide in slides_data:
    slide_latous = presentation.slide_layouts[1] #0 for title slide
    new_slide = presentation.slides.add_slide(slides_data)


    if slide['header']:
        title = new_slide.shapes.title
        title.text = slide["header"]
    if slide['content']:
        content = new_slide.shapes.placeholders[1]
        tf = content.text_frame
        tf.text = slide["content"]




presentation.save('generated.ppt')


