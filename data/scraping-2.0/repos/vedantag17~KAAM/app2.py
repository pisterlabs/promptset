from dotenv import load_dotenv
import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
import os
import pytesseract
import PyPDF2
from PIL import Image
import textract
from docx import Document
from pptx import Presentation
import pandas as pd
import tempfile

def ocr_image(image):
    text = pytesseract.image_to_string(image)
    return text

def read_text_from_docx(docx_file):
    doc = Document(docx_file)
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text

def read_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text + " "
    return text

def read_text_from_pdf(pdf_file):
    with open(pdf_file, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        text = "\n".join([page.extract_text() for page in reader.pages])
    return text

def read_data_from_spreadsheet(spreadsheet_file):
    _, file_extension = os.path.splitext(spreadsheet_file.name)
    if file_extension == ".xlsx":
        df = pd.read_excel(spreadsheet_file)
    else:
        df = pd.read_csv(spreadsheet_file)
    return df

def save_uploaded_file(uploaded_file):
    _, file_extension = os.path.splitext(uploaded_file.name)
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=file_extension)
    temp_file.write(uploaded_file.read())
    temp_file.close()
    return temp_file.name

def main():
    load_dotenv()
    st.set_page_config(page_title="KAAM")
    st.header("KKW Buddy!")

    uploaded_files = st.file_uploader("Upload Files", accept_multiple_files=True)

    if uploaded_files:
        user_question = st.text_input("Ask a question about the files:")
        
        if user_question:
            for file in uploaded_files:
                file_extension = os.path.splitext(file.name)[1].lower()

                if file_extension == ".pdf":
                    temp_file_path = save_uploaded_file(file)
                    text = read_text_from_pdf(temp_file_path)
                    os.remove(temp_file_path)

                    # Question-answering on PDF content
                    text_splitter = CharacterTextSplitter(
                        separator="\n",
                        chunk_size=1000,
                        chunk_overlap=200,
                        length_function=len
                    )
                    chunks = text_splitter.split_text(text)

                    embeddings = OpenAIEmbeddings()
                    knowledge_base = FAISS.from_texts(chunks, embeddings)

                    docs = knowledge_base.similarity_search(user_question)

                    llm = OpenAI()
                    chain = load_qa_chain(llm, chain_type="stuff")
                    response = chain.run(input_documents=docs, question=user_question)

                    st.write(f"File: {file.name}")
                    st.write(response)
                    st.write("---")

if __name__ == "__main__":
    main()

