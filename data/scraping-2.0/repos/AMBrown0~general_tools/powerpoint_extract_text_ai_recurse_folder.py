# -*- coding: utf-8 -*-
"""
Created on Thu Feb 23 17:59:29 2023

@author: andy
"""
# =============================================================================
# 
# =============================================================================
import tkinter as tk
import tkinter.filedialog as filedialog
import os
import os

import tkinter as tk
from tkinter import filedialog
import os
import pptx
from datetime import datetime
import openai
import re
import nltk
from nltk.corpus import stopwords
#nltk.download('stopwords')




# Function to select directory using tkinter dialog box
def select_directory():
    root = tk.Tk()
    root.withdraw()
    return filedialog.askdirectory()

# Function to recursively find all files with .pptx extension in the given directory
def find_pptx_files(directory):
    pptx_files = []
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.endswith(".pptx"):
                # Check that this is not a sync file (starts with ~$) and ignore this
                if not file.startswith(r"~$"):
                    pptx_files.append(os.path.join(root, file))
    return pptx_files



def call_openapi_recursive(question,action_text):
    
    #call openai
    response=call_openapi(question,action_text)
    if response == False:
        cut=int(len(question)/2)        
        h1=question[:cut]
        h2=question[cut:]
        h1_fulltext=f"{action_text}:{h1}"
        #print(f"H1 h1_fulltext:{h1_fulltext}")
        h2_fulltext=f"{action_text}:{h2}"
        #print(f"H2 h2_fulltext:{h2_fulltext}")
        r1=call_openapi_recursive(h1_fulltext,action_text)
        #print(f"r1={r1}")
        #print("======================================================")
        
        r2=call_openapi_recursive(h2_fulltext,action_text)
        #print(f"r2={r2}")
        #print("======================================================")
        complete_response = r1 + "\n\n" + r2
            
        #print(h2)
    else:
        complete_response=response
        
    #print(f"Reponse:{complete_response}")
    return complete_response
        
            


def call_openapi(question,action_text):
    full_text=f"{action_text}:{question}"
    try:
        openai.api_key =os.environ['CHATGPT_KEY']
        response = openai.Completion.create(
            engine="text-davinci-002",
            prompt=full_text,
            max_tokens=4000,
            temperature=0.5,
            top_p=1,
            frequency_penalty=0,
            presence_penalty=0,
            timeout=None,
        )
    except:
        return False
    response_text=response.choices[0].text
    #response_choices=len(response.choices)
    #print(f"responses={response_choices}")
    return response_text



def extract_powepoint_text(powerpoint_list):
    #headings = []
    #body_text = []
    slide_list=[]
    
    for filename in powerpoint_list:
        ppt = pptx.Presentation(filename)
        for slide in ppt.slides:
            
            # Extract the heading (assuming it's the first shape on the slide)
            
            try:
                heading=slide.shapes.title.text
            except:
                heading=""
            # try:
            #     if hasattr(slide.shapes[0],"text"):
            #         heading = slide.shapes[0].text_frame.text.strip()
            #         headings.append(heading)
            #     else:
            #             headings.append("")                        
            # except:
            #     headings.append("")
            
            # Add all text
            text = ""
            
            number_shapes_on_slide=len(slide.shapes)
            for i in range(0,number_shapes_on_slide):
                shape=slide.shapes[i]
                if hasattr(shape, "text"):
                    #print(slide.shapes[0].shape_type)
                    text += shape.text
                    #body_text.append(text)
            print(f"body_text:{text}\n\n")
            slide_dict={"file":filename,"heading":heading,"body_text":text}
            slide_list.append(slide_dict)
    return slide_list
    

    # for ppt in powerpoint_list:
    #     print(f"extracting text from:{ppt}")

    #     ppt = pptx.Presentation(ppt)
    #     for slide in ppt.slides:
    #         # Extract the heading (assuming it's the first shape on the slide)
    #         heading = slide.shapes[0].text_frame.text.strip()
    #         headings.append(heading)
        
    #     # Extract the body text (assuming it's the remaining shapes on the slide)
    #         text = ""
    #         for shape in slide.shapes[1:]:
    #             if hasattr(shape, "text"):
    #                 text += shape.text
    #                 body_text.append(text)
    #         slide_dict={"file":ppt,"heading":heading,"body_text":body_text}
    #         slide_list.append(slide_dict)
    # for slide in slide_list:
    #     print(f"file:{slide.get('file')} heading:{slide.get('body_text')}")
        
    
    
def clean_powerpoint(file_path):

    # Open the PowerPoint file
    prs = pptx.Presentation(file_path)
    
    
    for master in prs.slide_masters:
        # Loop through all the background objects on the master
        obj=master.background
        obj.fill.background()
        shapes=master.shapes
        for shape in shapes:
            shape_type=shape.shape_type
            
            if shape.shape_type is not  pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
                shape.element.getparent().remove(shape.element)
                print(f"deleting {shape_type}")
        # Delete the object
        #obj.element.getparent().remove(obj.element)
        
    
    # Loop through all the slides in the presentation
    for slide in prs.slides:
        # Get the master slide of the slide
        master = slide.slide_layout
      
        # Loop through all the shapes on the master slide
        for shape in master.shapes:
            
            
            #print(f"Slide[{master_count}]Shape_Type={shape_type}")
            
            if shape.shape_type is not  pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
                shape_type=shape.shape_type
                shape.element.getparent().remove(shape.element)
                print(f"deleting {shape_type}")
    # Save the presentation
    prs.save(file_path)
# Get the directory from the user using tkinter dialog box
directory = select_directory()

# Find all .pptx files in the selected directory recursively
pptx_files = find_pptx_files(directory)

all_slides_text_list=extract_powepoint_text(pptx_files)


#Write to text file
text_file_name="slide_summary"
now = datetime.now()
date_time = now.strftime("%Y-%m-%d-%H-%M-%S")
text_file_name=f"{directory}/{text_file_name}-{date_time}.txt"
question_file_name=f"{text_file_name}q-{date_time}.txt"
last_filename="NONE"

# with open(text_file_name, "w") as file:
#     for slide in all_slides_text_list:    
#         filename=slide.get("file")
#         if filename != last_filename:
#             file.write(f"\n\n================= {filename}==============\n")
#             last_filename=filename
#         heading=slide.get("heading")
#         text_body=slide.get("body_text")
#         text=f"{heading}\n{text_body}"
#         file.write(text)    
#         ai_text=f"Create a summary of these powerpoint bullet points:{heading}\n{text_body}"
#         r=call_openapi(ai_text)
#         file.write(r)    

# print(f"file={text_file_name}")
# # Print the file path of all .pptx files found
# 


# Ask AI to summarise each slide

responses=[]

with open(text_file_name, "w") as file:
    complete_slide_text=""
    for slide in all_slides_text_list:    
        
        
        filename=slide.get("file")
        if not filename == last_filename:
            # New slide deck
            slide_title_text=f"\n\n{filename}\n"
            file.write(slide_title_text)
            file.write(" =============================================================================\n\n")
            last_filename=filename
            
        # if filename != last_filename:
        #     file.write(f"\n\n================= {filename}==============\n")
        #     last_filename=filename
        heading=slide.get("heading")
        text_body=slide.get("body_text")
        #"detailed summary if thi
        ai_text=f"{heading}\n{text_body}"
        complete_slide_text=complete_slide_text+ai_text

        r=call_openapi_recursive(ai_text,"create a detailed summary of these powerpoint slide texts")
            
        responses.append(r)
        file.write(f"{heading}\n\n{r}\n\n")
        
     
        
        print(f"summary={r}")
        


# with open(question_file_name, "w") as file:
#     complete_slide_text=""
#     for response in responses:    
#         try:
#             ai_text=f"create 5 multiple choice quetions with answers :{response}"
#             r=call_openapi(ai_text)        
#             file.write(r)
#             print(f"multi-choice questions={r}")
            
#             ai_text=f"create 5 open quetions with answers:{response}"
#             r=call_openapi(ai_text)        
#             file.write(r)
#             print(f"open questions={r}")
#         except:
#                 print(f"too much information to respond this {response}")

# # Load the stop words for English language
# stop_words = set(stopwords.words('english'))

# # Define a text
# text = complete_slide_text

# # Tokenize the text into words
# words = nltk.word_tokenize(text)

# # Remove the stop words
# filtered_words = [word for word in words if word.lower() not in stop_words]

# print(filtered_words)
# recreated_text_no_stops=""
# #recreate text without stops
# for word in filtered_words:
#     recreated_text_no_stops=recreated_text_no_stops + f" {word}"
    


