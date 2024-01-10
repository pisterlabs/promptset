import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openai
import config
from PIL import Image, ImageEnhance
import requests
from io import BytesIO
import PyPDF2 
import re
import tiktoken
import numpy as np
from PIL import Image, ImageDraw
import os
import random
import cv2
import time
from googletrans import Translator
from pptxWithBoth import generate3PicturesSlide, generate2PicturesSlide, generate1PicturesSlide, firstSlide, generatePlan
from functions import question



# openai.api_key = "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"
prs = Presentation()
prs.slide_height = Pt(1080)
prs.slide_width = Pt(1920)

translator = Translator()

def GeneratePresentationWith2TypesOfPictures(topic, numOfSlides, font, usersname, plan=True, interestFact=True, timecode=1):
    print("start presentation")
    themes = []
    finalThemesFromPrompt = []
    language = translator.detect(topic).lang
    # if language != "en":
    #     themesFromPrompt = question(f"write to me like above a {numOfSlides-2} topics for {numOfSlides-2} slides about {translator.translate(topic).text}, without introducing and conclusion slide with 2-5 words. say about every slide what is better search pictures in network or generate with dalle. after topic write only 1 word: 'dalle' and write better prompt for falle for this picture devided by '-' if dalle is better. Write 'network' and write better querry for this picture devided by '-' if network is better to find picture for this slide. say about every slide how much picture is needed for this slide with only one integer digit devided by ' - ' in interval from 0 to 3. Then for every slide generate text for this slide about particular topic in 50 words devided by ' - '. ").replace('\n', ' - ').replace('"', '').split(' - ')
    # else:
    #     themesFromPrompt = question(f"write to me like above a {numOfSlides-2} topics for {numOfSlides-2} slides about {topic}, without introducing and conclusion slide with 2-5 words. say about every slide what is better search pictures in network or generate with dalle. after topic write only 1 word: 'dalle' and write better prompt for falle for this picture devided by '-' if dalle is better. Write 'network' and write better querry for this picture devided by '-' if network is better to find picture for this slide. say about every slide how much picture is needed for this slide with only one integer digit devided by ' - ' in interval from 0 to 3. Then for every slide generate text for this slide about particular topic in 50 words devided by ' - '. ").replace('\n', ' - ').replace('"', '').split(' - ')
    # for i in range(len(themesFromPrompt)):
    #     if themesFromPrompt[i] != '':
    #         finalThemesFromPrompt.append(themesFromPrompt[i])
    # # finalThemesFromPrompt = ['Colors', 'network', 'Rainbow colors', '2', 'An image showing the colors of the rainbow in order, from red to violet', 'The colors of the rainbow are a beautiful natural phenomenon. This image displays the colors in order, from red to violet, and can be used to teach children about the colors of the rainbow. ', 'Formation', 'dalle', 'Rainbow formation', '1', 'A picture of a rainbow forming in the sky, with clear distinction of the colors', 'Rainbows are formed when sunlight is refracted and reflected by water droplets in the air. This picture captures the moment of a rainbow forming in the sky, with clear distinction of the colors. ', 'Mythology', 'network', 'Rainbow mythology', '3', 'An image of a painting or artwork depicting a rainbow in mythology, such as the Norse Bifrost or the Greek Iris', 'Rainbows have been a part of mythology and folklore for centuries. This image shows a painting or artwork depicting a rainbow in mythology, such as the Norse Bifrost or the Greek Iris. ', 'Double Rainbow', 'dalle', 'Double rainbow', '2', 'A picture of a double rainbow, with both rainbows clearly visible', 'Double rainbows are a rare and beautiful sight. This picture captures the moment of a double rainbow, with both rainbows clearly visible and the colors vibrant and distinct. ', 'Scientific Explanation', 'network', 'Rainbow diagram', '1', 'A diagram or illustration explaining the scientific process behind the formation of a rainbow', 'Rainbows are a result of the refraction and reflection of light through water droplets in the air. This diagram or illustration explains the scientific process behind the formation of a rainbow, including the angles of refraction and reflection.']
    # for i in range(0, len(finalThemesFromPrompt), 6):
    #     themes.append(finalThemesFromPrompt[i:i+6])
    # print(themes)
    # [['1. Formation', 'dalle', '"coal formation"', '3', 'Coal is formed from the remains of ancient plants that were buried and subjected to high pressure and heat over millions of years. The process of coal formation is called coalification.'], ['2. Types', 'network', '"types of coal"', '4', 'There are four main types of coal: anthracite, bituminous, subbituminous, and lignite. Each type has different properties and uses, with anthracite being the highest quality and lignite being the lowest.'], ['3. Mining', 'dalle', '"coal mining"', '2', "Coal mining involves extracting coal from the earth's surface or underground. It can be a dangerous and environmentally damaging process, but it is also a vital source of energy for many countries."], ['4. Uses', 'network', '"coal uses"', '3', 'Coal is primarily used for electricity generation, but it is also used in steel production, cement manufacturing, and other industrial processes. It has been a key source of energy for centuries.'], ['5. Environmental impact', 'dalle', '"coal and the environment"', '4', 'Coal mining and burning have significant environmental impacts, including air and water pollution, habitat destruction, and greenhouse gas emissions. Efforts are being made to reduce these impacts and transition to cleaner energy sources.']]
    # themes = ['1. Formation of Coal', 'dalle', '2. Types of Coal', 'network', '3. Uses of Coal', 'dalle', '4. Environmental Impact of Coal', 'network', '5. Future of Coal', 'dalle']
    themes = [['Anatomy', 'dalle', 'Frog anatomy diagram', '2', "A diagram of a frog's anatomy, highlighting its unique features such as webbed feet and long tongue", 'Frogs have unique adaptations that allow them to thrive in their aquatic and terrestrial habitats. This diagram showcases the anatomy of a frog, including its webbed feet, long tongue, and bulging eyes. '], ['Habitat', 'network', 'Frog habitat image', '1', 'An image of a frog in its natural habitat, such as a pond or wetland', 'Frogs are found in a variety of habitats, from rainforests to deserts. This image captures a frog in its natural habitat, showcasing the importance of wetlands and other aquatic ecosystems for frog survival. '], ['Life cycle', 'dalle', 'Frog life cycle illustration', '2', "An illustration of the different stages of a frog's life cycle, from egg to tadpole to adult frog", "Frogs undergo a unique metamorphosis, transforming from aquatic tadpoles to terrestrial adults. This illustration depicts the different stages of a frog's life cycle, highlighting the importance of wetland habitats for breeding and development. "], ['Diet', 'network', 'Frog eating insect', '1', 'An image of a frog catching and eating an insect, showcasing its role as a predator in the food chain', 'Frogs are important predators in their ecosystems, feeding on insects and other small animals. This image captures a frog in action, highlighting its unique feeding behavior and role in maintaining a healthy ecosystem. '], ['Conservation', 'dalle', 'Endangered frog species', '3', 'A collection of images showcasing endangered frog species and the threats they face, such as habitat loss and pollution', 'Frogs are facing numerous threats, including habitat loss, pollution, and disease. This collection of images highlights some of the endangered frog species and the urgent need for conservation efforts to protect these important amphibians.']]
    
    os.mkdir(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}")
    slides = []
    blank_slide_layout = prs.slide_layouts[6]

    # FIRST SLIDE
    slides.append(prs.slides.add_slide(blank_slide_layout)) #slides[0] 13.333 in x 7.5 in
    firstSlide(slide=slides[0], topic=topic, usersname=usersname)

    # PLAN
    if(plan == True):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        generatePlan(slide=slides[1], themes=themes, language=language)
    

    # Main Slides
    # slides.append(prs.slides.add_slide(blank_slide_layout))
    
    for i in range(len(themes)):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        if themes[i][3] == '3':
            generate3PicturesSlide(slide=slides[2 + i], methodOfPictures=themes[i][1], PromptForPictures=themes[i][2],
                                    textForSlide=translator.translate(themes[i][5], dest=language).text, topic=translator.translate(themes[i][0], dest=language).text, varitation=1, usersname=usersname, timecode=timecode, 
                                    directory = str(usersname).replace(' ', '') + str(timecode),
                                    prompt2=themes[i][4])
        elif themes[i][3] == '2':
            generate2PicturesSlide(slide=slides[2 + i], methodOfPictures=themes[i][1], PromptForPictures=themes[i][2],
                                    textForSlide=translator.translate(themes[i][5], dest=language).text, topic=translator.translate(themes[i][0], dest=language).text, varitation=1, usersname=usersname, timecode=timecode, 
                                    directory = str(usersname).replace(' ', '') + str(timecode),
                                    prompt2=themes[i][4])
        else:
            generate1PicturesSlide(slide=slides[2 + i], methodOfPictures=themes[i][1], PromptForPictures=themes[i][2],
                                    textForSlide=translator.translate(themes[i][5], dest=language).text, topic=translator.translate(themes[i][0], dest=language).text, varitation=1, usersname=usersname, timecode=timecode, 
                                    directory = str(usersname).replace(' ', '') + str(timecode),
                                    prompt2=themes[i][4])
        print(f"{i + 3}'s slide Done by {themes[i][1]}")

    nameOfPresentation = str(usersname+topic+str(timecode)).replace(' ', '')

    prs.save(f'{nameOfPresentation}.pptx')
    slides = []
    
    os.rmdir(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}")


GeneratePresentationWith2TypesOfPictures(topic="Лягушки", numOfSlides=7, font="DoraDura", usersname="Suvernev Bogdan", timecode=time.time())