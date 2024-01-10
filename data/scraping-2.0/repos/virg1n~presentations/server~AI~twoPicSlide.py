import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openai
import config
from PIL import Image, ImageEnhance
import requests
from io import BytesIO
import numpy as np
from PIL import Image, ImageDraw
import os
import random
import cv2
import time
from googletrans import Translator
from functions import image1024, image256, imgToCircle, forBeautifullText, saveImg1024, saveImg256, question, getImgFromGoogle, forBeautifullTextBySimbols, saveImg512, cropPicture, cropPicture8x9, saveImg512WithDirectory, generateImage

# With Image Generation
def twoPicSlideSDFirst(slide, topic, timecode, directory, PromptForPictures, prompt2, textForSlide):
    topicForName = topic.replace(' ', '').replace('.', '')
    generateImage(f"firstFor2Pictures{topicForName}{timecode}", dir=directory, prompt=PromptForPictures)
    generateImage(f"secondFor2Pictures{topicForName}{timecode}", dir=directory, prompt=prompt2)


    # saveImg512WithDirectory(PromptForPictures, f"firstFor2Pictures{topicForName}{timecode}", dir=directory)
    # saveImg512WithDirectory(prompt2, f"secondFor2Pictures{topicForName}{timecode}", dir=directory)
    txBox = slide.shapes.add_textbox(Pt(0), Pt(50), width = Pt(1920), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(5), Pt(250), width = Pt(1920), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=65)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)

    slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png', Pt(200), Pt(530), width=Pt(512), height=Pt(512))
    slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png', Pt(1008), Pt(530), width=Pt(512), height=Pt(512))


def twoPicSlideSDSecond(slide, topic, timecode, directory, PromptForPictures, prompt2, textForSlide):
    topicForName = topic.replace(' ', '').replace('.', '')
    generateImage(f"firstFor2Pictures{topicForName}{timecode}", dir=directory, prompt=PromptForPictures)
    generateImage(f"secondFor2Pictures{topicForName}{timecode}", dir=directory, prompt=prompt2)

    # saveImg512WithDirectory(PromptForPictures, f"firstFor2Pictures{topicForName}{timecode}", dir=directory)
    # saveImg512WithDirectory(prompt2, f"secondFor2Pictures{topicForName}{timecode}", dir=directory)

    txBox = slide.shapes.add_textbox(Pt(0), Pt(200), width = Pt(1290), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(10), Pt(400), width = Pt(1270), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=40)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(50)

    slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png', Pt(1290), Pt(50), width=Pt(475), height=Pt(475))
    slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png', Pt(1290), Pt(545), width=Pt(475), height=Pt(475))


def twoPicSlideSDThird(slide, topic, timecode, directory, PromptForPictures, prompt2, textForSlide):
    topicForName = topic.replace(' ', '').replace('.', '')
    generateImage(f"firstFor2Pictures{topicForName}{timecode}", dir=directory, prompt=PromptForPictures)
    generateImage(f"secondFor2Pictures{topicForName}{timecode}", dir=directory, prompt=prompt2)


    # saveImg512WithDirectory(PromptForPictures, f"firstFor2Pictures{topicForName}{timecode}", dir=directory)
    # saveImg512WithDirectory(prompt2, f"secondFor2Pictures{topicForName}{timecode}", dir=directory)
    txBox = slide.shapes.add_textbox(Pt(630), Pt(200), width = Pt(1290), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(640), Pt(400), width = Pt(1270), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=40)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(50)

    slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png', Pt(155), Pt(50), width=Pt(475), height=Pt(475))
    slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png', Pt(155), Pt(545), width=Pt(475), height=Pt(475))


# With Images from google
def twoPicSlideGoogleFirstThreeFirst(slide, topic, textForSlide, directory):
    txBox = slide.shapes.add_textbox(Pt(0), Pt(50), width = Pt(1920), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(5), Pt(250), width = Pt(1920), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=64)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg', Pt(200), Pt(530), width=Pt(512), height=Pt(512))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png', Pt(200), Pt(530), width=Pt(512), height=Pt(512))

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.jpg', Pt(1008), Pt(530), width=Pt(512), height=Pt(512))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.png', Pt(1008), Pt(530), width=Pt(512), height=Pt(512))


def twoPicSlideGoogleFirstThreeSecond(slide, topic, textForSlide, directory):
    txBox = slide.shapes.add_textbox(Pt(0), Pt(200), width = Pt(1290), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(10), Pt(400), width = Pt(1270), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=40)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(50)

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg', Pt(1290), Pt(50), width=Pt(475), height=Pt(475))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png', Pt(1290), Pt(50), width=Pt(475), height=Pt(475))

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.jpg', Pt(1290), Pt(545), width=Pt(475), height=Pt(475))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.png', Pt(1290), Pt(545), width=Pt(475), height=Pt(475))


def twoPicSlideGoogleFirstThreeThird(slide, topic, textForSlide, directory):
    txBox = slide.shapes.add_textbox(Pt(630), Pt(200), width = Pt(1290), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(640), Pt(400), width = Pt(1270), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=40)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(50)

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg', Pt(155), Pt(50), width=Pt(475), height=Pt(475))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png', Pt(155), Pt(50), width=Pt(475), height=Pt(475))

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.jpg', Pt(155), Pt(545), width=Pt(475), height=Pt(475))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.png', Pt(155), Pt(545), width=Pt(475), height=Pt(475))


def twoPicSlideGoogleSecondOneFirst(slide, topic, textForSlide, directory):
    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg', Pt(0), Pt(0), width=Pt(960), height=Pt(1080))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png', Pt(0), Pt(0), width=Pt(960), height=Pt(1080))

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.jpg', Pt(960), Pt(0), width=Pt(960), height=Pt(1080))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.png', Pt(960), Pt(0), width=Pt(960), height=Pt(1080))
    rectangle = slide.shapes.add_picture("C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\AI\\rectangle2.png", Pt(453), Pt(310), width=Pt(1066), height=Pt(450))
    txBox = slide.shapes.add_textbox(Pt(453), Pt(314), width = Pt(1066), height = Pt(450))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic[2:]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(65)

    txBox1 = slide.shapes.add_textbox(Pt(453), Pt(390), width = Pt(1066), height = Pt(450))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=42)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)


