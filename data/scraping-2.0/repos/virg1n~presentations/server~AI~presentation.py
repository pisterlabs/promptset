import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openai
import config
from PIL import Image, ImageEnhance
import requests
from io import BytesIO
import PyPDF2 
import re
import tiktoken
import numpy as np
from PIL import Image, ImageDraw
import os


# openai.api_key = "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"
openai.api_key = "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"
prs = Presentation()

def imgToCircle(path):
    img=Image.open(path).convert("RGB")
    npImage=np.array(img)
    h,w=img.size

    # Create same size alpha layer with circle
    alpha = Image.new('L', img.size,0)
    draw = ImageDraw.Draw(alpha)
    draw.pieslice([0,0,h,w],0,360,fill=255)

    # Convert alpha Image to numpy array
    npAlpha=np.array(alpha)

    # Add alpha layer to RGB
    npImage=np.dstack((npImage,npAlpha))

    # Save with alpha
    Image.fromarray(npImage).save(path)

def image(body):
    response = openai.Image.create(
        prompt=body,
        n=1,
        size="1024x1024"
    )
    image_url = response['data'][0]['url']
    return image_url

def image256(body):
    response = openai.Image.create(
        prompt=body,
        n=1,
        size="256x256"
    )
    image_url = response['data'][0]['url']
    return image_url

def question(body):
    response = openai.ChatCompletion.create(
        model='gpt-3.5-turbo',
        messages=[
            {'role': 'user', 'content': body}
        ],
        temperature=0
    )
    return response['choices'][0]['message']['content']

def forBeautifullText(string):
    i = 0
    newString = str()
    while i < len(string):
        newString += string[i:i+35] + ' \\n '
        i += 35
    return newString

def saveImg(prompt, name):
    url = image(prompt)
    response = requests.get(url)
    img = Image.open(BytesIO(response.content))
    img.save(f'./Ai/{name}.png')    


def saveImg256(prompt, name):
    url = image256(prompt)
    response = requests.get(url)
    img = Image.open(BytesIO(response.content))
    img.save(f'./Ai/{name}.png')    


def firstSlide(title, subtext, font, timecode):
    saveImg256(prompt=title + "Darkened Image", name=f"firstPage{timecode}")
    img_path = f'./Ai/firstPage{timecode}.png'
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(0) 
    width = height = Inches(10)
    pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)
    circle = slide.shapes.add_picture("./AI/circle.png", Inches(-3.1), Inches(-6.2), width=width, height=height)
    circle2 = slide.shapes.add_picture("./AI/circle.png", Inches(7.2), Inches(6.1), width=width/2, height=height/2)

    left = Inches(0.4)
    top = Inches(0.4)
    width = Inches(3)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.name = 'Comic Sans' #Arial Rounded MT Bold
    p.font.size = Pt(40)
    p = tf.add_paragraph()
    p.text = subtext
    p.font.size = Pt(15)
    p.font.name = 'Comic Sans'
    p.font.color.rgb = RGBColor(0, 0, 0)

    print(1)


#   Second Slide
    slide2 = prs.slides.add_slide(blank_slide_layout)
    # topics = question(f"Give me a 3 topics with 1 word maximum in each for slide with plan about {title}").replace('\n', '').replace('2', '').replace('3', '').split('.')[1:]
    
    left = top = Inches(0) 
    width = height = Inches(10)
    left = Inches(0)
    top = Inches(0.4)
    width = Inches(10)
    height = Inches(1)
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Our Plan"
    p.font.name = 'Comic Sans' #Impact
    p.font.size = Pt(40)
    p.alignment = PP_ALIGN.CENTER
    arrayOfPrompts = question(f"Give me a 3 topics with 1 word maximum in each for slide with plan about {title}").replace('\n', '').replace('2', '').replace('3', '').split('.')[1:]
    saveImg256(arrayOfPrompts[0] + title, f"lolimg0{timecode}")
    saveImg256(arrayOfPrompts[1] + title, f"lolimg1{timecode}")
    saveImg256(arrayOfPrompts[2] + title, f"lolimg2{timecode}")
    imgToCircle(f'./AI/lolimg0{timecode}.png')
    imgToCircle(f'./AI/lolimg1{timecode}.png')
    imgToCircle(f'./AI/lolimg2{timecode}.png')
    pic1 = slide2.shapes.add_picture(f"./AI/lolimg0{timecode}.png", Inches(0.5), Inches(2.1), width=Inches(2.1), height=Inches(2.1))
    pic2 = slide2.shapes.add_picture(f"./AI/lolimg1{timecode}.png", Inches(4), Inches(2.1), width=Inches(2.1), height=Inches(2.1))
    pic3 = slide2.shapes.add_picture(f"./AI/lolimg2{timecode}.png", Inches(7.4), Inches(2.1), width=Inches(2.1), height=Inches(2.1))
    txBox1 = slide2.shapes.add_textbox(Inches(0.9), Inches(4), Inches(3), Inches(1))
    tf = txBox1.text_frame
    p = tf.add_paragraph()
    p.text = arrayOfPrompts[0]
    p.font.name = 'Comic Sans' #Impact
    p.font.size = Pt(20)

    txBox2 = slide2.shapes.add_textbox(Inches(4.5), Inches(4), Inches(3), Inches(1))
    tf = txBox2.text_frame
    p = tf.add_paragraph()
    p.text = arrayOfPrompts[1]
    p.font.name = 'Comic Sans' #Impact
    p.font.size = Pt(20)


    txBox3 = slide2.shapes.add_textbox(Inches(7.9), Inches(4), Inches(3), Inches(1))
    tf = txBox3.text_frame
    p = tf.add_paragraph()
    p.text = arrayOfPrompts[2]
    p.font.name = 'Comic Sans' #Impact
    p.font.size = Pt(20)

    print(2)
# Third Slide
    slide3 = prs.slides.add_slide(blank_slide_layout)
    left = Inches(5.5)
    top = Inches(0.4)
    width = Inches(4.5)
    height = Inches(3)
    txBox = slide3.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = arrayOfPrompts[0]
    p.font.name = 'Comic Sans' #Impact
    p.font.size = Pt(40)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0, 0, 0)
    txBox = slide3.shapes.add_textbox(left, Inches(1.5), width, Inches(3.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = question(f"say something about {arrayOfPrompts[0] + title} in 45 words")
    p.font.name = 'Comic Sans' #Impact
    p.font.size = Pt(20)
    saveImg(arrayOfPrompts[0] + title, f"painting{timecode}")
    picPainting = slide3.shapes.add_picture(f"./AI/painting{timecode}.png", Inches(-4.5), Inches(0), width=Inches(10), height=Inches(10))
    print(3)


    # 4th slide
    slide4 = prs.slides.add_slide(blank_slide_layout)
    saveImg(arrayOfPrompts[1] + title, f"photograpy2{timecode}")
    saveImg(arrayOfPrompts[1] + title, f"photograpy{timecode}")
    picphotograpy2 = slide4.shapes.add_picture(f"./AI/photograpy2{timecode}.png", Inches(0), Inches(0), width=Inches(10), height=Inches(10))
    picphotograpy = slide4.shapes.add_picture(f"./AI/photograpy{timecode}.png", Inches(5), Inches(0), width=Inches(10), height=Inches(10))
    picphotograpy3 = slide4.shapes.add_picture("./AI/rectangle.png", Inches(0), Inches(-1), width=Inches(11.3), height=Inches(8.8))
    left = Inches(2.8)
    top = Inches(2)
    width = Inches(5.3)
    height = Inches(3)
    txBox = slide4.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = arrayOfPrompts[1]
    p.font.name = 'Comic Sans' #Impact
    p.font.size = Pt(40)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0, 0, 0)
    p = tf.add_paragraph()
    p.text = question(f"say something about {arrayOfPrompts[1] + title} in 35 words")
    p.font.size = Pt(17)

    print(4)
    # 5th slide
    slide5 = prs.slides.add_slide(blank_slide_layout)
    saveImg256(arrayOfPrompts[2] + title, f"writing2{timecode}")
    saveImg256(arrayOfPrompts[2] + title, f"writing{timecode}")
    picwriting2 = slide5.shapes.add_picture(f"./AI/writing2{timecode}.png", Inches(1), Inches(0.5), width=Inches(3), height=Inches(3))
    picwriting = slide5.shapes.add_picture(f"./AI/writing{timecode}.png", Inches(1), Inches(4), width=Inches(3), height=Inches(3))
    left = Inches(4)
    top = Inches(1)
    width = Inches(6)
    height = Inches(3)
    txBox = slide5.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = arrayOfPrompts[2]
    p.font.name = 'Comic Sans' #Impact
    p.font.size = Pt(40)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = question(f"say something about {arrayOfPrompts[2] + title} in 45 words")
    p.font.name = 'Comic Sans' #Impact
    p.font.size = Pt(20)

    print(5)

    # Conclusion
    slide6 = prs.slides.add_slide(blank_slide_layout)
    lastpic1 = slide6.shapes.add_picture(f"./AI/writing2{timecode}.png", Inches(7), Inches(1), width=Inches(3), height=Inches(3))
    lastpic2 = slide6.shapes.add_picture(f"./AI/writing{timecode}.png", Inches(3), Inches(0.5), width=Inches(3), height=Inches(3))
    lastpic3 = slide6.shapes.add_picture(f"./AI/photograpy{timecode}.png", Inches(1), Inches(0.5), width=Inches(3), height=Inches(3))
    lastpic4 = slide6.shapes.add_picture(f"./AI/photograpy2{timecode}.png", Inches(4), Inches(4), width=Inches(3), height=Inches(3))
    lastpic5 = slide6.shapes.add_picture(f"./AI/painting{timecode}.png", Inches(7.5), Inches(3), width=Inches(3), height=Inches(3))
    lastpic6 = slide6.shapes.add_picture(f"./AI/lolimg1{timecode}.png", Inches(5.8), Inches(4), width=Inches(3), height=Inches(3))
    lastpic7 = slide6.shapes.add_picture(f"./AI/lolimg2{timecode}.png", Inches(1), Inches(2), width=Inches(3), height=Inches(3))
    lastpic8 = slide6.shapes.add_picture(f"./AI/lolimg0{timecode}.png", Inches(1), Inches(3.8), width=Inches(3), height=Inches(3))
    lastpic9 = slide6.shapes.add_picture(f"./AI/rectangle.png", Inches(0), Inches(-1), width=Inches(11.3), height=Inches(8.8))
    left = Inches(2.8)
    top = Inches(2)
    width = Inches(5.3)
    height = Inches(3)
    txBox = slide6.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Conclusion"
    p.font.name = 'Comic Sans' #Impact
    p.font.size = Pt(40)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0, 0, 0)
    p = tf.add_paragraph()
    p.text = question(f"write a text for conclusion slide about {title} in 35 words")
    prs.save(f'{subtext+title}.pptx')
    os.remove(f"./AI/writing2{timecode}.png")
    os.remove(f"./AI/writing{timecode}.png")
    os.remove(f"./AI/photograpy{timecode}.png")
    os.remove(f"./AI/photograpy2{timecode}.png")
    os.remove(f"./AI/painting{timecode}.png")
    os.remove(f"./AI/lolimg1{timecode}.png")
    os.remove(f"./AI/lolimg2{timecode}.png") 
    os.remove(f"./AI/lolimg0{timecode}.png")
    os.remove(f"./AI/firstPage{timecode}.png")
# firstSlide(title="Hobbies", subtext="Suvernev Bogdan", font = "Comic Sans")
print(312)