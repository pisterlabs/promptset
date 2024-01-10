

# ----------------------------------------------------------------------------
# Project: Semantic Search Module for the Alter Brain project
# File:    lib__embedded_context.py
# 
# This lib is the Semantic Search Module for the Alter Brain project. It implements a 
# system for understanding and processing natural language to facilitate 
# information retrieval based on semantics rather than traditional keyword-based search.
# 
# Author:  Michel Levy Provencal
# Brightness.ai - 2023 - contact@brightness.fr
# ----------------------------------------------------------------------------

import pandas as pd
import os
import csv
import openai
from openai import OpenAI
#from openai.embeddings_utils import get_embedding
from transformers import GPT2TokenizerFast
from dotenv import load_dotenv
import random
import numpy as np
import sys
import time
import requests
import os.path
import PyPDF2
import docx
import json
import pptx
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
import pytesseract
from PIL import Image
from openpyxl import load_workbook
import requests
from bs4 import BeautifulSoup
from markdownify import markdownify as md
from urllib.parse import urlparse, urljoin
from lib__env import *
os.environ["TOKENIZERS_PARALLELISM"] = "false"

"""
#############################################################################################################
    
    ## TOOLS
    
#############################################################################################################
"""


# Fonction pour générer un nom de fichier unique
def generate_unique_filename(prefix, suffix):
    """
    Generates a unique filename by appending a random number between 1 and 9999 to the given prefix,
    and then appending a specified suffix.
    :param prefix: The prefix of the filename.
    :param suffix: The suffix of the filename (usually the file extension).
    :return: A string representing a unique filename with the format 'prefix_randomNumber.suffix'.
    """ 
    random_number = random.randint(1, 9999)
    return f"{prefix}_{random_number}.{suffix}"

# ----------------------------------------------------------------------------
# Fonction pour convertir un fichier PDF en texte
def convert_pdf_to_text(file_path):
    """
    Converts a PDF file to text using PyPDF2.
    This function reads a PDF file, extracts the text from each page, and then concatenates the extracted text from all pages into a single string.
    :param file_path: Path to the PDF file.
    :return: The text extracted from the PDF file.
    """
    with open(file_path, "rb") as file:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text()
        return text


# ----------------------------------------------------------------------------
# Fonction pour convertir un fichier .docx en texte
def convert_docx_to_text(file_path):
    """
    Converts the contents of a .docx file into plain text.
    
    This function takes as input a path to a .docx file, opens the file,
    and extracts the text from each paragraph in the document. The extracted
    text from each paragraph is then joined together with newline characters
    in between each paragraph to form a single string of text.

    :param file_path: The path to the .docx file.
    :return: A single string containing the text of the .docx file.
    """
    doc = docx.Document(file_path)
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text

# ----------------------------------------------------------------------------
# Fonction pour convertir un fichier CSV en texte
def convert_csv_to_text(file_path):
    """
    Converts a CSV file into a text format.
    :param file_path: The path to the CSV file.
    :return: A string representation of the CSV file.
    """  
    with open(file_path, "r") as file:
        csv_reader = csv.reader(file)
        text = "\n".join([",".join(row) for row in csv_reader])
        return text

# ----------------------------------------------------------------------------
# Fonction pour convertir un fichier JSON en texte
def convert_json_to_text(file_path):
    """
    Converts a JSON file into a formatted text string.
    :param file_path: The path to the JSON file.
    :return: A string representing the JSON data, formatted with indentation.
    """
    with open(file_path, "r") as file:
        data = json.load(file)
        text = json.dumps(data, indent=4)
        return text


# ----------------------------------------------------------------------------
# Fonction pour convertir un fichier Excel en texte
def convert_excel_to_text(file_path):
    """
    Converts an Excel workbook into a text format. Each cell is separated by a comma,
    and each row is separated by a new line.

    :param file_path: The path of the Excel file.
    :return: A string representing the content of the Excel file.
    """
    workbook = load_workbook(file_path)
    text = ""
    for sheet in workbook:
        for row in sheet.values:
            text += ",".join([str(cell) for cell in row])
            text += "\n"
    return text



# ----------------------------------------------------------------------------
# Fonction pour convertir un fichier .pptx en texte
def convert_pptx_to_text(file_path):
    """
    Converts the content of a PowerPoint presentation (.pptx) into text.
    :param file_path: The path to the .pptx file.
    :return: A string containing the text content of the presentation.
    """ 
    presentation = pptx.Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text
                    text += "\n"
    return text


# ----------------------------------------------------------------------------
# Fonction pour convertir un fichier XML en texte
def convert_xml_to_text(file_path):
    """
    Convertit le contenu d'un fichier XML en texte brut.
    :param file_path: Le chemin vers le fichier XML.
    :return: Le texte extrait du fichier XML.
    """
    tree = ET.parse(file_path)
    root = tree.getroot()
    text = ET.tostring(root, encoding="utf-8", method="text").decode("utf-8")
    return text

# ----------------------------------------------------------------------------
# Fonction pour convertir un fichier HTML en texte
def convert_html_to_text(file_path):
    """
    Converts an HTML file into a plain text by removing all the HTML tags.
    :param file_path: The path to the HTML file.
    :return: The text content of the HTML file.
    """
    with open(file_path, "r") as file:
        soup = BeautifulSoup(file, "html.parser")
        text = soup.get_text()
        return text

# ----------------------------------------------------------------------------
# Fonction pour convertir une image en texte à l'aide de l'OCR
def convert_image_to_text(file_path):
    """
    Converts an image into text using Optical Character Recognition (OCR).
    :param file_path: Path to the image file.
    :return: Extracted text from the image.
    """
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image, lang="eng")
    return text



def convert_text_to_text(file_path):
    """
    Converts the content of a text file to UTF-8 and returns a text string.
    :param file_path: The path to the text file.
    :return: The text extracted from the text file.
    """
    with open(file_path, "r", encoding="utf-8") as file:
        text = file.read()
    return text


""" VERSION WITH COMPREHENSSION LIST
# ----------------------------------------------------------------------------
# Function that concat all files contained in a folder in a text
def concat_files_in_text(path):
    
    # Concatenates the files of a directory into a single text.
    # :param path: Directory path.
    # :return: Concatenated text.
       
    files = [os.path.join(path, f) for f in os.listdir(path) if os.path.isfile(os.path.join(path, f))]
    texts = []
    for file in files:
        with open(file, 'r') as f:
            texts.append(f.read())
    return ' '.join(texts)

"""

def concat_files_in_text(path):
    """
    Concatenates the files of a directory into a single text.
    :param path: Directory path.
    :return: Concatenated text.
    """
    files = []
    for f in os.listdir(path):
        full_path = os.path.join(path, f)
        if os.path.isfile(full_path):
            files.append(full_path)
    
    texts = []
    for file in files:
        with open(file, 'r') as f:
            file_content = f.read()
            texts.append(file_content)
            
    return ' '.join(texts)



# ----------------------------------------------------------------------------
# Function that split the txt file into blocks. Uses batch (limit) of 2000 words with gpt3.5 and 4000 with gpt4
def split_text_into_blocks(text, limit=4000):
    
    """
    This function splits a given text into chunks, or "blocks", each containing a certain number of words specified by the 'limit' parameter. It's particularly designed for use with GPT-3.5 (limit of 2000 words) and GPT-4 (limit of 4000 words).
    Each block is constructed by sequentially adding words from the input text until the block size (the number of words in the block) reaches the limit. If adding another word would exceed the limit, the function checks for the last sentence or line delimiter in the current block (a period or newline character), then separates the block at that delimiter.
    If there is no delimiter in the current block, the entire block is added to the list of blocks and the next word starts a new block. If a delimiter is found, the block is split at the delimiter, and the remaining text (if any) is added to the next block along with the next word.
    The function returns a list of blocks.
    :param text: The input text to be split into blocks.
    :param limit: The maximum number of words allowed in each block. Default is 4000.
    :return: A list of text blocks obtained from the input text.
    """
    ### TODO : Adapt the limit to other LLMs (LLAMAS, All-GPT, etc.)
    
    blocks = []
    current_block = ""
    words = text.split()

    for word in words:
        if len(current_block + word) + 1 < limit:
            current_block += word + " "
        else:
            last_delimiter_index = max(current_block.rfind(". "), current_block.rfind("\n"))

            if last_delimiter_index == -1:
                blocks.append(current_block.strip())
                current_block = word + " "
            else:
                delimiter = current_block[last_delimiter_index]
                blocks.append(current_block[:last_delimiter_index + (1 if delimiter == '.' else 0)].strip())
                current_block = current_block[last_delimiter_index + (2 if delimiter == '.' else 1):].strip() + " " + word + " "

    if current_block.strip():
        blocks.append(current_block.strip())

    return blocks


# ----------------------------------------------------------------------------
# Function that write blocks into filename
def write_blocks_to_csv(blocks, path, filename):
    """
    This function takes a list of blocks (data items) and a filename as input parameters, then writes these blocks into a CSV file specified by the given filename.
    The blocks are written row by row in the CSV file, with each block making up a single row. 
    The CSV file is created with a specific encoding (UTF-8), and using specific settings for delimiter and quotechar for CSV data formatting.
    :param blocks: A list of data items to be written into the CSV file.
    :param filename: The name of the CSV file where the data should be written.
    """
    with open(path + filename, "w", newline="", encoding="utf-8") as csvfile:
        csvwriter = csv.writer(csvfile, delimiter=',', quotechar='"', quoting=csv.QUOTE_MINIMAL)
        
        # Write the header
        csvwriter.writerow(['Datas'])

        for block in blocks:
            csvwriter.writerow([block])




def get_embedding(text, engine="text-embedding-ada-002"):
    """
    This function takes in a piece of text and a model engine as input parameters, and returns an embedding for the input text.
    It utilizes OpenAI's Embedding API to generate the embedding based on the specified model.

    The function first replaces newline characters in the input text with spaces, as the embedding models typically 
    handle single continuous strings of text.

    :param text: The input text for which to generate an embedding.
    :param engine: The model engine to use for generating the embedding. Default is 'text-embedding-ada-002'.
    :return: The generated embedding for the input text.
    """
    
    text = text.replace("\n", " ")
    client = OpenAI(api_key=os.environ['OPENAI_API_KEY'])
    
    # DEBUG
    print("get embedding for " + text)
    
    response = client.embeddings.create(input=[text], model=engine)
    response_dict = response.model_dump()  # Conversion de la réponse en dictionnaire
    return response_dict['data'][0]['embedding']





# Updated function to create embeddings with the new OpenAI SDK
def create_embeddings(path, filename):
    """
    This function reads text data from a specified CSV file and creates embeddings for each text entry using OpenAI's
    Embedding API. It then saves the generated embeddings back to a new CSV file.

    The function uses a GPT-2 tokenizer to tokenize the text data. It then filters out rows where the number of tokens
    exceeds 8000 and keeps the last 2000 records. The function also drops rows with missing values from the data.

    The embeddings are generated using the 'text-embedding-ada-002' model by default, and the generated embeddings are 
    saved as a new column in the DataFrame.

    The function finally saves the DataFrame, with the embeddings, to a new CSV file.

    :param path: The directory path where the input and output CSV files are located.
    :param filename: The name of the input CSV file from which to read the text data.
    """
    
    load_dotenv(DOTENVPATH)
    tokenizer = GPT2TokenizerFast.from_pretrained("gpt2")
       
    # Open the input CSV file and read it into a Pandas DataFrame.
    df_full = pd.read_csv(path + filename, sep=';', on_bad_lines='skip', encoding='utf-8')

    # Rename the dataframe's columns
    df = df_full[['Datas']]
    # Remove rows with missing values 
    df = df.dropna()

    # Count the number of tokens in each row and filter the DataFrame
    df['n_tokens'] = df.Datas.apply(lambda x: len(tokenizer.encode(x)))
    df = df[df.n_tokens < 8000].tail(2000)

    df['ada_embedding'] = df.Datas.apply(get_embedding)

    # Write the DataFrame to a new CSV file
    df.to_csv(path + "emb_" + filename, index=False)

    return path



# ----------------------------------------------------------------------------
# Function that reads and processes a CSV file and returns a DataFrame
def read_and_process_csv(index_filename):
    """
    This function takes as input the filename of a CSV file and reads this file into a Pandas DataFrame. 
    It then processes the 'ada_embedding' column of the DataFrame, converting the string representations 
    of embeddings stored in this column back into actual Numpy array objects.

    The function first reads the CSV file using Pandas' read_csv function, creating a DataFrame where each 
    row corresponds to a data item from the CSV file and each column corresponds to a field in the data items.

    It then applies the eval function to each item in the 'ada_embedding' column to convert the string 
    representations of the embeddings back into list objects. These lists are then further converted into 
    Numpy arrays using the np.array function. This processed 'ada_embedding' column replaces the original 
    column in the DataFrame.

    :param index_filename: The filename of the CSV file to read.
    :return: The DataFrame created from the CSV file, with the 'ada_embedding' column processed.
    """
    
    df = pd.read_csv(index_filename)
    df['ada_embedding'] = df.ada_embedding.apply(eval).apply(np.array)
    return df



# ----------------------------------------------------------------------------
# Function that gets an embedding vector for a given text
def get_search_vector(text):
    """
    This function takes as input a piece of text and returns an embedding vector for the input text. 
    It utilizes the 'get_embedding' function to generate the 
    embedding vector.

    The function is a convenience wrapper around the 'get_embedding' function, simplifying its use by 
    directly passing the input text and relying on the 'get_embedding' function's default parameters 
    for generating the embedding.

    :param text: The input text for which to generate an embedding vector.
    :return: The embedding vector for the input text.
    """

    return get_embedding(text)



# ----------------------------------------------------------------------------
# Function that finds similar rows in a DataFrame based on an input vector
def find_similar_rows(df, searchvector, n_results):
    """
    This function takes as input a DataFrame, a search vector, and a number of results to return.
    It calculates the cosine similarity between the search vector and the embeddings in the DataFrame. 
    The rows with the highest cosine similarity are then sorted and the top 'n_results' rows are returned.

    The function adds a new column, 'similarities', to the DataFrame. For each row, it computes the dot 
    product between the 'ada_embedding' of the row and the search vector, which is equivalent to calculating 
    the cosine similarity when the vectors are normalized.

    The rows in the DataFrame are then sorted in descending order of 'similarities', and the top 'n_results' 
    rows are returned as a new DataFrame.

    :param df: The input DataFrame, each row of which should have an 'ada_embedding' column containing a vector.
    :param searchvector: The vector to compare against the 'ada_embedding' of each row in the DataFrame.
    :param n_results: The number of top similar rows to return.
    :return: A DataFrame containing the top 'n_results' similar rows to the input vector.
    """

    df['similarities'] = df.ada_embedding.apply(lambda x: np.dot(x, searchvector))
    res = df.sort_values('similarities', ascending=False).head(n_results)
    return res


# ----------------------------------------------------------------------------
# Function that validates a DataFrame and extracts the combined data
def validate_and_get_combined(res):
    """
    This function takes a DataFrame as input, performs several validation checks on it, 
    and then extracts and returns the combined data from the 'Datas' column of the DataFrame.

    The function first checks if the 'Datas' column exists in the DataFrame. If it does not, 
    a ValueError is raised. It then checks if the DataFrame is empty. If it is, a ValueError 
    is raised. Finally, it checks if the index of the DataFrame is of type 'int64'. If it is 
    not, a ValueError is raised.

    Once these validation checks have passed, the function concatenates all the strings in the 
    'Datas' column of the DataFrame, with each string separated by a newline character. This combined 
    string is then returned.

    :param res: The input DataFrame to validate and extract combined data from.
    :return: A string consisting of all the data from the 'Datas' column of the DataFrame, concatenated with newline characters.
    """
    
    if 'Datas' not in res.columns:
        raise ValueError("La colonne 'Datas' n'existe pas dans le DataFrame")

    if res.empty:
        raise ValueError("Le DataFrame est vide")

    if res.index.dtype != 'int64':
        raise ValueError("L'index du DataFrame n'est pas de type entier")

    return '\n'.join(res['Datas'].values)





# ----------------------------------------------------------------------------
# Function that converts all the files in a folder to text and store them in a new subfolder named 'text'
def create_text_folder(folder_path):
    """
    TODO: Write documentation
    """
    
    # Folder containing the files to convert
    source_folder = folder_path

    # Destination folder for the converted text files
    destination_folder = folder_path + "text_tmp"

    # List of supported file formats
    supported_formats = {
        ".pdf": convert_pdf_to_text,
        ".docx": convert_docx_to_text,
        ".csv": convert_csv_to_text,
        ".json": convert_json_to_text,
        ".xls": convert_excel_to_text,
        ".xlsx": convert_excel_to_text,
        ".pptx": convert_pptx_to_text,
        ".xml": convert_xml_to_text,
        ".html": convert_html_to_text,
        ".jpg": convert_image_to_text,
        ".jpeg": convert_image_to_text,
        ".png": convert_image_to_text,
        ".txt": convert_text_to_text
    }
    
    # Create the destination folder if it doesn't exist
    if not os.path.exists(destination_folder):
        os.makedirs(destination_folder)

    # Iterate through all the files in the source folder
    for file_name in os.listdir(source_folder):
        source_file_path = os.path.join(source_folder, file_name)
        if os.path.isfile(source_file_path):
            file_name_without_ext, file_ext = os.path.splitext(file_name)
            if file_ext.lower() in supported_formats:
                converter = supported_formats[file_ext.lower()]
                text = converter(source_file_path)
                destination_file_name = generate_unique_filename(file_name_without_ext, "txt")
                destination_file_path = os.path.join(destination_folder, destination_file_name)
                with open(destination_file_path, "w", encoding="utf-8") as file:
                    file.write(text)

    return(str(destination_folder))


# ----------------------------------------------------------------------------
# Function that converts an url to text
def get_text_from_url(url):
    """
    Récupère le texte d'une page web à partir de son URL.
    
    Args:
        url (str): L'URL de la page web à récupérer.
        
    Returns:
        str: Le texte de la page web, ou une chaîne vide en cas d'erreur.
    """
    text = ""

    headers = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3"
    }

    try:
        # Récupérer le contenu HTML de l'URL
        response = requests.get(url, headers=headers)
        response.raise_for_status()
    except requests.exceptions.RequestException as e:
        print(f"Erreur lors de la requête vers {url} : {e}")
        return text

    try:
        soup = BeautifulSoup(response.text, 'html.parser')

        # Extraire le texte de la structure HTML
        text = soup.get_text()
    except Exception as e:
        print(f"Erreur lors de l'analyse HTML pour {url} : {e}")
        return text

    return text



"""
#############################################################################################################
    
    ## FUNCTIONS TO INDEX & SEARCH EMBEDDINGS
    
#############################################################################################################
"""



## ----------------------------------------------------------------------------
## Function that creates a csv index file containing embeddings, from a folder named path. 
## The index is stored in the same folder, and named : emb_csv_XXX.csv
## The function returns the name of the index file : emb_csv_XXX.csv
## (while the concatenated text is stored in a txt file named txt_XXX.txt and the csv file containing the blocks is named csv_XXX.csv)
def build_index(folder_path):
    """
    This function reads multiple text files from a specified folder, concatenates the text, 
    splits the concatenated text into blocks of specified length, writes these blocks into a CSV file, 
    creates embeddings for each block using the create_embeddings function, and finally, saves the 
    embeddings back to the CSV file.

    The function first generates a random number which is used to create unique filenames for the 
    intermediate text and CSV files. It then reads and concatenates all the text files from the 
    specified folder into a single string of text.

    The function then calls the split_text_into_blocks function to split the text into blocks of 
    up to 4000 characters each. The resulting list of blocks is written to a new CSV file.

    The function then calls the create_embeddings function to create embeddings for each text block 
    in the CSV file. The embeddings are saved back to the CSV file.

    Finally, the function returns the name of the CSV file containing the embeddings.

    :param folder_path: The directory path where the text files are located and where the CSV file will be saved.
    :return: The name of the CSV file containing the embeddings.
    """
    
    # transform files into text, create a subfolder named 'txt' and save the text files in it. 
    text_folder_path = create_text_folder(folder_path)
    
    # Concatenate files in text
    text = concat_files_in_text(text_folder_path)

    # OLD VERSION : Save text in a csv file named csv(random number).csv
    # random_num = random.randint(1000,9999)  # Generates a random number between 1000 and 9999

    # Call the function split_text_into_blocks() to split the text into blocks
    blocks = split_text_into_blocks(text, limit=4000)

    # Call the function write_blocks_to_csv() to write the blocks into a csv file
    write_blocks_to_csv(blocks, folder_path, 'index.csv')

    # Create embeddings for the csv file
    brain_id = create_embeddings(folder_path, 'index.csv')
    
    return(brain_id)


## ----------------------------------------------------------------------------
## Function that creates embdeedings from the content of an url
def build_index_url(url):
    
    index_text = get_text_from_url(url)
    
    #créer un fichier temporaire    
    timestamp = time.strftime("%Y%m%d-%H%M%S")
    folder_path = "datas/" + timestamp + "/"

    # Créer le nouveau dossier
    os.makedirs(folder_path, exist_ok=True)

    # enregistrer le texte dans un fichier txt
    with open(folder_path + 'url_index.txt', 'w', encoding='utf-8') as f:
        f.write(index_text)
            
    build_index(folder_path)
    
    return(timestamp)

# ----------------------------------------------------------------------------
# Function that finds the context for a given query in an index file
def find_context(text, index_filename, n_results=5):
    """
    This function takes as input a piece of text, the filename of a CSV file containing indexed data, 
    and an optional number of results to return. It finds the most similar data items to the input 
    text in the indexed data and returns the combined data from these items.

    The function first loads environment variables from a .env file, including the OpenAI API key. 
    It then reads and processes the indexed data from the CSV file into a DataFrame.

    The function creates an embedding for the input text using the get_search_vector function. 
    This embedding is compared with the embeddings of the data items in the DataFrame to find 
    the most similar items.

    The most similar items are sorted by their similarity scores, and the top 'n_results' items are 
    selected. The combined data from these items is extracted and returned.

    :param text: The input text for which to find similar data items.
    :param index_filename: The filename of the CSV file containing the indexed data.
    :param n_results: The number of most similar data items to return. Default is 5.
    :return: The combined data from the most similar data items.
    """
    
    
    
    load_dotenv(".env")  # Load the environment variables from the .env file.
    openai.api_key = os.environ.get("OPENAI_API_KEY")
    if not os.path.exists(index_filename):
        return ""
    
    df = read_and_process_csv(index_filename)

    searchvector = get_search_vector(text)

    res = find_similar_rows(df, searchvector, n_results)

    return validate_and_get_combined(res)




# ----------------------------------------------------------------------------
# Function that queries the OpenAI language model with a context and query
def query_extended_llm(text, index_filename, model="gpt-4"):
    """
    This function takes as input a piece of text, the filename of a CSV file containing indexed data, 
    and an optional AI model to use. It queries the OpenAI language model with a context derived from 
    the most similar data items to the input text, and a prompt derived from the input text itself. 
    It returns the response from the language model.

    The function first finds the context for the input text using the find_context function. It then 
    loads environment variables from a .env file, including the OpenAI API key. 

    The function then enters a loop where it attempts to query the language model with the context and 
    the input text as a prompt. If it encounters an exception during this process, it waits for 5 seconds 
    and then tries again, up to a maximum of 10 attempts.

    If the function is able to successfully query the language model, it returns the model's response as 
    a string. If it is unable to do so after 10 attempts, it prints an error message and terminates the 
    program.

    :param text: The input text for which to query the language model.
    :param index_filename: The filename of the CSV file containing the indexed data.
    :param model: The AI model to use for the query. Default is 'gpt-4'.
    :return: The response from the language model.
    """
    context = find_context(text, index_filename)
    client = OpenAI(api_key=os.environ['OPENAI_API_KEY'])  # Initialisation du client OpenAI

    attempts = 0
    prompt = "Context : " + context + "\n\n" + "Query : " + text

    while attempts < 10:
        try:
            response = client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "user", "content": prompt}
                ]
            )
            message = response.choices[0].message.content
            return message.strip()

        except Exception as e:
            error_code = type(e).__name__
            error_reason = str(e)
            attempts += 1
            print(f"Erreur : {error_code} - {error_reason}. Nouvel essai dans 5 secondes...")
            time.sleep(int(attempts)*2)

    print("Erreur : Echec de la création de la completion après 10 essais")
    sys.exit()
