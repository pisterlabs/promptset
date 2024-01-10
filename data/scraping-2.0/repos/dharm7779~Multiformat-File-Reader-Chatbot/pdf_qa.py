import textract  # For extracting text from various document formats
from pptx import Presentation  # For reading PPTX files
from docx import Document  # For reading DOCX files
from langchain.document_loaders import PDFPlumberLoader  # Document text extraction
from langchain.text_splitter import CharacterTextSplitter, TokenTextSplitter  # Text splitting
from transformers import pipeline  # Hugging Face Transformers for NLP
# Prompt templates for language models
from langchain.prompts import PromptTemplate
from langchain.chat_models import ChatOpenAI  # Chat-based language models
from langchain.vectorstores import Chroma  # Vector database for embeddings
from langchain.chains import RetrievalQA  # Retrieval-based question-answering
from langchain import HuggingFacePipeline  # Hugging Face Transformers pipeline
# Embeddings from Hugging Face models
from langchain.embeddings import HuggingFaceInstructEmbeddings, HuggingFaceEmbeddings
from langchain.embeddings.openai import OpenAIEmbeddings  # OpenAI embeddings
from openpyxl import load_workbook  # Reading Excel files
from langchain.llms import OpenAI  # Language models from OpenAI
from PyPDF2 import PdfReader  # Reading PDF files
from constants import *  # Constants and configurations
from transformers import AutoTokenizer , AutoModelForCausalLM # Tokenization for Transformers
import torch  # PyTorch for deep learning
import os  # Operating system related functions
import re  # Regular expressions for text processing


def extract_text_from_excel(xlsx_file):
    workbook = load_workbook(xlsx_file)
    text = ""
    for sheet in workbook.sheetnames:
        for row in workbook[sheet].iter_rows(values_only=True):
            text += " ".join(map(str, row)) + "\n"
    return text

# Function to extract text from PowerPoint (PPTX) file

# Function to extract text from Word (DOCX) file


def extract_text_from_word_docx(docx_file):
    doc = Document(docx_file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

# Function to extract text from PowerPoint (PPTX) file


def extract_text_from_presentation(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text


class PdfQA:

    def __init__(self, config: dict = {}):
        self.config = config
        self.embedding = None
        self.vectordb = None
        self.llm = None
        self.qa = None
        self.retriever = None

    # The following class methods are useful to create global GPU model instances
    # This way we don't need to reload models in an interactive app,
    # and the same model instance can be used across multiple user sessions
    @classmethod
    def create_instructor_xl(cls):
        device = "cuda" if torch.cuda.is_available() else "cpu"
        return HuggingFaceInstructEmbeddings(model_name=EMB_INSTRUCTOR_XL, model_kwargs={"device": device})

    @classmethod
    def create_sbert_mpnet(cls):
        device = "cuda" if torch.cuda.is_available() else "cpu"
        return HuggingFaceEmbeddings(model_name=EMB_SBERT_MPNET_BASE, model_kwargs={"device": device})

    @classmethod
    def create_flan_t5_xl(cls, load_in_8bit=False):
        return pipeline(
            task="text2text-generation",
            model="google/flan-t5-xl",
            max_new_tokens=2000,
            model_kwargs={"device_map": "auto", "load_in_8bit": load_in_8bit,
                          "max_length": 512, "temperature": 0.}
        )

    @classmethod
    def create_flan_t5_small(cls, load_in_8bit=False):
        # Local flan-t5-small for inference
        # Wrap it in HF pipeline for use with LangChain
        model = "google/flan-t5-small"
        tokenizer = AutoTokenizer.from_pretrained(model)
        return pipeline(
            task="text2text-generation",
            model=model,
            tokenizer=tokenizer,
            max_new_tokens=1000,
            model_kwargs={"device_map": "auto", "load_in_8bit": load_in_8bit,
                          "max_length": 512, "temperature": 0.}
        )

    @classmethod
    def create_flan_t5_base(cls, load_in_8bit=True):
        # Wrap it in HF pipeline for use with LangChain
        model = "google/flan-t5-base"
        tokenizer = AutoTokenizer.from_pretrained(model)
        return pipeline(
            task="text2text-generation",
            model=model,
            tokenizer=tokenizer,
            max_new_tokens=1000,
            model_kwargs={"device_map": "auto", "load_in_8bit": load_in_8bit,
                          "max_length": 1000, "temperature": 0.}
        )

    @classmethod
    def create_flan_t5_large(cls, load_in_8bit=False):
        # Wrap it in HF pipeline for use with LangChain
        model = "google/flan-t5-large"
        tokenizer = AutoTokenizer.from_pretrained(model)
        return pipeline(
            task="text2text-generation",
            model=model,
            tokenizer=tokenizer,
            max_new_tokens=1000,
            model_kwargs={"device_map": "auto", "load_in_8bit": load_in_8bit,
                          "max_length": 512, "temperature": 0.}
        )

    @classmethod
    def create_fastchat_t5_xl(cls, load_in_8bit=False):
        return pipeline(
            task="text2text-generation",
            model="lmsys/fastchat-t5-3b-v1.0",
            max_new_tokens=1000,
            model_kwargs={"device_map": "auto", "load_in_8bit": load_in_8bit,
                          "max_length": 512, "temperature": 0.01}
        )

    @classmethod
    def create_falcon_instruct_small(cls, load_in_8bit=False):
        model = "tiiuae/falcon-7b-instruct"

        tokenizer = AutoTokenizer.from_pretrained(model)
        hf_pipeline = pipeline(
            task="text-generation",
            model=model,
            tokenizer=tokenizer,
            trust_remote_code=True,
            max_new_tokens=100,
            model_kwargs={
                "device_map": "auto",
                "load_in_8bit": load_in_8bit,
                "max_length": 512,
                "temperature": 0.01,
                "torch_dtype": torch.bfloat16,
            }
        )
        return hf_pipeline

    def init_embeddings(self) -> None:
        # OpenAI ada embeddings API
        if self.config["embedding"] == EMB_OPENAI_ADA:
            self.embedding = OpenAIEmbeddings()
        elif self.config["embedding"] == EMB_INSTRUCTOR_XL:
            # Local INSTRUCTOR-XL embeddings
            if self.embedding is None:
                self.embedding = PdfQA.create_instructor_xl()
        elif self.config["embedding"] == EMB_SBERT_MPNET_BASE:
            # this is for SBERT
            if self.embedding is None:
                self.embedding = PdfQA.create_sbert_mpnet()
        else:
            self.embedding = None  # DuckDb uses sbert embeddings
            # raise ValueError("Invalid config")

    def init_models(self) -> None:
        """ Initialize LLM models based on config """
        load_in_8bit = self.config.get("load_in_8bit", False)
        # OpenAI GPT 3.5 API
        if self.config["llm"] == LLM_OPENAI_GPT35:
            # OpenAI GPT 3.5 API
            pass
        elif self.config["llm"] == LLM_FLAN_T5_SMALL:
            if self.llm is None:
                self.llm = PdfQA.create_flan_t5_small(
                    load_in_8bit=load_in_8bit)
        elif self.config["llm"] == LLM_FLAN_T5_BASE:
            if self.llm is None:
                self.llm = PdfQA.create_flan_t5_base(load_in_8bit=load_in_8bit)
        elif self.config["llm"] == LLM_FLAN_T5_LARGE:
            if self.llm is None:
                self.llm = PdfQA.create_flan_t5_large(
                    load_in_8bit=load_in_8bit)
        elif self.config["llm"] == LLM_FLAN_T5_XL:
            if self.llm is None:
                self.llm = PdfQA.create_flan_t5_xl(load_in_8bit=load_in_8bit)
        elif self.config["llm"] == LLM_FLAN_T5_XXL:
            if self.llm is None:
                self.llm = PdfQA.create_flan_t5_xxl(load_in_8bit=load_in_8bit)
        elif self.config["llm"] == LLM_FASTCHAT_T5_XL:
            if self.llm is None:
                self.llm = PdfQA.create_fastchat_t5_xl(
                    load_in_8bit=load_in_8bit)
        elif self.config["llm"] == LLM_FALCON_SMALL:
            if self.llm is None:
                self.llm = PdfQA.create_falcon_instruct_small(
                    load_in_8bit=load_in_8bit)

        else:
            raise ValueError("Invalid config")


# ...

    def vector_db_documents(self, document) -> None:
        """
        Creates a vector database for the embeddings and persists them or loads a vector database from the persist directory.
        """
        document_type = self.config.get("document_type", None)
        document_path = self.config.get("document_path", None)
        persist_directory = self.config.get("persist_directory", None)

        if persist_directory and os.path.exists(persist_directory):
            # Load from the persist db
            self.vectordb = Chroma(
                persist_directory=persist_directory, embedding_function=self.embedding)
        elif document_path and os.path.exists(document_path):
            # 1. Extract the documents
            if document_type == "pdf":
                # Handle PDF
                pdf_reader = PdfReader(document_path)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text()
            elif document_type == "docx":
                # Handle Word document
                text = extract_text_from_word_docx(document_path)
            elif document_type == "pptx":
                # Handle Presentation (PPTX) file
                text = extract_text_from_presentation(document_path)
            elif document_type == "xlsx":
                # Handle Excel file
                text = extract_text_from_excel(document_path)
            else:
                raise ValueError("Unsupported document type")

            # 2. Split the texts
            text_splitter = CharacterTextSplitter(
                chunk_size=10000, chunk_overlap=0)
            texts = text_splitter.split_text(text)

            # 3. Create Embeddings and add to chroma store
            # TODO: Validate if self.embedding is not None
            self.vectordb = Chroma.from_texts(
                texts=texts, embedding=self.embedding, persist_directory=persist_directory)
        else:
            raise ValueError("No supported file found")

    def retreival_qa_chain(self):
        """
        Creates retrieval qa chain using vectordb as retrivar and LLM to complete the prompt
        """
        # TODO: Use custom prompt
        self.retriever = self.vectordb.as_retriever(search_kwargs={"k": 3})

        if self.config["llm"] == LLM_OPENAI_GPT35:
            # Use ChatGPT API
            self.qa = RetrievalQA.from_chain_type(llm=OpenAI(model_name=LLM_OPENAI_GPT35, temperature=0.), chain_type="stuff",
                                                  retriever=self.vectordb.as_retriever(search_kwargs={"k": 3}))
        else:
            hf_llm = HuggingFacePipeline(
                pipeline=self.llm, model_id=self.config["llm"])

            self.qa = RetrievalQA.from_chain_type(
                llm=hf_llm, chain_type="stuff", retriever=self.retriever)
            if self.config["llm"] == LLM_FLAN_T5_SMALL or self.config["llm"] == LLM_FLAN_T5_BASE or self.config["llm"] == LLM_FLAN_T5_LARGE:
                question_t5_template = """
                context: {context}
                question: {question}
                answer: 
                """
                QUESTION_T5_PROMPT = PromptTemplate(
                    template=question_t5_template, input_variables=[
                        "context", "question"]
                )
                self.qa.combine_documents_chain.llm_chain.prompt = QUESTION_T5_PROMPT
            self.qa.combine_documents_chain.verbose = True
            self.qa.return_source_documents = True

    def answer_query(self, question: str) -> str:
        """
        Answer the question
        """

        answer_dict = self.qa({"query": question, })
        print(answer_dict)
        answer = answer_dict["result"]
        if self.config["llm"] == LLM_FASTCHAT_T5_XL:
            answer = self._clean_fastchat_t5_output(answer)
        return answer

    def _clean_fastchat_t5_output(self, answer: str) -> str:
        # Remove <pad> tags, double spaces, trailing newline
        answer = re.sub(r"<pad>\s+", "", answer)
        answer = re.sub(r"  ", " ", answer)
        answer = re.sub(r"\n$", "", answer)
        return answer
