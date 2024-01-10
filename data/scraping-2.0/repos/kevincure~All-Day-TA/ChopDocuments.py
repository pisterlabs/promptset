# This code takes all pdfs in documents, scans them, then generates paragraph embeddings on a sliding scale of 200 tokens (roughly 150 words)
# Run this before you run EmbedDocuments.py or app.py
# You need an OpenAI key saved in APIkey.txt
# Note that if your PDFs are not searchable, this won't work - use a third party tool to convert them to txt or doc first.  You
#   can look at the "-originaltext.csv" file created here and scan real quick to see if the text looks corrupted for any of your docs


import os
import time
import chardet
from PyPDF2 import PdfReader 
import nltk
import pandas as pd
import numpy as np
import json
import io
import re
import openai
import shutil
from pptx import Presentation
# you need to pip install python-docx, not docx
import docx
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

# Set the desired chunk size and overlap size
# chunk_size is how many tokens we will take in each block of text
# overlap_size is how much overlap. So 200, 100 gives you chunks of between the 1st and 200th word, the 100th and 300th, the 200 and 400th...
# I have in no way optimized these
chunk_size = 200
overlap_size = 100

# load user settings and api key
def read_settings(file_name):
    settings = {}
    with open(file_name, "r") as f:
        for line in f:
            key, value = line.strip().split("=")
            settings[key] = value
    return settings
settings = read_settings("settings.txt")
filedirectory = settings["filedirectory"]
classname = settings["classname"]
professor = settings["professor"]
assistants = settings["assistants"]
classdescription = settings["classdescription"]
assistant_name = settings['assistantname']
instruct = settings['instructions']
num_chunks = int(settings['num_chunks'])
# get API_key
with open("APIkey.txt", "r") as f:
    openai.api_key = f.read().strip()
# Check if the subfolder exists, if not, create it
output_folder = "Textchunks"
if not os.path.exists(output_folder):
    os.makedirs(output_folder)


# Loop through all pdf, txt, tex, ppt, pptx, in the "documents" folder
for filename in os.listdir(filedirectory):
    # Create an empty DataFrame to store the text and title of each document
    df = pd.DataFrame(columns=["Title", "Text"])
    print("Loading " + filename)
    if filename.endswith(".pdf"):
        # Open the PDF file in read-binary mode
        filepath = os.path.join(filedirectory, filename)
        reader = PdfReader(filepath)

        # Extract the text from each page of the PDF
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"

        # Add the text and title to the DataFrame
        title = os.path.splitext(filename)[0]  # Remove the file extension from the filename
        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
        df = pd.concat([df, new_row], ignore_index=True)

    elif filename.endswith(".ppt") or filename.endswith(".pptx"):
        filepath = os.path.join(filedirectory, filename)
        ppt = Presentation(filepath)
        text = ''
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        title = os.path.splitext(filename)[0]  # Remove the file extension from the filename
        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
        df = pd.concat([df, new_row], ignore_index=True)

    elif filename.endswith(".doc") or filename.endswith(".docx"):
        # Open the DOC/DOCX file in binary mode and read the raw data
        filepath = os.path.join(filedirectory, filename)
        doc = docx.Document(filepath)

        # Convert the file to UTF-8 and extract the text
        text = ''
        for paragraph in doc.paragraphs:
            text += paragraph.text

        # Add the text and title to the DataFrame
        title = os.path.splitext(filename)[0]  # Remove the file extension from the filename
        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
        df = pd.concat([df, new_row], ignore_index=True)

    elif filename.endswith(".txt"):
        # Open the text file and read its contents
        filepath = os.path.join(filedirectory, filename)
        with open(filepath, "r", encoding="utf-8") as file:
            text = file.read()

        # Add the text and title to the DataFrame
        title = os.path.splitext(filename)[0]  # Remove the file extension from the filename
        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
        df = pd.concat([df, new_row], ignore_index=True)
        
    elif filename.endswith(".tex"):
        # Use regular expressions to extract regular text from the LaTeX file
        filepath = os.path.join(filedirectory, filename)
        with open(filepath, "r", encoding="utf-8") as file:
            text = file.read()

        # Replace special characters
        text = text.replace('\\$', '$')
        text = text.replace('\\\\', '\n')  # Replace \\ with newline for paragraph breaks
        # Remove comments
        text = re.sub(r'%.*?\n', '', text)

        def replace_math_expression(match):
            # Remove $ or $$ signs but keep the expression
            return match.group(1)

        # Modified regular expression to match both $...$ and $$...$$
        text= re.sub(r'\${1,2}(.*?)\${1,2}', replace_math_expression, text)

        # Remove \begin{} ... \end{} blocks
        text = re.sub(r'\\begin{.*?}.*?\\end{.*?}', '', text, flags=re.DOTALL)

        # Remove common LaTeX commands
        commands = [
            r'\\textbf{.*?}', r'\\textit{.*?}', r'\\emph{.*?}', r'\\underline{.*?}',  # Formatting
            r'\\cite{.*?}', r'\\ref{.*?}',  # References
            r'\\label{.*?}',  # Labels
            # Add more commands as needed
        ]
        for command in commands:
            text = re.sub(command, '', text)
        
        # Add the text and title to the DataFrame
        title = os.path.splitext(filename)[0] # Remove the file extension from the filename
        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
        df = pd.concat([df, new_row], ignore_index=True)

    # Create summaries to append to each chunk of what this text is about
    # Loop through the rows and create overlapping chunks for each text
    chunks = []
    summary_chunks = []
    chunk_counter = 0
    for i, row in df.iterrows():
        # Tokenize the text for the current row
        tokens = nltk.word_tokenize(row['Text'])

        # Loop through the tokens and create overlapping chunks
        for j in range(0, len(tokens), chunk_size - overlap_size):
            # Get the start and end indices of the current chunk
            start = j
            end = j + chunk_size
            # create summaries
            if chunk_counter % 5 == 0:
                # Define the extended chunk range
                extended_start = max(0, start - 500)
                extended_end = min(len(tokens), start + 500)
                # Create the extended chunk
                summary_chunks = ' '.join(tokens[extended_start:extended_end])
                # Apply the summer function to the extended chunk and store the result
                send_to_gpt = []
                response = []
                current_summary = []
                instructions = "Consider this text from portion of a reading, transcript, slides or handout for " + classname + ", a " + classdescription + ".  Give a SHORT ONE SENTENCE summary of what this specific block of text is about, assuming the user already knows the document it comes from and the class is relates to. The format should be a list of NO MORE THAN THREE ideas covered in the block of text, likely for the only time in this class, separated by commas, like 'Context: ...' where again, the response is a SHORT ONE SENTENCE summary, such as 'Context: marginal costs applied to new firms, example of steel' or 'Context: melting point of steel, relation to aluminum, underlying atomic reason'"
                send_to_gpt.append({"role": "system", "content": instructions})
                send_to_gpt.append({"role": "user", "content": summary_chunks})
                response = openai.ChatCompletion.create(
                    messages=send_to_gpt,
                    temperature=0.1,
                    max_tokens=50,
                    model="gpt-3.5-turbo"
                )
                current_summary = response["choices"][0]["message"]["content"]
                print(current_summary)

            chunk_counter += 1

            # Create the current chunk by joining the tokens within the start and end indices
            chunk = ' '.join(tokens[start:end])

            # Add the article title to the beginning of the chunk
            chunk_with_title = "Source: " + row['Title'] + ". " + current_summary + " " + chunk

            # Append the current chunk to the list of chunks, along with the corresponding title
            chunks.append([row['Title'], chunk_with_title])

    # Convert the list of chunks to a dataframe
    df_chunks = pd.DataFrame(chunks, columns=['Title', 'Text'])

    # Truncate the filename if it's too long, e.g., limit to 250 characters
    max_filename_length = 250
    if len(filename) > max_filename_length:
        filename = filename[:max_filename_length]

    # Remove the file extension from the filename
    filename_without_extension = os.path.splitext(filename)[0]

    # Save the df_chunks to the output_folder subfolder with the new file name
    output_file = os.path.join(output_folder, filename_without_extension + "-originaltext.csv")
    df_chunks.to_csv(output_file, encoding='utf-8', escapechar='\\', index=False)

    print("Saving " + filename)

# move files to old directory
destination_directory = '../Already Chopped Documents'
for filename in os.listdir(filedirectory):
    source_path = os.path.join(filedirectory, filename)
    destination_path = os.path.join(destination_directory, filename)

    # Move the file to the destination directory
    shutil.move(source_path, destination_path)
print(f"Moved chopped documents to old directory")





