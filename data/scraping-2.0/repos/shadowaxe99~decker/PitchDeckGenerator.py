import openai
from pptx import Presentation

class PitchDeckGenerator:
    def __init__(self):
        self.presentation = Presentation()

    def generate_content(self, prompt):
        response = openai.Completion.create(engine="text-davinci-003", prompt=prompt, max_tokens=100)
        return response.choices[0].text.strip()

    def create_slide(self, title, content):
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title_placeholder.text = title
        content_placeholder.text = content

    def create_pitch_deck(self, influencer):
        slides_content = [
            ("Introduction", f"Generate an introduction for a pitch deck for {influencer}'s brand."),
            ("Team Overview", f"Generate a team overview for a pitch deck for {influencer}'s brand."),
            ("Problem Statement", f"Generate a problem statement for a pitch deck for {influencer}'s brand."),
            ("Solution", f"Generate a description of the solution for a pitch deck for {influencer}'s brand."),
            ("Market Size", f"Generate a description of the market size for a pitch deck for {influencer}'s brand."),
            ("Business Model", f"Generate a description of the business model for a pitch deck for {influencer}'s brand."),
            ("Marketing Strategy", f"Generate a description of the marketing strategy for a pitch deck for {influencer}'s brand."),
            ("Financial Projections", f"Generate financial projections for a pitch deck for {influencer}'s brand."),
            ("Conclusion", f"Generate a conclusion for a pitch deck for {influencer}'s brand.")
        ]

        for title, prompt in slides_content:
            content = self.generate_content(prompt)
            self.create_slide(title, content)

        self.presentation.save(f'{influencer}_pitch_deck.pptx')

def create_investor_specific_pitch_deck(self, influencer, investor):
        preferences = investor.get_preferences()
        self.create_pitch_deck(influencer, preferences)
    def main(self, brand):
        self.create_pitch_deck(brand)