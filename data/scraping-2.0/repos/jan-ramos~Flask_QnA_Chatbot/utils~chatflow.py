from PyPDF2 import PdfReader
import docx2txt
from pptx import Presentation
from flask import render_template,session
from langchain.memory import ConversationBufferMemory
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.chains import ConversationalRetrievalChain


def document_loader(batch_files):
    text = ""
    for file in batch_files:

        if file.mimetype== 'application/pdf':
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text()
            session['validation-file'] = True

        elif file.mimetype == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            doc =  docx2txt.process(file)
            text += doc
            session['api_key'] = True

        elif file.mimetype == 'text/plain':
            text += str(file.read(),"utf-8")
            session['validation-file'] = True

        elif file.mimetype == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
            session['validation-file'] = True
        else:
            session['validation-file'] = False
            return render_template('upload.html')
    
    return text

def split_text(document):
    text_splitter = RecursiveCharacterTextSplitter(
    chunk_size = 1000,
    chunk_overlap = 200,
    length_function = len)

    pages = text_splitter.split_text(document)

    return pages

def vector_storage(pages,api_key):
    embeddings = OpenAIEmbeddings(openai_api_key=api_key)
    storage = FAISS.from_texts(pages,embeddings)
    
    return storage

def get_conversation_chain(storage,api_key):
    llm = ChatOpenAI(model_name = 'gpt-3.5-turbo',openai_api_key=api_key)
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    chain = ConversationalRetrievalChain.from_llm(llm=llm, retriever=storage.as_retriever(),memory = memory)

    return chain

