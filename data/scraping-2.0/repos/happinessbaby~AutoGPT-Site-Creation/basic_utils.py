# import requests
import os
import fitz
from pathlib import Path
import pypandoc
import uuid
import markdown
import csv
import bs4
import urllib.request
from urllib.request import Request, urlopen
import uuid
from pptx import Presentation
from langchain.document_loaders import UnstructuredURLLoader
from typing import Any, List, Union, Dict
from docxtpl import DocxTemplate
from langchain.document_transformers import Html2TextTransformer
from langchain.document_loaders import AsyncHtmlLoader
import asyncio
from bs4 import BeautifulSoup
from playwright.async_api import async_playwright
    

def convert_to_txt(file, output_path):
    file_ext = os.path.splitext(file)[1]
    if (file_ext)=='.txt':
        os.rename(file, output_path)
    if (file_ext=='.pdf'): 
        convert_pdf_to_txt(file, output_path)
    elif (file_ext=='.odt' or file_ext=='.docx'):
        convert_doc_to_txt(file, output_path)
    elif (file_ext==".log"):
        convert_log_to_txt(file, output_path)
    elif (file_ext==".pptx"):
        convert_pptx_to_txt(file, output_path)
    elif (file_ext==".ipynb"):
        convert_ipynb_to_txt(file, output_path)

def convert_ipynb_to_txt(file, output_path):
    os.rename(file, output_path)

        

def convert_log_to_txt(file, output_path):
    with open(file, "r") as f:
        content = f.read()
        print(content)
        with open(output_path, "w") as f:
            f.write(content)
            f.close()

def convert_pptx_to_txt(pptx_file, output_path):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text+=shape.text+'\n'
    with open(output_path, 'w') as f:
        f.write(text)
        f.close()

#TODO: needs to find the best pdf to txt converter that takes care of special characters best (such as the dash between dates)
def convert_pdf_to_txt(pdf_file, output_path):
    pdf = fitz.open(pdf_file)
    text = ""
    for page in pdf:
        text+=page.get_text() + '\n'
    with open(output_path, 'w') as f:
        f.write(text)
        f.close()

#TODO: needs to find the best docx to txt converter that takes care of special characters best
def convert_doc_to_txt(doc_file, output_path):
    pypandoc.convert_file(doc_file, 'plain', outputfile=output_path)

def read_txt(file):
    try:
        with open(file, 'r', errors='ignore') as f:
            text = f.read()
            return text
    except Exception as e:
        raise e
    
def markdown_table_to_dict(markdown_table):
    # Convert Markdown to HTML
    html = markdown.markdown(markdown_table)
    print(html)

    # Parse HTML table using csv module
    rows = csv.reader(html.split('\n'), delimiter='|')

    # Extract header and data rows
    header = next(rows)
    data = [row for row in rows if len(row) > 1]

    # Convert rows to dictionary
    result = []
    for row in data:
        item = {}
        for i, value in enumerate(row):
            key = header[i].strip()
            item[key] = value.strip()
        result.append(item)

    return result

class AppURLopener(urllib.request.FancyURLopener):
    version = "Mozilla/5.0"

def retrieve_web_content(link, save_path="./web_data/test.txt"):

    req = Request(
        url=link, 
        headers={'User-Agent': 'Mozilla/5.0'}
    )
    try: 
        webpage=str(urllib.request.urlopen(link).read())
    except Exception: 
        # webpage = urlopen(req).read()
        opener = AppURLopener()
        webpage = opener.open(link)
    soup = bs4.BeautifulSoup(webpage, features="lxml")

    content = soup.get_text()

    if content:
        with open(save_path, 'w') as file:
            file.write(content)
            file.close()
            print('Content retrieved and written to file.')
            return True
    else:
        print('Failed to retrieve content from the URL.')
        return False
    
# this one is better than the above function 
def html_to_text(urls:List[str], save_path="./web_data/test.txt"):
    
    loader = AsyncHtmlLoader(urls)
    docs = loader.load()
    html2text = Html2TextTransformer()
    docs_transformed = html2text.transform_documents(docs)
    content = docs_transformed[0].page_content              
    if content:
        with open(save_path, 'w') as file:
            file.write(content)
            file.close()
            print('Content retrieved and written to file.')
            return True
    else:
        print('Failed to retrieve content from the URL.')
        return False



# def write_to_docx_template(doc: Any, field_name: List[str], field_content: Dict[str, str], res_path) -> None:
#     context = {key: None for key in field_name}
#     for field in field_name:
#         if field_content[field] != -1:
#             context[field] = field_content[field]
#     doc.render(context)
#     doc.save(res_path)
#     print(f"Succesfully written {field_name} to {res_path}.")

# source code: https://github.com/trancethehuman/entities-extraction-web-scraper/blob/main/scrape.py
async def ascrape_playwright(url, tags: list[str] = ["h1", "h2", "h3", "span"]) -> str:
    """
    An asynchronous Python function that uses Playwright to scrape
    content from a given URL, extracting specified HTML tags and removing unwanted tags and unnecessary
    lines.
    """
    print("Started scraping...")
    results = ""
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        try:
            page = await browser.new_page()
            await page.goto(url)

            page_source = await page.content()

            # results = remove_unessesary_lines(extract_tags(remove_unwanted_tags(
            #     page_source), tags))
            results = extract_tags(remove_unwanted_tags(
                page_source), tags)
            print("Content scraped")
        except Exception as e:
            results = f"Error: {e}"
        await browser.close()
    # return page_source
    return results 

def remove_unwanted_tags(html_content, unwanted_tags=["script", "style"]):
    """
    This removes unwanted HTML tags from the given HTML content.
    """
    soup = BeautifulSoup(html_content, 'html.parser')

    for tag in unwanted_tags:
        for element in soup.find_all(tag):
            element.decompose()

    return str(soup)


def extract_tags(html_content, tags: list[str]):
    """
    This takes in HTML content and a list of tags, and returns a string
    containing the text content of all elements with those tags, along with their href attribute if the
    tag is an "a" tag.
    """
    soup = BeautifulSoup(html_content, 'html.parser')
    text_parts = []

    for tag in tags:
        elements = soup.find_all(tag)
        for element in elements:
            # If the tag is a link (a tag), append its href as well
            if tag == "a":
                href = element.get('href')
                if href:
                    text_parts.append(f"{element.get_text()} ({href})")
                else:
                    text_parts.append(element.get_text())
            else:
                text_parts.append(element.get_text())

    return ' '.join(text_parts)


def remove_unessesary_lines(content):
    # Split content into lines
    lines = content.split("\n")

    # Strip whitespace for each line
    stripped_lines = [line.strip() for line in lines]

    # Filter out empty lines
    non_empty_lines = [line for line in stripped_lines if line]

    # Remove duplicated lines (while preserving order)
    seen = set()
    deduped_lines = [line for line in non_empty_lines if not (
        line in seen or seen.add(line))]

    # Join the cleaned lines without any separators (remove newlines)
    cleaned_content = " ".join(deduped_lines)

    return cleaned_content     




if __name__=="__main__":
    # retrieve_web_content(
    #     "https://www.google.com/search?q=software+engineer+jobs+in+chicago+&oq=jobs+in+chicago&aqs=chrome.0.0i131i433i512j0i512l9.1387j1j7&sourceid=chrome&ie=UTF-8&ibp=htl;jobs&sa=X&ved=2ahUKEwikxaml1dqBAxULkmoFHRzvD2MQudcGKAF6BAgSEC8&sxsrf=AM9HkKmG-P_UI-ha1ySTJJklAvltPyKEtA:1696363178524#fpstate=tldetail&htivrt=jobs&htidocid=AMKKBKD-6xYovEnvAAAAAA%3D%3D",
    #     save_path = f"./uploads/link/chicago0.txt")
    html_to_text(
        "https://www.betterup.com/blog/skills-for-resume",
        # save_path =f"./uploads/link/software08.txt")
        save_path = f"./web_data/{str(uuid.uuid4())}.txt")





    

    




