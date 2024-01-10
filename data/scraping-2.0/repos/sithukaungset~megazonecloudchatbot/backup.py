from dotenv import load_dotenv
import os
import pandas as pd
import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains import RetrievalQA
from langchain.chat_models import AzureChatOpenAI
import openai
from langchain import LLMChain
from langchain.prompts.chat import (
    ChatPromptTemplate,
    SystemMessagePromptTemplate,
    HumanMessagePromptTemplate,
)
from pptx import Presentation
from langchain.llms import AzureOpenAI
import tiktoken
import sqlite3
import fitz  # PyMuPDF
import pytesseract
from pdfminer.high_level import extract_text
from pdf2image import convert_from_path
import re
#import cv2
import base64
import requests
import json
import sympy as sp
import io
from pptx.util import Inches
from PIL import Image
import time



    

#Powerpoint Processor
class PowerPointProcessor:

    AZURE_ENDPOINT = "https://formtestlsw.cognitiveservices.azure.com/formrecognizer/v2.1/prebuilt/receipt/analyze"
    AZURE_HEADERS = {
        "Ocp-Apim-Subscription-Key" : "2fe1b91a80f94bb2a751f7880f00adf6",
        "Content-Type" : "image/png"
    }

    def ocr_with_azure(self, image):
        """
        Use Azure Form Recognizer to extract text from the given image.

        """
        img_stream = io.BytesIO()
        image.save(img_stream, format='PNG')
        img_bytes = img_stream.getvalue()

        # Make the API request
        response = requests.post(
            self.AZURE_ENDPOINT, headers=self.AZURE_HEADERS, data = img_bytes)
        response_data = response.json()

        # Extract text from the response
        text_data = []
        for page in response_data.get('analyzeResult', {}).get('readResults', []):
            for line in page.get('lines', []):
                text_data.append(line.get("text", ""))

        return "\n".join(text_data)

    def extract_text_from_ppt(self, ppt_stream):
        prs = Presentation(ppt_stream)
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                # Extracting text from shapes
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                # Extracting text from tables
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text += cell.text + "\n"

                # Extracting text from images using Azure Form Recognizer 
                if shape.shape_type == 13: # 13 is the shape type for Picuture
                    img_stream = shape.image.blob
                    img = Image.open(io.BytesIO(img_stream))
                    extracted_text = self.ocr_with_azure(img)
                    text = extract_text + "\n"

        return text.strip() # return text, removing extra spaces

# PDF Processor
class PDFProcessor:
    AZURE_ENDPOINT = "https://formtestlsw.cognitiveservices.azure.com/formrecognizer/v2.1/prebuilt/receipt/analyze"
    AZURE_HEADERS = {
        "Ocp-Apim-Subscription-Key": "2fe1b91a80f94bb2a751f7880f00adf6",
        "Content-Type": "image/png"
    }

    def ocr_with_azure(self, image):
        """
        Use Azure Form Recognizer to extract text from the given image.
        """

        # Convert the image to a PNG format
        _, img_encoded = cv2.imencode('.png', image)
        img_bytes = img_encoded.tobytes()
        
        # Make the API request
        response = requests.post(
            self.AZURE_ENDPOINT, headers=self.AZURE_HEADERS, data = img_bytes)
        response_data = response.json()

        # Extract from the response. Depending on the structure of the response data,
        text_data = []
        for page in response.data.get('analyzeResult', {}).get('readResults', []):
            for line in page.get('lines', []):
                text_data.append(line.get('text', ''))
        
        return "\n".join(text_data)
    
    def ocr_pdf(self, pdf_path):
        """
        Convert a PDF into images and then use Azure to extract text.
        Return the combined text from all pages.
        """
        

        # Convert PDF to a list of images
        images = convert_from_path(pdf_path)

        # OCR each image to extract text using Azure
        texts = [self.ocr_with_azure(img) for img in images]

        # Combine the texts from all pages
        combined_text = "\n".join(texts)

        return combined_text
    
    def extract_text_from_pdf(self, pdf_stream):
        doc = fitz.open(stream=pdf_stream, filetype='pdf')
        text = ""
        for page in doc:
            text += page.get_text("text", clip=page.rect, flags=fitz.TEXT_PRESERVE_LIGATURES)
        return text
    
    def remove_headers_and_footers(self, text):
        # A very basic method: remove the first and last line from each page, assuming they might be headers/footers.
        # This might need more advanced logic, possibly using patterns or machine learning models.
        # Split text into pages
        pages = text.split("\n\n")
        # For each page, remove the first and last lines if they exist
        cleaned_pages = []
        for page in pages:
            lines = page.split('\n')
            # Check if the page has more than 2 lines, if so, remove the first and last lines.
            # Otherwise, just use the lines as they are.
            cleaned_page = lines[1:-1] if len(lines) > 2 else lines
            cleaned_pages.append("\n".join(cleaned_page))

        # Join the cleaned pages back into a single text
        cleaned_text = "\n".join(cleaned_pages)
        return cleaned_text
        

    def enhanced_segment_content(self, text):
        # Define potential section headers and their variations
        sections = {
            "introduction": ["introduction", "intro", "background"],
            "methods": ["methods", "methodology", "experimental", "experiment","materials and methods"],
            "results": ["results","findings", "outcome"],
            "discussion": ["discussion","analysis"],
            "references": ["references", "bibliography","citations"],
            "acknowledgments": ["acknowledgments", "acknowledgement","thanks","gratitude"]
        }

        segments = {key: None for key in sections.keys()}
    
        # Convert the text to lower case for case insensitive search
        lower_text = text.lower()
        
        # For each section, find the starting index using its potential headers
        indices = {}
        for section, patterns in sections.items():
            indices[section] = float('inf')  # initialize with "infinity"
            for pattern in patterns:
                idx = lower_text.find(pattern)
                if idx != -1 and idx < indices[section]:  # Update with the smallest index found
                    indices[section] = idx
                    
        # Sort sections by their starting index
        sorted_sections = sorted(indices.items(), key=lambda x: x[1])
        
        # Extract content for each section based on the detected starting indices
        for i, (section, start_idx) in enumerate(sorted_sections):
            if start_idx == float('inf'):  # If section was not found
                continue
            
            # Set end index to start of next section or end of text
            end_idx = sorted_sections[i+1][1] if i+1 < len(sorted_sections) else len(text)
            segments[section] = text[start_idx:end_idx].strip()
        
        return segments
    

    # Mathematical Section for preprocessing Maths related texts

    def identify_math_expressions(self, text):
        """
        A simple function to identity mathematical expressions.
        This is a very basic way to identify mathematical content using the presence of '='
        Depending on our needs, this might need to be refined.
        """
        math_expressions = []
        for line in text.split('\n'):
            if '=' in line:
                math_expressions.append(line)
            return math_expressions
        
    def process_math_expressions(self, math_expressions):
        """
        Process the identified mathematical expressions using Azure OCR or any other processing methods
        
        1. Evaluation, 2. Simplication, 3. Transcription to LaTex"""
        processed_expressions = {}
    
        for expression in math_expressions:
            print(expression)
            expression_data = {}

            # Evaluation
            try: 
                expression_data["evaluated"] = sp.sympify(expression).evalf()
            except Exception as e:
                expression_data["evaluated"] = str(e)

            # Simplification
            try:
                expression_data["simplified"] = sp.sympify(expression).simplify()
            except Exception as e:
                expression_data["simplified"] = str(e)

            # Transcription to Latex
            try:
                expression_data["latex"] = sp.latex(sp.sympify(expression))
            except Exception as e:
                expression_data["latex"] = str(e)

            processed_expressions[expression] = expression_data
        
        return processed_expressions
    
    # Now when we call 'process math expressions', it will return a dictionary with each original
    # expression as the key.The value for each key will be another dictionary containing the evaluated,
    # simplified and LaTex transcribed results.


    def process_pdf_stream(self, pdf_stream):
        # Extract text from the PDF using fitz
        text = self.extract_text_from_pdf(pdf_stream)

        # Remove headers and footers
        text = self.remove_headers_and_footers(text)

        # Segment content using the enhanced method
        segments =self.enhanced_segment_content(text)

        math_expressions = self.identify_math_expressions(text)
        self.process_math_expressions(math_expressions)

        return segments



        
# Tabular data preprocessing
 
class TabularDataProcessor:
    def __init__(self, data):
        self.data = data

    def preprocess(self):
        self.data = self.data.fillna('Unknown')  # Fill missing values
        # apply the preprocess_math_expression on possible mathematical strings
        self.data = self.data.applymap(self.preprocess_math_expression)
        self.data = self.data.applymap(lambda s: s.lower() if type(
            s) == str else s)  # convert text to lowercase

    # def transform_to_sentences(self):
    #     # Create a list to store the sentences
    #     sentences = []

    #     # Iterate over each row in the DataFrame
    #     for index, row in self.data.iterrows():
    #         # Create a sentence for each row
    #         sentence = ','.join(
    #             [f'{col} is {val}' for col, val in row.items()])
    #         sentences.append(sentence)

    #     return sentences
    
    def transform_to_sentences(self):
        # Create a list to store the sentences
        sentences = []

        # Define a dictionary for column descriptions
        column_descriptions = {
            'age': 'The age of the individual',
            'first_name': 'The first name',
            'last_name': 'The last name',
            # ... add more descriptions as needed
        }

        # Iterate over each row in the DataFrame
        for index, row in self.data.iterrows():
            sentence_parts = []
            
            # Handle specific combined cases
            if 'first_name' in row and 'last_name' in row:
                sentence_parts.append(f"The individual's name is {row['first_name']} {row['last_name']}")
            else:
                for col, val in row.items():
                    if col in column_descriptions:
                        sentence_parts.append(f"{column_descriptions[col]} is {val}")
                    else:
                        # Default behavior if no special description is found
                        sentence_parts.append(f"{col} is {val}")
            
            # Combine all parts for the current row to form a complete sentence
            sentence = '. '.join(sentence_parts)
            sentences.append(sentence)

        return sentences


    def get_num_sheets(self, excel_file):
        # Get the number of sheets in the excel file
        return len(excel_file.sheet_names)

    def get_sheet_names(self, excel_file):
        # Get the name of the sheets in the excel file
        return excel_file.sheet_names

    def process_all_sheets(self, excel_file):
        # Initialize text
        text = ""

        # Iterate over all sheets in the excel file
        for sheet in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet)

            # Process the data in the sheet
            self.data = df
            self.preprocess()
            text += " .".join(self.transform_to_sentences()) + ". "

        return text
    
    

    # MATHPIX_ENDPOINT = "https://api.mathpix.com/v3/text"
    # MATHPIX_HEADERS = {
    #     "app_id": "MY_APP_ID",
    #     "app_key": "MY_APP_KEY",
    #     "Content-type": "application/json"
    # }

    # def ocr_with_mathpix(self, image):
    #     """
    #     Use Mathpix API to extract text from the given image.
    #     """
    #     # Convert the image to a base64 string
    #     _, img_encoded = cv2.imencode('.png', image)
    #     img_str = base64.b64encode(img_encoded).decode('utf-8')

    #     # Create payload for API request
    #     payload = {
    #         "src": f"data:image/png;base64, {img_str}"
    #     }
    #     response = requests.post(
    #         self.MATHPIX_ENDPOINT, headers=self.MATHPIX_HEADERS, data=json.dumps(payload))
    #     response_data = response.json()

    #     # Extract text from the response, handle errors appropriately
    #     return response_data('text', '')

    # def ocr_pdf(self, pdf_path):
    #     """
    #     Convert a PDF into images and then use Mathpix to extract text.
    #     Return the combined text from all pages.
    #     """

    #     # Convert PDF to a list of images
    #     images = convert_from_path(pdf_path)

    #     # OCR each image to extract text using Mathpix
    #     texts = [self.ocr_with_mathpix(img) for img in images]

    #     # Combine the texts from all pages
    #     combined_text = "/n".join(texts)

    #     return combined_text

    # def ocr_pdf(self, pdf_path):
    #     """ 
    #     Convert a PDF into images and then use OCR to extract text.
    #     Return the combined text from all pages.
    #     """

    #     # Convert PDF to a list of images
    #     images = convert_from_path(pdf_path)

    #     # OCR each image to extract text
    #     texts = [pytesseract.image_to_string(img) for img in images]

    #     # Combine the texts from all pages
    #     combined_text = "/n".join(texts)

    #     return combined_text

    # def process_pdf(self, pdf_path):
    #     """
    #     Process a PDF file using OCR and then preprocess any potential mathematical expressions.
    #     """
    #     # Extract text from the PDF using OCR
    #     ocr_text = self.ocr_pdf(pdf_path)

    #     # Split the OCR text into lines and preprocess each line individually
    #     preprocessed_lines = [self.preprocess_math_expression(
    #         line) for line in ocr_text.split("n")]

    #     # Combined the preprocessed lines back into a single string
    #     preprocessed_text = "\n".join(preprocessed_lines)

    #     return preprocessed_text

    # def preprocess_math_expression(self, expression):
    #     # Check if the input is a string, if not, return as is
    #     if not isinstance(expression, str):
    #         return expression

    #     # Remove any extra spaces
    #     expression = expression.lower().strip()

    #     # Add spaces around operators for better tokenization
    #     expression = re.sub(r'(\+|\-|\*|\/|\=|\(|\))', r' \1 ', expression)

    #     # Remove any extra spaces around numbers and variables
    #     expression = re.sub(r'/s+', ' ', expression).strip()

    #     return expression


def translate(text, target_language='ko'):
    # Use the translation API
    # This function should return translated text
    translated_text = text  # replace this with the translation API
    return translated_text


def main():
    # Establish a connection to the database (will create it if it doesn't exist)
    conn = sqlite3.connect('chat_history.db')

    # Create a cursor object
    c = conn.cursor()

    # Create table
    c.execute('''CREATE TABLE IF NOT EXISTS chat_history
                 (question TEXT, answer TEXT)''')

    # Save (commit) the changes
    conn.commit()

    st.set_page_config(page_title="Megazone Cloud ChatBot")
    st.markdown("<h1 style='text-align: center; color: lightgreen;'>Megazone Cloud ChatBot 💬</h1>",
                unsafe_allow_html=True)

    # load environment variables
    load_dotenv()

    OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
    OPENAI_API_BASE = os.getenv("OPENAI_API_BASE")
    # OPENAI_DEPLOYMENT_NAME = os.getenv("OPENAI_DEPLOYMENT_NAME")
    # OPENAI_MODEL_NAME = st.selectbox(
    #   'Select GPT Model', ('GPT35Turbo', 'GPT48K', 'GPT432K'))  # added model selection
    OPENAI_MODEL_NAMES = os.getenv("OPENAI_MODEL_NAMES").split(',')
    OPENAI_DEPLOYMENT_NAMES = os.getenv("OPENAI_DEPLOYMENT_NAMES").split(',')
    OPENAI_EMBEDDING_DEPLOYMENT_NAME = os.getenv(
        "OPENAI_EMBEDDING_DEPLOYMENT_NAME")
    OPENAI_EMBEDDING_MODEL_NAME = os.getenv("OPENAI_EMBEDDING_MODEL_NAME")
    OPENAI_DEPLOYMENT_VERSION = os.getenv("OPENAI_DEPLOYMENT_VERSION")
    OPENAI_MODEL_NAME = st.selectbox('Select GPT Model', OPENAI_MODEL_NAMES)
    OPENAI_DEPLOYMENT_NAME = st.selectbox(
        'Select GPT Deployment name', OPENAI_DEPLOYMENT_NAMES)

    # init Azure OpenAI
    openai.api_type = "azure"
    openai.api_version = OPENAI_DEPLOYMENT_VERSION
    openai.api_base = OPENAI_API_BASE
    openai.api_key = OPENAI_API_KEY
    # init openai
    llm = AzureChatOpenAI(deployment_name=OPENAI_DEPLOYMENT_NAME,
                          model_name=OPENAI_MODEL_NAME,
                          openai_api_base=OPENAI_API_BASE,
                          openai_api_version=OPENAI_DEPLOYMENT_VERSION,
                          openai_api_key=OPENAI_API_KEY)

    embeddings = OpenAIEmbeddings(
        deployment=OPENAI_EMBEDDING_DEPLOYMENT_NAME, model=OPENAI_EMBEDDING_MODEL_NAME, chunk_size=1)

    # Select chat mode
    natural_chat_mode = st.checkbox('Switch to Natural Chat Mode')

    if natural_chat_mode:
        prompt_template = st.text_input("Custom Prompt 🎯:")
        user_input = st.text_input("Type your message here 🤖:")
        # Create a placeholder for the chat history
        chat_placeholder = st.empty()

        # Fetch all records from the database
        c.execute("SELECT * FROM chat_history")
        rows = c.fetchall()

        # Display the chat history
        chat_history = "<h2>Chat History:</h2>"
        for row in rows:
            st.markdown(f"<strong>User :</strong> {row[0]}<br><strong>ChatBot :</strong> {row[1]}<br><br>",
                        unsafe_allow_html=True)

        chat_placeholder.markdown(chat_history, unsafe_allow_html=True)

        if user_input:
            response = openai.ChatCompletion.create(
                engine=OPENAI_DEPLOYMENT_NAME,
                messages=[
                    {"role": "system",
                     # "content": "Assistant is a large language model trained by OpenAI."},
                     "content": prompt_template},
                    {"role": "user", "content": user_input}
                ]
            )
            # st.markdown(
            #     f'### Answer: \n {response["choices"][0]["message"]["content"]}', unsafe_allow_html=True)
            def typewriter_effect(text, delay=0.1):
                typewritten_text = ""
                for char in text:
                    typewritten_text += char
                    time.sleep(delay)
                return typewritten_text
            
            if "choices" in response and response["choices"]:
                response_content = response["choices"][0]["message"]["content"]
                
                # Use the typewriter effect function
                typewritten_response = typewriter_effect(response_content)
                    
                                        
                # Update the chat history with the new message
                chat_history += f"<strong>User :</strong> {user_input}<br><strong>ChatBot :</strong> {response_content}<br><br>"
                chat_placeholder.markdown(chat_history, unsafe_allow_html=True)

                # Insert the question and answer into the database
                c.execute("INSERT INTO chat_history VALUES (?,?)", (user_input, response_content))
                conn.commit()



            if st.button('Translate to Korean'):
                translated_text = translate(result)
            # Insert the question and answer into the database
            c.execute("INSERT INTO chat_history VALUES (?,?)",
                      (user_input, response["choices"][0]["message"]["content"]))

            # Commit the insert
            conn.commit()
            # Update the chat history with the new message
            chat_history = f"<strong>User :</strong> {user_input}<br><strong>ChatBot :</strong> {response['choices'][0]['message']['content']}<br><br>" + chat_history
            chat_placeholder.markdown(chat_history, unsafe_allow_html=True)

    else:

        # Create an instance of the PDFProcessor
        pdf_processor = PDFProcessor()
        # upload file
        uploaded_file = st.file_uploader("Upload your file", type=[
            "pdf", "csv", "txt", "xlsx", "xls", "ppt", "pptx"])

        # extract the text
        if uploaded_file is not None:
            file_details = {"FileName": uploaded_file.name,
                            "FileType": uploaded_file.type, "FileSize": uploaded_file.size}
            st.write(file_details)

            if file_details["FileType"] == "application/pdf":
                with st.spinner('Processing the PDF...'):
                    # Using the process_pdf_stream method to extract and segment text from the PDF
                    segments = pdf_processor.process_pdf_stream(uploaded_file.read())
            
                    # for section, content in segments.items():
                    #     if content:
                    #         st.write(f"{section.capitalize()}:\n{content}\n")
                    text = "\n".join(filter(None, segments.values()))


            elif file_details["FileType"] in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
                with st.spinner('Reading the PowerPoint file...'):
                    # Create an instance of the Powerpoint Processor
                    ppt_processor = PowerPointProcessor()
                    text = ppt_processor.extract_text_from_ppt(uploaded_file)

            elif file_details["FileType"] == "text/plain":
                with st.spinner('Reading the TXT file...'):
                    text = uploaded_file.read().decode("utf-8")

            elif file_details["FileType"] in ["application/vnd.ms-excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"]:
                with st.spinner('Reading the Excel file...'):
                    excel_file = pd.ExcelFile(uploaded_file)

                    # text = process_all_sheets(excel_file)
                    # Create an instance of TabularDataProcessor
                    processor = TabularDataProcessor(None)
                    # df = pd.read_excel(uploaded_file)

                    # Get the number of sheets and their names
                    num_sheets = processor.get_num_sheets(excel_file)
                    sheet_names = processor.get_sheet_names(excel_file)

                    # sheet_names = excel_file.sheet_names

                    st.write(f"Number of sheets: {num_sheets}")
                    st.write(f"Sheet Names: {sheet_names}")

                    text = processor.process_all_sheets(excel_file)

                    # for sheet in sheet_names:
                    #     df = pd.read_excel(excel_file, sheet_name=sheet)
                    #     processor.data = df  # Set the data for the processor
                    #     processor.preprocess()
                    #     # text = " ".join(map(str, df.values))
                    #     text = ". ".join(
                    #         processor.transform_to_sentences()) + ""

            elif file_details["FileType"] == "text/csv":
                with st.spinner('Reading the CSV file...'):
                    
                    df = pd.read_csv(uploaded_file)
                    text = " ".join(map(str, df.values))
            else:
                st.error("File type not supported.")

            # split into chunks
            text_splitter = CharacterTextSplitter(
                separator="\n", chunk_size=1000, chunk_overlap=200, length_function=len)
            chunks = text_splitter.split_text(text)

            # load the faiss vector store we saved into memory
            with st.spinner('Creating knowledge base...'):
                vectorStore = FAISS.from_texts(chunks, embeddings)

            # use the faiss vector store we saved to search the local document
            retriever = vectorStore.as_retriever(
                search_type="similarity", search_kwargs={"k": 2})

            # use the vector store as a retriever
            qa = RetrievalQA.from_chain_type(
                llm=llm, chain_type="stuff", retriever=retriever, return_source_documents=False)

            # show user input
            prompt_template = st.text_input("Custom Prompt 🎯:")
            user_question = st.text_input("Ask a question 🤖:")

            chat_placeholder = st.empty()

            # Fetch all records from the database
            c.execute("SELECT * FROM chat_history")
            rows = c.fetchall()

            # Display the chat history
            chat_history = "<h2>Chat History:</h2>"
            for row in rows:
                st.markdown(f"<strong>User :</strong> {row[0]}<br><strong>Chat Bot :</strong> {row[1]}<br><br>",
                            unsafe_allow_html=True)

            chat_placeholder.markdown(chat_history, unsafe_allow_html=True)

            if user_question:
                result = qa({"query": user_question})
                # Display the result in a more noticeable way
                
                # st.markdown(
                #     f'### Answer: \n {result["result"]}', unsafe_allow_html=True)
                if "choices" in result and result["choices"]:
                    response_content = result["choices"][0]["message"]["content"]
                    response_id = f"response_{hash(response_content)}"

                    # Update the chat history with the new message
                    chat_history += f"<strong>User :</strong> {user_input}<br><strong>ChatBot :</strong> <div id='{response_id}'></div><br><br>"
                    chat_placeholder.markdown(chat_history, unsafe_allow_html=True)

                    # Insert the question and answer into the database
                    c.execute("INSERT INTO chat_history VALUES (?,?)", (user_input, response_content))
                    conn.commit()

                    # Use the typewrite function for the typewriter effect
                    typewritten_response = typewriter_effect(response_content)
                    st.components.v1.html(typewritten_response, height=400, scrolling=True)

                else:
                    st.write("Sorry, I couldn't generate a response for that question.")



                # Insert the question and answer into the database
                c.execute("INSERT INTO chat_history VALUES (?,?)",
                          (user_question, result["result"]))

                # Commit the insert
                conn.commit()
                chat_history = f"<strong>User :</strong> {user_question}<br><strong>ChatBot :</strong> {result['result']}<br><br>" + chat_history
                chat_placeholder.markdown(chat_history, unsafe_allow_html=True)


if __name__ == '__main__':
    main()
