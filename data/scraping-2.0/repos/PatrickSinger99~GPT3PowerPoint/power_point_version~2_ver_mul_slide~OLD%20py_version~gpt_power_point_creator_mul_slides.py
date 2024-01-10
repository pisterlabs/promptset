from pptx import Presentation
from pptx.util import Cm, Pt
import PIL
from pptx.dml.color import RGBColor
from google_images_download import google_images_download
import os
import openai
import shutil

# Data
slide_height = Cm(19.05)
slide_width = Cm(25.4)

# Openai key
with open("openai_key.txt") as file:
    key = file.read()
    openai.api_key = key


def add_picture_bottom_right(slide, img_path, margin=Cm(1), max_pic_height=Cm(14), max_pic_width=Cm(10),
                             slide_height=Cm(19.05), slide_width=Cm(25.4)):

    if not isinstance(margin, Cm):
        margin = Cm(margin)

    if not isinstance(max_pic_height, Cm):
        max_pic_height = Cm(max_pic_height)

    if not isinstance(max_pic_width, Cm):
        max_pic_width = Cm(max_pic_width)

    if not isinstance(slide_height, Cm):
        slide_height = Cm(slide_height)

    if not isinstance(slide_width, Cm):
        slide_width = Cm(slide_width)

    pic_width, pic_height = PIL.Image.open(img_path).size
    pic_width = Cm(round(pic_width / 47.25, 2))
    pic_height = Cm(round(pic_height / 47.25, 2))

    # Convert width if too big
    if pic_width > max_pic_width:
        ratio = max_pic_width / pic_width
        pic_width = max_pic_width
        pic_height *= ratio

    # Convert Height if too big
    if pic_height > max_pic_height:
        ratio = max_pic_height / pic_height
        pic_height = max_pic_height
        pic_width *= ratio

    image = slide.shapes.add_picture(img_path, slide_width - margin - pic_width, slide_height - margin - pic_height,
                                     height=pic_height, width=pic_width)

    return image


def create_power_point_slides_from_gpt(dictionary, img_dict, pp_title):

    # Create File and Slide
    prs = Presentation()
    for key in dictionary:
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Add Title
        title = slide.shapes.title
        title.text = key

        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59, 89, 152)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.name = "Calibri Light"

        # Add Bulletpoints
        subtitle = slide.placeholders[1]
        subtitle.text = dictionary[key].replace(".", ".\n")

        for line in subtitle.text_frame.paragraphs:
            line.font.size = Pt(25)
            line.font.name = "Calibri Light"

        subtitle.width = Cm(12)
        subtitle.height = Cm(19.05 - 5.2)
        subtitle.top = Cm(4.2)
        subtitle.left = Cm(1)

        try:
            img = add_picture_bottom_right(slide, img_dict[key], max_pic_width=Cm(11))
        except:
            pass

    prs.save(pp_title + "_mul_slides.pptx")
