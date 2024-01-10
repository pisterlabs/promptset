import re
import io
import tempfile
import os
import openai
import json
import shutil
import docx
import pptx
# for http parsing
import requests
from bs4 import BeautifulSoup
from pdfminer.converter import TextConverter
from pdfminer.pdfinterp import PDFPageInterpreter
from pdfminer.pdfinterp import PDFResourceManager
from pdfminer.pdfpage import PDFPage
from pdfminer.layout import LAParams
from subprocess import Popen
from subprocess import DEVNULL

laparams = LAParams()
laparams.char_margin = 1
laparams.word_margin = 2

## pdf
def extract_text_by_page(pdf_path):
    with open(pdf_path, 'rb') as fh:
        for page in PDFPage.get_pages(fh,
                                      caching=True,
                                      check_extractable=True):
            resource_manager = PDFResourceManager()
            fake_file_handle = io.StringIO()
            converter = TextConverter(resource_manager, fake_file_handle, laparams=laparams)
            #converter = TextConverter(resource_manager, fake_file_handle, laparams=LAParams())
            page_interpreter = PDFPageInterpreter(resource_manager, converter)
            page_interpreter.process_page(page)

            text = fake_file_handle.getvalue()
            yield text

            # close open handles
            converter.close()
            fake_file_handle.close()


def extract_text(pdf_path):
    text = ""
    for page in extract_text_by_page(pdf_path):
        text += page
    chapters = extract_chapters(text)
    return chapters


# if we have "chapters" headings
#def extract_chapters(text):
#    # split the text by chapter heading
#    chapters = re.split(r'(?i)(chapter|section)\s+\d+', text)
#    return chapters

def extract_chapters(text, max_words=1000):
    # split the text into words
    words = text.split()
    # initialize a list to store the chunks
    chunks = []
    chunk = ""
    for word in words:
        if len(chunk.split()) + len(word.split()) <= max_words:
            chunk += " " + word
        else:
            chunks.append(chunk)
            chunk = word
    chunks.append(chunk)
    return chunks



def create_summary(text, max_tokens_completion = 100, prompt_in = ''):
      prompt = prompt_in + text
      out = openai.Completion.create(
        model="text-davinci-003",
        prompt = prompt,
        max_tokens=max_tokens_completion,
        temperature=0.7
      )
      json_object = json.loads(str(out))
      response_summary = json_object['choices'][0]['text']
      response_summary = response_summary.rstrip(',')
      response_summary = response_summary.rstrip('.')
      return response_summary


def generate_summaries(chapters, min_words_summary = 20):
    summaries = []
    for chapter in chapters:
        summary = create_summary(chapter, max_tokens_completion = 100, prompt_in = 'Create a brief summary from the inpute text in bullet points.' +
                                                                                    'Start every bullet point with "\item ". Do not output incomplete sentences. This is the input text: ')
        summaries.append(summary)
    if len(summaries[-1].split()) <= min_words_summary:
        summaries = summaries[:-1]
    summaries[-1] += '.'
    return summaries

def summarize_pdf(summaries, overall_summary, file_out):
    # create the LaTeX file
    out_file_raw, file_extension = os.path.splitext(file_out)
    pdf_path = os.path.dirname(file_out)
    tex_file = out_file_raw + '.tex'
    with open(tex_file, 'w') as f:
        # write the document preamble
        f.write(r'\documentclass[12pt]{article}')
        f.write('\n')
        f.write(r'\usepackage[utf8]{inputenc}')
        f.write('\n')
        f.write(r'\usepackage{amsmath}')
        f.write('\n')
        f.write(r'\usepackage{amsfonts}')
        f.write('\n')
        f.write(r'\usepackage{amssymb}')
        f.write('\n')
        f.write(r'\usepackage{graphicx}')
        f.write('\n')
        f.write(r'\begin{document}')
        f.write('\n')
        # write the number of chapters
        f.write(r'\section*{{Number of Chapters: {}}}'.format(len(summaries)))
        f.write('\n')
        # write the overall summary
        f.write(r'\subsection*{Overall Summary}')
        f.write('\n')
        f.write(overall_summary)
        f.write('\n')
        f.write('\\newpage')
        f.write(r'\section*{Chapter Summaries}')
        f.write('\n')
        # write the chapter summaries
        for i, summary in enumerate(summaries):
            f.write(r'\subsection*{{Chapter {} Summary}}'.format({i + 1}))
            f.write('\n')
            f.write('\\begin{itemize}')
            f.write('\n')
            f.write(summary.replace("\\item", "\\item "))
            f.write('\end{itemize}')
            f.write('\n')

        # end the document
        f.write(r'\end{document}')
    # compile the LaTeX file into a PDF
    command_args = ' -interaction=nonstopmode -output-directory=' + pdf_path +  ' ' + tex_file
    command = 'pdflatex' + command_args + ' 2> /dev/null'
    os.system(command)
    #Popen(['pdflatex', 'summaries.tex'], stdin=DEVNULL, stdout=DEVNULL, stderr=DEVNULL)
    return file_out, overall_summary

def pdf_to_summary(file_in, file_out):
        chapters = extract_text(file_in)
        summaries = generate_summaries(chapters, min_words_summary = 10)
        overall_summary = create_summary(text = " ".join(summaries).replace('\\item', ''), max_tokens_completion = 400, prompt_in = 'From the given text, generate a concise overall summary: ')
        out = summarize_pdf(summaries, overall_summary, file_out = file_out)
        return out


# summary from text file
def txt_to_summary(file_in, file_out):
        # Read the contents of the file into a single string
        with open(file_in, 'r') as file:
          text = file.read()
        chapters = extract_chapters(text, max_words = 1000)
        summaries = generate_summaries(chapters, min_words_summary = 10)
        overall_summary = create_summary(text = " ".join(summaries).replace('\\item', ''), max_tokens_completion = 400, prompt_in = 'From the given text, generate a concise overall summary: ')
        out = summarize_pdf(summaries, overall_summary, file_out = file_out)
        return out

# summary from word file
def docx_to_summary(file_in, file_out):
        document = docx.Document(file_in)
        # Create an empty string
        text = ""
        # Iterate over the paragraphs in the document
        for paragraph in document.paragraphs:
        # Add the text of each paragraph to the string
          text += paragraph.text
        chapters = extract_chapters(text, max_words = 1000)
        summaries = generate_summaries(chapters, min_words_summary = 10)
        overall_summary = create_summary(text = " ".join(summaries).replace('\\item', ''), max_tokens_completion = 400, prompt_in = 'From the given text, generate a concise overall summary: ')
        out = summarize_pdf(summaries, overall_summary, file_out = file_out)
        return out


# powerpoint
def pptx_to_summary(file_in, file_out):
        presentation = pptx.Presentation(file_in)
        # Create an empty string
        text = ""
        # Iterate over the slides in the presentation
        for slide in presentation.slides:
          # Iterate over the shapes on the slide
          for shape in slide.shapes:
            # Check if the shape is a text box
            if shape.has_text_frame:
              # Add the text of the text box to the string
              text += shape.text
        chapters = extract_chapters(text, max_words = 1000)
        summaries = generate_summaries(chapters, min_words_summary = 10)
        overall_summary = create_summary(text = " ".join(summaries).replace('\\item', ''), max_tokens_completion = 400, prompt_in = 'From the given text, generate a concise overall summary: ')
        out = summarize_pdf(summaries, overall_summary, file_out = file_out)
        return out


# url
def url_to_summary(url_in, file_out):
        # Make an HTTP GET request to the webpage
        response = requests.get(url_in)
        # Parse the HTML of the webpage
        soup = BeautifulSoup(response.text, 'html.parser')
        # Find the text on the webpage
        text = soup.get_text()
        chapters = extract_chapters(text, max_words = 1000)
        summaries = generate_summaries(chapters, min_words_summary = 10)
        overall_summary = create_summary(text = " ".join(summaries).replace('\\item', ''), max_tokens_completion = 400, prompt_in = 'From the given text, generate a concise overall summary: ')
        out = summarize_pdf(summaries, overall_summary, file_out = file_out)
        return out

