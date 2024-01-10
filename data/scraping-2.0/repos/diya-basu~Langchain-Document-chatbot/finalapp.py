import os
import docx2txt
import PyPDF2
import textract
from io import BytesIO
import streamlit as st
from pptx import Presentation
from dotenv import load_dotenv
from langchain import HuggingFaceHub
from nltk.tokenize import sent_tokenize
from langchain.vectorstores import FAISS  
from langchain.embeddings import HuggingFaceEmbeddings 
from langchain.chains.question_answering import load_qa_chain
from langchain.text_splitter import RecursiveCharacterTextSplitter
import easyocr
import cv2
from matplotlib import pyplot as plt
import numpy as np
import tempfile
from langchain.agents import create_csv_agent
from langchain.llms import OpenAI

def extract_text_from_docx(docx_bytes):
    return docx2txt.process(BytesIO(docx_bytes))

def extract_text_from_pdf(pdf_bytes):
    pdf_text = ""
    pdf_reader = PyPDF2.PdfReader(BytesIO(pdf_bytes))
    for page in pdf_reader.pages:
        pdf_text += page.extract_text()
    return pdf_text

def extract_text_from_ppt(ppt_bytes):
    ppt_text = ""
    presentation = Presentation(BytesIO(ppt_bytes))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                ppt_text += shape.text + "\n"
    return ppt_text

def extract_text_from_image(img):
    reader = easyocr.Reader(['en'])
    result = reader.readtext(img)
    img = cv2.imread(img)
    spacer = 100
    detected_sentence=""
    for detection in result: 
        top_left = tuple(detection[0][0])
        bottom_right = tuple(detection[0][2])
        text = detection[1]
        detected_sentence+=text+" "
        spacer+=15
    return detected_sentence
    

def process_uploaded_files(docs):
    all_text = ""
    for doc in docs:
        if doc.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc_text = extract_text_from_docx(doc.read())
        elif doc.type == "application/pdf":
            doc_text = extract_text_from_pdf(doc.read())
        elif doc.type == "application/vnd.ms-powerpoint":
            doc_text = extract_text_from_ppt(doc.read())
        elif doc.type in ("image/jpeg", "image/png", "image/jpg"):
            doc_text = extract_text_from_image(doc.read())
        elif doc.type == "text/csv":
            # Handle CSV processing here if needed
            temp_dir = tempfile.mkdtemp()
            temp_csv_path = os.path.join(temp_dir, "uploaded_csv.csv")
            with open(temp_csv_path, "wb") as f:
               f.write(doc.read())

            agent = create_csv_agent(
            OpenAI(temperature=0), temp_csv_path, verbose=True)

        # Save the uploaded CSV file to the temporary location
        else:
            doc_text = textract.process(doc.read()).decode("utf-8", errors="ignore")


        sentences = sent_tokenize(doc_text)
        all_text += "\n".join(sentences) + "\n"
    return all_text


def get_vectorstore(chunks):
    embeddings = HuggingFaceEmbeddings()

    if not chunks:
        return None
     
    try:
        vectorstore = FAISS.from_documents(chunks, embeddings)
        return vectorstore
    
    except Exception as e:
        return None  


def main():
    load_dotenv()
    os.environ["HUGGINGFACEHUB_API_TOKEN"] = "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"
    st.set_page_config(page_title="Query your PDFs", page_icon=":scroll:")
    st.header("The ultimate PDF whisperer 💬")
    
    # upload files
    pdfs = st.file_uploader("Upload your PDFs", type=["docx", "pdf", "ppt", "txt"], accept_multiple_files=True)
    
    # process each uploaded PDF
    if pdfs is not None:
        text = process_uploaded_files(pdfs)
        sentences = sent_tokenize(text)
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.create_documents(sentences)

        if chunks is not None:        
            db = get_vectorstore(chunks)

        user_question = st.text_input(f"Ask a question about PDF:")

        if user_question:
            docs = db.similarity_search(user_question)
            # llm=HuggingFaceHub(repo_id="google/flan-t5-xxl", model_kwargs={"temperature":1, "max_length":512})
            llm=HuggingFaceHub(repo_id="databricks/dolly-v2-3b", model_kwargs={"temperature":1, "max_length":500})
            chain = load_qa_chain(llm, chain_type="stuff")
            response = chain.run(input_documents=docs, question=user_question)
            st.write(response)

if __name__ == '__main__':
    main()