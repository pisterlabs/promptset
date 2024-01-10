import os
import streamlit as st
import docx2txt
import PyPDF2
import textract
from pptx import Presentation
from io import BytesIO
from dotenv import load_dotenv
from langchain.document_loaders import TextLoader  #for textfiles
from langchain.text_splitter import CharacterTextSplitter #text splitter
from langchain.embeddings import HuggingFaceEmbeddings #for using HugginFace models
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import FAISS  #facebook vectorizationfrom langchain.chains.question_answering import load_qa_chain
from langchain.chains.question_answering import load_qa_chain
from langchain import HuggingFaceHub
from langchain.document_loaders import UnstructuredPDFLoader  #load pdf
from langchain.indexes import VectorstoreIndexCreator #vectorize db index with chromadb
from langchain.chains import RetrievalQA
from langchain.document_loaders import UnstructuredURLLoader  #load urls into docoument-loader
import textwrap
from nltk.tokenize import sent_tokenize
from transformers import AutoTokenizer, AutoModelForQuestionAnswering


def extract_text_from_docx(docx_bytes):
    return docx2txt.process(BytesIO(docx_bytes))

def extract_text_from_pdf(pdf_bytes):
    pdf_text = ""
    pdf_reader = PyPDF2.PdfReader(BytesIO(pdf_bytes))
    for page in pdf_reader.pages:
        pdf_text += page.extract_text()
    return pdf_text

def extract_text_from_ppt(ppt_bytes):
    ppt_text = ""
    presentation = Presentation(BytesIO(ppt_bytes))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                ppt_text += shape.text + "\n"
    return ppt_text

def process_uploaded_files(docs):
    all_text = ""
    for doc in docs:
        if doc.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc_text = extract_text_from_docx(doc.read())
        elif doc.type == "application/pdf":
            doc_text = extract_text_from_pdf(doc.read())
        elif doc.type == "application/vnd.ms-powerpoint":
            doc_text = extract_text_from_ppt(doc.read())
        else:
            doc_text = textract.process(doc.read()).decode("utf-8", errors="ignore")

        sentences = sent_tokenize(doc_text)
        all_text += "\n".join(sentences) + "\n"
    return all_text

tokenizer = AutoTokenizer.from_pretrained("google/flan-t5-xxl")
model = AutoModelForQuestionAnswering.from_pretrained("google/flan-t5-xxl")


def get_vectorstore(chunks):
    embedder=HuggingFaceEmbeddings()

    if not chunks:
        return None
     
    try:
        vectorstore = FAISS.from_documents(chunks, embedder)
        return vectorstore
    
    except Exception as e:
        return None  

def wrap_text_preserve_newlines(text, width=110):
    lines = text.split('\n')
    wrapped_lines = [textwrap.fill(line, width=width) for line in lines]
    wrapped_text = '\n'.join(wrapped_lines)
    return wrapped_text

def main():
    load_dotenv()
    os.environ["HUGGINGFACEHUB_API_TOKEN"] = "Your API Token here"
    st.set_page_config(page_title="Query your PDFs", page_icon=":scroll:")
    st.header("The ultimate PDF whisperer 💬")
    
    # upload files
    pdfs = st.file_uploader("Upload your PDFs", type=["docx", "pdf", "ppt", "txt"], accept_multiple_files=True)
    
    # process each uploaded PDF
    if pdfs is not None:
        text = process_uploaded_files(pdfs)
        sentences = sent_tokenize(text)
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.create_documents(sentences)

        if chunks is not None:        
            db = get_vectorstore(chunks)

        user_question = st.text_input(f"Ask a question about PDF:")

        if user_question:
          
            docs = db.similarity_search(user_question)
            llm=llm=HuggingFaceHub(repo_id="google/flan-t5-xxl", model_kwargs={"temperature":0.5, "max_length":512})
            chain = load_qa_chain(llm, chain_type="stuff")
            response = chain.run(input_documents=docs,question=user_question)
            st.write(response)

if __name__ == '__main__':
    main()
