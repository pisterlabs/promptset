import os
import streamlit as st
import docx2txt
import requests
import PyPDF2
import textract
from pptx import Presentation
from io import BytesIO
from dotenv import load_dotenv
import numpy as np
import time
from langchain.document_loaders import TextLoader  #for textfiles
from langchain.text_splitter import CharacterTextSplitter #text splitter
from langchain.embeddings import HuggingFaceEmbeddings #for using HugginFace models
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import FAISS  #facebook vectorizationfrom langchain.chains.question_answering import load_qa_chain
from langchain.chains.question_answering import load_qa_chain
from langchain import HuggingFaceHub
from langchain.document_loaders import UnstructuredPDFLoader  #load pdf
from langchain.indexes import VectorstoreIndexCreator #vectorize db index with chromadb
from langchain.chains import RetrievalQA
from langchain.document_loaders import UnstructuredURLLoader  #load urls into docoument-loader
import textwrap

RATE_LIMIT = 2
last_request_time = 0

def check_rate_limit():
    global last_request_time
    current_time = time.time()
    time_elapsed = current_time - last_request_time

    if time_elapsed < 60 / RATE_LIMIT:
        return False
    
    last_request_time = current_time
    return True

def extract_text_from_docx(docx_bytes):
    return docx2txt.process(BytesIO(docx_bytes))

def extract_text_from_pdf(pdf_bytes):
    pdf_text = ""
    pdf_reader = PyPDF2.PdfReader(BytesIO(pdf_bytes))
    for page in pdf_reader.pages:
        pdf_text += page.extract_text()
    return pdf_text

def extract_text_from_ppt(ppt_bytes):
    ppt_text = ""
    presentation = Presentation(BytesIO(ppt_bytes))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                ppt_text += shape.text + "\n"
    return ppt_text

def process_uploaded_files(docs):
    all_text = ""
    for doc in docs:
        if doc.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc_text = extract_text_from_docx(doc.read())
        elif doc.type == "application/pdf":
            doc_text = extract_text_from_pdf(doc.read())
        elif doc.type == "application/vnd.ms-powerpoint":
            doc_text = extract_text_from_ppt(doc.read())
        else:
            doc_text = textract.process(doc.read()).decode("utf-8", errors="ignore")
        all_text += doc_text + "\n"  
    return all_text


def get_vectorstore(chunks):
    embeddings = HuggingFaceEmbeddings()

    if not chunks:
        return None
     
    try:
        vectorstore = FAISS.from_documents(chunks, embeddings)
        return vectorstore
    
    except Exception as e:
        return None  

def wrap_text_preserve_newlines(text, width=110):
    lines = text.split('\n')
    wrapped_lines = [textwrap.fill(line, width=width) for line in lines]
    wrapped_text = '\n'.join(wrapped_lines)
    return wrapped_text

def main():
    load_dotenv()
    os.environ["HUGGINGFACEHUB_API_TOKEN"] = "Your API Token here"
    st.set_page_config(page_title="Query your PDFs", page_icon=":scroll:")
    st.header("The ultimate PDF whisperer 💬")
    
    # upload files
    pdfs = st.file_uploader("Upload your PDFs", type=["docx", "pdf", "ppt", "txt"], accept_multiple_files=True)
    
    # process each uploaded PDF
    if pdfs is not None:
        text = process_uploaded_files(pdfs)
        text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=10)
        chunks = text_splitter.create_documents(text)
        chunks = text_splitter.split_documents(chunks)

        if chunks is not None:        
            # knowledge_base = get_vectorstore(chunks)
            # print(chunks)
            db = get_vectorstore(chunks)

        user_question = st.text_input(f"Ask a question about PDF:")

        if user_question:
            docs = db.similarity_search(user_question)
            response = wrap_text_preserve_newlines(str(docs[0].page_content))
            st.write(response)

if __name__ == '__main__':
    main()
