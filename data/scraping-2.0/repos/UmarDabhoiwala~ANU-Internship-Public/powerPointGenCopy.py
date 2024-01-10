import re
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Pt, Cm
from midjourney import imageMaker
import random
from halo import Halo
import openai
import config
import time

openai.api_key = config.OPENAI_API_KEY

spinner = Halo(text='Loading', spinner='dots')


def PowerPointGen (themeId = 0,subject= "", numSlides = 3, imageYesOrNo = False): 
    
    themeId = int(themeId)
    
    systemPrompt = """You are a PowerPoint Generation Assistant, Help Create Slides for a Presentation. Respond only with the relevant information."""
    
    if subject == "":
        initialPrompt = f"Respond only with the relevant information. Create a powerpoint title and {numSlides} relevant sub titles for the slides. Enclose the title and sub titles in double quotes."
    else:
        initialPrompt = f"Respond only with the relevant information. Create a powerpoint title relating to the topic {subject}, and {numSlides} relevant sub titles for the slides. Enclose the title and sub titles in double quotes."
        
       
    chatMSG = [{"role": "system", "content": systemPrompt},
            {"role": "user", "content": initialPrompt}]  
        
    completion = openai.ChatCompletion.create(
        model="gpt-3.5-turbo", 
        messages= chatMSG
    )
    
    response = completion['choices'][0]["message"]["content"]
    
    chatMSG.append({"role": "system", "content": response})
    
    
    
    regex_pattern = r'"([^"]+)"'
    

    lines = response.splitlines()
    slideTitles = []

    for line in lines:
        line = re.findall(regex_pattern, line)
        if len(line) > 0:
            slideTitles.append(line[0])

    themes = ["blue", "dividend", "green", "integral", "ion", "paralax","red","slate","wavy","random"]
    if themeId != 0:
        custom = True
        theme = themes[themeId-1]
        if theme == "random":
            themeRandInt = random.randint(0, (len(themes) -1))
            theme = themes[themeRandInt]
        prs = Presentation(f'powerpointTemplates/{theme}.pptx')
    else: 
        custom = False
        prs = Presentation()
        
        
    Layout = prs.slide_layouts[0] 
    first_slide = prs.slides.add_slide(Layout)
    titleText = first_slide.shapes.title
    titleText.text = slideTitles[0]


    font = titleText.text_frame.paragraphs[0].runs[0].font
    font.size = Pt(42)

    textbox = first_slide.shapes[1]
    sp = textbox.element
    sp.getparent().remove(sp)

    slideTitles.pop(0)
    prs.save("files/example.pptx")
    
    for i, x in enumerate(slideTitles):
        print(x)
        choices = [1, 2, 3]
        if imageYesOrNo: 
            weights = [0.5, 0.3, 0.2]
        else:
            weights = [1, 0, 0]

        selected_choice = random.choices(choices, weights=weights)[0]
        
        if selected_choice == 1:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = x
            slideTextPrompt = f"Expand upon {x} with a small paragraph and a few concise dot points. Do Not go over 100 words. Only respond with relevant information do not respond with {x}, 'paragraph', or any other title"

            chatMSG.append({"role": "user", "content": slideTextPrompt})
        
            spinner.start("Asking ChatGPT for slide content")
            time.sleep(20)
            completion = openai.ChatCompletion.create(
                model="gpt-3.5-turbo", 
                messages= chatMSG
                )
            spinner.succeed("ChatGPT Finished")
        
            response = completion['choices'][0]["message"]["content"]
            chatMSG.append({"role": "system", "content": response})
            response = response.lstrip()
            
            slide.placeholders[1].text = response
            frame = slide.placeholders[1].text_frame
            
            frame.word_wrap = True
            
            for paragraph in frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)
            
            frame.margin_left = Pt(28.35)  # 1 inch
            frame.margin_right = Pt(28.35)  # 1 inch
            frame.margin_top = Pt(28.35)  # 1 inch
            frame.margin_bottom = Pt(28.35)  # 1 inch
                    
            
        elif selected_choice == 2:
            slide = prs.slides.add_slide(prs.slide_layouts[3])
            title = slide.shapes.title
            title.text = x
            slideTextPrompt = f"Expand upon {x} with a few concise dot points, only respond with relevant information do not respond with {x} or any other title. No more than 100 words"
            picturePrompt = f"Write a description for a picture that would be relevant to {x}. Start the prompt with photo of and then concisely describe the picture in less than 30 words"
            
            chatMSG.append({"role": "user", "content": slideTextPrompt})
            spinner.start("Generating Text and Image")
            
            time.sleep(20)
            completion = openai.ChatCompletion.create(
                model="gpt-3.5-turbo", 
                    messages=chatMSG
                )
        
            response = completion['choices'][0]["message"]["content"]
            chatMSG.append({"role": "system", "content": response})
            
            chatMSG.append({"role": "user", "content": picturePrompt})
            time.sleep(20)
            completion = openai.ChatCompletion.create(
                model="gpt-3.5-turbo", 
                    messages=chatMSG
                )
            
            
            spinner.succeed("Generated Text and Image Prompts")
            
            pictureResponse = completion['choices'][0]["message"]["content"]
            chatMSG.append({"role": "system", "content": pictureResponse})
            
            #Adding Text     
            slide.placeholders[1].text = response
            frame = slide.placeholders[1].text_frame
            for paragraph in frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    
                    
            frame.word_wrap = True
            
            frame.margin_left = Pt(28.35)  # 1 inch
            frame.margin_right = Pt(28.35)  # 1 inch
            frame.margin_top = Pt(28.35)  # 1 inch
            frame.margin_bottom = Pt(28.35)  # 1 inch
            
            #Adding Picture
                
            textbox = slide.shapes[2]
            sp = textbox.element
            sp.getparent().remove(sp)
            
            if custom:
                left = Cm(20)
            else: 
                left = Cm(12.67)
            top = Cm(4.5)
            height = Cm(12.5)
            
            imageMaker(pictureResponse, "files/output.png")
            slide.shapes.add_picture("files/output.png", left, top, height=height)
                
        elif selected_choice == 3:
            slide = prs.slides.add_slide(prs.slide_layouts[8])
            picturePrompt = f"Write a description for a picture that would be relevant to {x}. Start the prompt with photo of and then concisely describe the picture in less than 30 words"
            title = slide.shapes.title
            title.text = x
            
            spinner.start("Generating Image Prompt")
            
            chatMSG.append({"role": "user", "content": picturePrompt})
            time.sleep(20)
            completion = openai.ChatCompletion.create(
                model="gpt-3.5-turbo", 
                messages=chatMSG
                )
            spinner.succeed("Generated Image Prompt")
        
            pictureResponse = completion['choices'][0]["message"]["content"]
            chatMSG.append({"role": "system", "content": pictureResponse})
            
            placeholder = slide.placeholders[1]
            imageMaker(pictureResponse, "files/output.png")
            
            picture = placeholder.insert_picture('files/output.png')
        else: 
            print ("error")
        
        prs.save("files/example.pptx")

    #If template    
    if custom:
        xml_slides = prs.slides._sldIdLst  
        slides = list(xml_slides)
        xml_slides.remove(slides[0]) 
        
        prs.save("files/example.pptx")
