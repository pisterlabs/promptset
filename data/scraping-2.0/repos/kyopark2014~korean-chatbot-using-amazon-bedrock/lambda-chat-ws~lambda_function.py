import json
import boto3
import os
import time
import datetime
from io import BytesIO
import PyPDF2
import csv
import traceback
import re
from urllib import parse

from langchain.prompts import PromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from langchain.chains.summarize import load_summarize_chain
from langchain.llms.bedrock import Bedrock
from langchain.chains import ConversationChain
from langchain.memory import ConversationBufferWindowMemory
from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler
from botocore.config import Config

from langchain.vectorstores.faiss import FAISS
from langchain.vectorstores.opensearch_vector_search import OpenSearchVectorSearch
from langchain.embeddings import BedrockEmbeddings
from langchain.chains import RetrievalQA
from langchain.chains import LLMChain
from langchain.retrievers import AmazonKendraRetriever
from langchain.chains import ConversationalRetrievalChain
from multiprocessing import Process, Pipe
from googleapiclient.discovery import build

s3 = boto3.client('s3')
s3_bucket = os.environ.get('s3_bucket') # bucket name
s3_prefix = os.environ.get('s3_prefix')
meta_prefix = "metadata/"
callLogTableName = os.environ.get('callLogTableName')
kendra_region = os.environ.get('kendra_region', 'us-west-2')
profile_of_LLMs = json.loads(os.environ.get('profile_of_LLMs'))
isReady = False   
rag_method = os.environ.get('rag_method', 'RetrievalPrompt') # RetrievalPrompt, RetrievalQA, ConversationalRetrievalChain

opensearch_account = os.environ.get('opensearch_account')
opensearch_passwd = os.environ.get('opensearch_passwd')
enableReference = os.environ.get('enableReference', 'false')
debugMessageMode = os.environ.get('debugMessageMode', 'false')
opensearch_url = os.environ.get('opensearch_url')
path = os.environ.get('path')
doc_prefix = s3_prefix+'/'
speech_prefix = 'speech/'

useParallelUpload = os.environ.get('useParallelUpload', 'true')
useParallelRAG = os.environ.get('useParallelRAG', 'true')
kendraIndex = os.environ.get('kendraIndex')
kendra_method = os.environ.get('kendraMethod')
roleArn = os.environ.get('roleArn')
top_k = int(os.environ.get('numberOfRelevantDocs', '8'))
selected_LLM = 0
capabilities = json.loads(os.environ.get('capabilities'))
print('capabilities: ', capabilities)
MSG_LENGTH = 100
MSG_HISTORY_LENGTH = 20
speech_generation = True
history_length = 0
token_counter_history = 0
allowDualSearch = os.environ.get('allowDualSearch')
allowDualSearchWithMulipleProcessing = True

# google search api
googleApiSecret = os.environ.get('googleApiSecret')
secretsmanager = boto3.client('secretsmanager')
try:
    get_secret_value_response = secretsmanager.get_secret_value(
        SecretId=googleApiSecret
    )
    #print('get_secret_value_response: ', get_secret_value_response)
    secret = json.loads(get_secret_value_response['SecretString'])
    #print('secret: ', secret)
    google_api_key = secret['google_api_key']
    google_cse_id = secret['google_cse_id']
    #print('google_cse_id: ', google_cse_id)    

except Exception as e:
    raise e

# websocket
connection_url = os.environ.get('connection_url')
client = boto3.client('apigatewaymanagementapi', endpoint_url=connection_url)
print('connection_url: ', connection_url)

HUMAN_PROMPT = "\n\nHuman:"
AI_PROMPT = "\n\nAssistant:"
def get_parameter(model_type, maxOutputTokens):
    if model_type=='titan': 
        return {
            "maxTokenCount":1024,
            "stopSequences":[],
            "temperature":0,
            "topP":0.9
        }
    elif model_type=='claude':
        return {
            "max_tokens_to_sample":maxOutputTokens, # 8k    
            "temperature":0.1,
            "top_k":250,
            "top_p":0.9,
            "stop_sequences": [HUMAN_PROMPT]            
        }

map_chain = dict() # For RAG
map_chat = dict() # For general conversation  


# Multi-LLM
def get_llm(profile_of_LLMs, selected_LLM):
    profile = profile_of_LLMs[selected_LLM]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'LLM: {selected_LLM}, bedrock_region: {bedrock_region}, modelId: {modelId}')
        
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }            
        )
    )
    parameters = get_parameter(profile['model_type'], int(profile['maxOutputTokens']))
    # print('parameters: ', parameters)

    # langchain for bedrock
    llm = Bedrock(
        model_id=modelId, 
        client=boto3_bedrock, 
        # streaming=True,
        # callbacks=[StreamingStdOutCallbackHandler()],
        model_kwargs=parameters)
    return llm

def sendMessage(id, body):
    try:
        client.post_to_connection(
            ConnectionId=id, 
            Data=json.dumps(body)
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('err_msg: ', err_msg)
        raise Exception ("Not able to send a message")

def sendResultMessage(connectionId, requestId, msg):    
    result = {
        'request_id': requestId,
        'msg': msg,
        'status': 'completed'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, result)

def sendDebugMessage(connectionId, requestId, msg):
    debugMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'debug'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, debugMsg)

def sendErrorMessage(connectionId, requestId, msg):
    errorMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'error'
    }
    print('error: ', json.dumps(errorMsg))
    sendMessage(connectionId, errorMsg)

def isKorean(text):
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+')
    word_kor = pattern_hangul.search(str(text))
    # print('word_kor: ', word_kor)

    if word_kor and word_kor != 'None':
        print('Korean: ', word_kor)
        return True
    else:
        print('Not Korean: ', word_kor)
        return False

def get_prompt_template(query, conv_type, rag_type):    
    if isKorean(query):
        if conv_type == "normal": # for General Conversation
            prompt_template = """\n\nHuman: 다음의 <history>는 Human과 Assistant의 친근한 이전 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.

            <history>
            {history}
            </history>            

            <question>            
            {input}
            </question>
            
            Assistant:"""
        elif conv_type=='qa' and rag_type=='faiss' and isReady==False: # for General Conversation
            prompt_template = """\n\nHuman: 다음의 <history>는 Human과 Assistant의 친근한 이전 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.

            <history>
            {history}
            </history>

            <question>            
            {input}
            </question>
                    
            Assistant:"""

        elif (conv_type=='qa' and rag_type=='all') or (conv_type=='qa' and rag_type=='opensearch') or (conv_type=='qa' and rag_type=='kendra') or (conv_type=='qa' and rag_type=='faiss' and isReady):  
            # for RAG, context and question
            #prompt_template = """\n\nHuman: 다음의 참고자료(<context>)를 참조하여 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.
            #prompt_template = """\n\nHuman: 참고자료로 부터 구체적인 세부 정보를 충분히 제공합니다. 참고자료는 <context></context> XML tags안에 있습니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.
            prompt_template = """\n\nHuman: 다음의 <context> tag안의 참고자료를 이용하여 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.
            
            <context>
            {context}
            </context>

            <question>            
            {question}
            </question>

            Assistant:"""
                
        elif conv_type == "translation":  # for translation, input
            prompt_template = """\n\nHuman: 다음의 <article>를 English로 번역하세요. 머리말은 건너뛰고 본론으로 바로 들어가주세요. 또한 결과는 <result> tag를 붙여주세요.

            <article>
            {input}
            </article>
                        
            Assistant:"""

        elif conv_type == "sentiment":  # for sentiment, input
            prompt_template = """\n\nHuman: 아래의 <example> review와 Extracted Topic and sentiment 인 <result>가 있습니다.
            <example>
            객실은 작지만 깨끗하고 편안합니다. 프론트 데스크는 정말 분주했고 체크인 줄도 길었지만, 직원들은 프로페셔널하고 매우 유쾌하게 각 사람을 응대했습니다. 우리는 다시 거기에 머물것입니다.
            </example>
            <result>
            청소: 긍정적, 
            서비스: 긍정적
            </result>

            아래의 <review>에 대해서 위의 <result> 예시처럼 Extracted Topic and sentiment 을 만들어 주세요..

            <review>
            {input}
            </review>

            Assistant:"""

        elif conv_type == "extraction":  # information extraction
            prompt_template = """\n\nHuman: 다음 텍스트에서 이메일 주소를 정확하게 복사하여 한 줄에 하나씩 적어주세요. 입력 텍스트에 정확하게 쓰여있는 이메일 주소만 적어주세요. 텍스트에 이메일 주소가 없다면, "N/A"라고 적어주세요. 또한 결과는 <result> tag를 붙여주세요.

            <text>
            {input}
            </text>

            Assistant:"""

        elif conv_type == "pii":  # removing PII(personally identifiable information) containing name, phone number, address
            prompt_template = """\n\nHuman: 아래의 <text>에서 개인식별정보(PII)를 모두 제거하여 외부 계약자와 안전하게 공유할 수 있도록 합니다. 이름, 전화번호, 주소, 이메일을 XXX로 대체합니다. 또한 결과는 <result> tag를 붙여주세요.
            
            <text>
            {input}
            </text>
        
            Assistant:"""

        elif conv_type == "grammar":  # Checking Grammatical Errors
            prompt_template = """\n\nHuman: 다음의 <article>에서 문장의 오류를 찾아서 설명하고, 오류가 수정된 문장을 답변 마지막에 추가하여 주세요.

            <article>
            {input}
            </article>
            
            Assistant: """

        elif conv_type == "step-by-step":  # compelex question 
            prompt_template = """\n\nHuman: 다음은 Human과 Assistant의 친근한 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. 아래 문맥(context)을 참조했음에도 답을 알 수 없다면, 솔직히 모른다고 말합니다. 여기서 Assistant의 이름은 서연입니다.

            {input}

            Assistant: 단계별로 생각할까요?

            Human: 예, 그렇게하세요.
            
            Assistant:"""

        elif conv_type == "like-child":  # Child Conversation (few shot)
            prompt_template = """\n\nHuman: 다음 대화를 완성하기 위해 "A"로 말하는 다음 줄을 작성하세요. Assistant는 유치원 선생님처럼 대화를 합니다.
            
            Q: 이빨 요정은 실제로 있을까?
            A: 물론이죠, 오늘 밤 당신의 이를 감싸서 베개 밑에 넣어두세요. 아침에 뭔가 당신을 기다리고 있을지도 모릅니다.
            Q: {input}

            Assistant:"""      

        elif conv_type == "timestamp-extraction":  # Child Conversation (few shot)
            prompt_template = """\n\nHuman: 아래의 <text>는 시간을 포함한 텍스트입니다. 친절한 AI Assistant로서 시간을 추출하여 아래를 참조하여 <example>과 같이 정리해주세요.
            
            - 년도를 추출해서 <year>/<year>로 넣을것 
            - 월을 추출해서 <month>/<month>로 넣을것
            - 일을 추출해서 <day>/<day>로 넣을것
            - 시간을 추출해서 24H으로 정리해서 <hour>/<hour>에 넣을것
            - 분을 추출해서 <minute>/<minute>로 넣을것

            이때의 예제는 아래와 같습니다.
            <example>
            2022년 11월 3일 18시 26분
            </example>
            <result>
                <year>2022</year>
                <month>11</month>
                <day>03</day>
                <hour>18</hour>
                <minute>26</minute>
            </result>

            결과에 개행문자인 "\n"과 글자 수와 같은 부가정보는 절대 포함하지 마세요.

            <text>
            {input}
            </text>

            Assistant:"""  

        elif conv_type == "funny": # for free conversation
            prompt_template = """\n\nHuman: 다음의 <history>는 Human과 Assistant의 친근한 이전 대화입니다. 모든 대화는 반말로하여야 합니다. Assistant의 이름은 서서이고 10살 여자 어린이 상상력이 풍부하고 재미있는 대화를 합니다. 때로는 바보같은 답변을 해서 재미있게 해줍니다.

            <history>
            {history}
            </history>

            <question>            
            {input}
            </question>
            
            Assistant:"""     

        elif conv_type == "get-weather":  # getting weather (function calling)
            prompt_template = """\n\nHuman: In this environment you have access to a set of tools you can use to answer the user's question.

            You may call them like this. Only invoke one function at a time and wait for the results before invoking another function:
            
            <function_calls>
            <invoke>
            <tool_name>$TOOL_NAME</tool_name>
            <parameters>
            <$PARAMETER_NAME>$PARAMETER_VALUE</$PARAMETER_NAME>
            ...
            </parameters>
            </invoke>
            </function_calls>

            Here are the tools available:
            <tools>
            {tools_string}
            </tools>

            Human:
            {user_input}

            Assistant:"""                  
                
        else:
            prompt_template = """\n\nHuman: 다음은 Human과 Assistant의 친근한 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant는 모르는 질문을 받으면 솔직히 모른다고 말합니다. 여기서 Assistant의 이름은 서연입니다. 
        
            <question>            
            {question}
            </question>

            Assistant:"""

    else:  # English
        if conv_type == "normal": # for General Conversation
            prompt_template = """\n\nHuman: Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor.

            <history>
            {history}
            </history>
            
            <question>            
            {input}
            </question>

            Assistant:"""

        elif conv_type=='qa' and rag_type=='faiss' and isReady==False: # for General Conversation
            prompt_template = """\n\nHuman: Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor.

            <history>
            {history}
            </history>
            
            <question>            
            {input}
            </question>

            Assistant:"""           

        elif (conv_type=='qa' and rag_type=='all') or (conv_type=='qa' and rag_type=='opensearch') or (conv_type=='qa' and rag_type=='kendra') or (conv_type=='qa' and rag_type=='faiss' and isReady):  # for RAG
            prompt_template = """\n\nHuman: Here is pieces of context, contained in <context> tags. Provide a concise answer to the question at the end. If you don't know the answer, just say that you don't know, don't try to make up an answer. 
            
            <context>
            {context}
            </context>

            Go directly into the main points without the preamble. Do not include any additional information like newline characters "\n" or character counts in the result.
                        
            <question>
            {question}
            </question>

            Assistant:"""

        elif conv_type=="translation": 
            prompt_template = """\n\nHuman: Here is an article, contained in <article> tags. Translate the article to Korean. Put it in <result> tags.
            
            <article>
            {input}
            </article>
                        
            Assistant:"""
        
        elif conv_type == "sentiment":  # for sentiment, input
            prompt_template = """\n\nHuman: Here is <example> review and extracted topics and sentiments as <result>.

            <example>
            The room was small but clean and comfortable. The front desk was really busy and the check-in line was long, but the staff were professional and very pleasant with each person they helped. We will stay there again.
            </example>

            <result>
            Cleanliness: Positive, 
            Service: Positive
            </result>

            <review>
            {input}
            </review>
            
            Assistant:"""

        elif conv_type == "pii":  # removing PII(personally identifiable information) containing name, phone number, address
            prompt_template = """\n\nHuman: We want to de-identify some text by removing all personally identifiable information from this text so that it can be shared safely with external contractors.
            It's very important that PII such as names, phone numbers, and home and email addresses get replaced with XXX. Put it in <result> tags.

            Here is the text, inside <text></text> XML tags.

            <text>
            {input}
            </text>

            Assistant:"""

        elif conv_type == "extraction":  # for sentiment, input
            prompt_template = """\n\nHuman: Please precisely copy any email addresses from the following text and then write them, one per line.  Only write an email address if it's precisely spelled out in the input text.  If there are no email addresses in the text, write "N/A".  Do not say anything else.  Put it in <result> tags.

            {input}

            Assistant:"""

        elif conv_type == "grammar":  # Checking Grammatical Errors
            prompt_template = """\n\nHuman: Here is an article, contained in <article> tags:

            <article>
            {input}
            </article>

            Please identify any grammatical errors in the article. Also, add the fixed article at the end of answer.
            
            Assistant: """

        elif conv_type == "step-by-step":  # compelex question 
            prompt_template = """\n\nHuman: Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor.
            
            {input}

            Assistant: Can I think step by step?

            Human: Yes, please do.

            Assistant:"""
        
        elif conv_type == "like-child":  # Child Conversation (few shot)
            prompt_template = """\n\nHuman: Please complete the conversation by writing the next line, speaking as "A". You will be acting as a kindergarten teacher.

            Q: Is the tooth fairy real?
            A: Of course, sweetie. Wrap up your tooth and put it under your pillow tonight. There might be something waiting for you in the morning.
            Q: {input}

            Assistant:"""       

        elif conv_type == "funny": # for free conversation
            prompt_template = """\n\nHuman: 다음의 <history>는 Human과 Assistant의 친근한 이전 대화입니다. Assistant의 이름은 서서이고 10살 여자 어린이입니다. 상상력이 풍부하고 재미있는 대화를 잘합니다. 때론 바보같은 답변을 합니다.

            <history>
            {history}
            </history>

            <question>            
            {input}
            </question>
            
            Assistant:"""     

        else: # normal
            prompt_template = """\n\nHuman: Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor named Seoyeon.

            Human: {input}

            Assistant:"""

            # Use the following pieces of context to provide a concise answer to the question at the end. If you don't know the answer, just say that you don't know, don't try to make up an answer. 
            # The following is a friendly conversation between a human and an AI. The AI is talkative and provides lots of specific details from its context. If the AI does not know the answer to a question, it truthfully says it does not know.
    
    return PromptTemplate.from_template(prompt_template)

def store_document_for_faiss(docs, vectorstore_faiss):
    print('store document into faiss')    
    vectorstore_faiss.add_documents(docs)       
    print('uploaded into faiss')

from opensearchpy import OpenSearch
def delete_index_if_exist(index_name):
    client = OpenSearch(
        hosts = [{
            'host': opensearch_url.replace("https://", ""), 
            'port': 443
        }],
        http_compress = True,
        http_auth=(opensearch_account, opensearch_passwd),
        use_ssl = True,
        verify_certs = True,
        ssl_assert_hostname = False,
        ssl_show_warn = False,
    )

    if client.indices.exists(index_name):
        print('remove index: ', index_name)
        response = client.indices.delete(
            index=index_name
        )
        print('response(remove): ', response)    
    else:
        print('no index: ', index_name)

def store_document_for_opensearch(bedrock_embeddings, docs, userId, documentId):
    #index_name = "rag-index-"+userId+'-'+documentId
    index_name = "rag-index-"+documentId
    #index_name = index_name.replace(' ', '_') # remove spaces
    #index_name = index_name.replace(',', '_') # remove commas
    #index_name = index_name.lower() # change to lowercase
    print('index_name: ', index_name)
    
    delete_index_if_exist(index_name)

    try:
        new_vectorstore = OpenSearchVectorSearch(
            index_name=index_name,  
            is_aoss = False,
            #engine="faiss",  # default: nmslib
            embedding_function = bedrock_embeddings,
            opensearch_url = opensearch_url,
            http_auth=(opensearch_account, opensearch_passwd),
        )
        response = new_vectorstore.add_documents(docs)
        print('response of adding documents: ', response)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                
        raise Exception ("Not able to request to LLM")

    print('uploaded into opensearch')
    
# store document into Kendra
def store_document_for_kendra(path, doc_prefix, s3_file_name, documentId):
    print('store document into kendra')
    encoded_name = parse.quote(s3_file_name)
    source_uri = path+doc_prefix+encoded_name    
    #print('source_uri: ', source_uri)
    ext = (s3_file_name[s3_file_name.rfind('.')+1:len(s3_file_name)]).upper()
    print('ext: ', ext)

    # PLAIN_TEXT, XSLT, MS_WORD, RTF, CSV, JSON, HTML, PDF, PPT, MD, XML, MS_EXCEL
    if(ext == 'PPTX'):
        file_type = 'PPT'
    elif(ext == 'TXT'):
        file_type = 'PLAIN_TEXT'         
    elif(ext == 'XLS' or ext == 'XLSX'):
        file_type = 'MS_EXCEL'      
    elif(ext == 'DOC' or ext == 'DOCX'):
        file_type = 'MS_WORD'
    else:
        file_type = ext

    kendra_client = boto3.client(
        service_name='kendra', 
        region_name=kendra_region,
        config = Config(
            retries=dict(
                max_attempts=10
            )
        )
    )

    documents = [
        {
            "Id": documentId,
            "Title": s3_file_name,
            "S3Path": {
                "Bucket": s3_bucket,
                "Key": s3_prefix+'/'+s3_file_name
            },
            "Attributes": [
                {
                    "Key": '_source_uri',
                    'Value': {
                        'StringValue': source_uri
                    }
                },
                {
                    "Key": '_language_code',
                    'Value': {
                        'StringValue': "ko"
                    }
                },
            ],
            "ContentType": file_type
        }
    ]
    print('document info: ', documents)

    result = kendra_client.batch_put_document(
        IndexId = kendraIndex,
        RoleArn = roleArn,
        Documents = documents       
    )
    # print('batch_put_document(kendra): ', result)
    print('uploaded into kendra')

# load documents from s3 for pdf and txt
def load_document(file_type, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)
    
    if file_type == 'pdf':
        Byte_contents = doc.get()['Body'].read()
        reader = PyPDF2.PdfReader(BytesIO(Byte_contents))
        
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text())
        contents = '\n'.join(texts)
        
    elif file_type == 'pptx':
        Byte_contents = doc.get()['Body'].read()
            
        from pptx import Presentation
        prs = Presentation(BytesIO(Byte_contents))

        texts = []
        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = text + shape.text
            texts.append(text)
        contents = '\n'.join(texts)
        
    elif file_type == 'txt':        
        contents = doc.get()['Body'].read().decode('utf-8')

    elif file_type == 'docx':
        Byte_contents = doc.get()['Body'].read()
            
        import docx
        doc_contents =docx.Document(BytesIO(Byte_contents))

        texts = []
        for i, para in enumerate(doc_contents.paragraphs):
            if(para.text):
                texts.append(para.text)
                # print(f"{i}: {para.text}")        
        contents = '\n'.join(texts)
            
    # print('contents: ', contents)
    new_contents = str(contents).replace("\n"," ") 
    print('length: ', len(new_contents))

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100,
        separators=["\n\n", "\n", ".", " ", ""],
        length_function = len,
    ) 

    texts = text_splitter.split_text(new_contents) 
                
    return texts

# load csv documents from s3
def load_csv_document(path, doc_prefix, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)

    lines = doc.get()['Body'].read().decode('utf-8').split('\n')   # read csv per line
    print('lins: ', len(lines))
        
    columns = lines[0].split(',')  # get columns
    #columns = ["Category", "Information"]  
    #columns_to_metadata = ["type","Source"]
    print('columns: ', columns)
    
    docs = []
    n = 0
    for row in csv.DictReader(lines, delimiter=',',quotechar='"'):
        # print('row: ', row)
        #to_metadata = {col: row[col] for col in columns_to_metadata if col in row}
        values = {k: row[k] for k in columns if k in row}
        content = "\n".join(f"{k.strip()}: {v.strip()}" for k, v in values.items())
        doc = Document(
            page_content=content,
            metadata={
                'name': s3_file_name,
                'page': n+1,
                'uri': path+doc_prefix+parse.quote(s3_file_name)
            }
            #metadata=to_metadata
        )
        docs.append(doc)
        n = n+1
    print('docs[0]: ', docs[0])

    return docs

def get_summary(llm, texts):    
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+') 
    word_kor = pattern_hangul.search(str(texts))
    print('word_kor: ', word_kor)
    
    if word_kor:
        #prompt_template = """\n\nHuman: 다음 텍스트를 간결하게 요약하세오. 텍스트의 요점을 다루는 글머리 기호로 응답을 반환합니다.
        #prompt_template = """\n\nHuman: 아래 <text>는 문서에서 추출한 텍스트입니다. 친절한 AI Assistant로서 아래와 같이 정리해주세요.
        
        #- 50자 미안의 제목을 <title>Name: </title> 안에 넣을것
        #- 300자 미안의 설명을 <description>설명: </description> 안에 넣을것
        #- 500자 미만의 내용 요약을 <summarization>요약: </summarization> 안에 넣을것
        #- 10자 미안의 애용과 과련된 테그 5개를 <tag></tag> 테그 안에 생성할 것

        #모든 생성 결과는 한국어로 해주세요. 결과에 개행문자인 "\m"과 글자 수와 같은 부가정보는 절대 포함하지 마세요.
        #생성이 어렵거나 해당 내용을 모르는 경우 "None"로 결과를 생성하세요.

        prompt_template = """\n\nHuman: 다음 텍스트를 요약해서 500자 이내로 설명하세오.
        
        <text>
        {text}
        </text
        
        Assistant:"""        
    else:         
        prompt_template = """\n\nHuman: Write a concise summary of the following:

        {text}
        
        Assistant:"""
    
    PROMPT = PromptTemplate(template=prompt_template, input_variables=["text"])
    chain = load_summarize_chain(llm, chain_type="stuff", prompt=PROMPT)

    docs = [
        Document(
            page_content=t
        ) for t in texts[:5]
    ]
    summary = chain.run(docs)
    print('summary: ', summary)

    if summary == '':  # error notification
        summary = 'Fail to summarize the document. Try agan...'
        return summary
    else:
        # return summary[1:len(summary)-1]   
        return summary
    
def load_chat_history(userId, allowTime, conv_type, rag_type):
    dynamodb_client = boto3.client('dynamodb')

    response = dynamodb_client.query(
        TableName=callLogTableName,
        KeyConditionExpression='user_id = :userId AND request_time > :allowTime',
        ExpressionAttributeValues={
            ':userId': {'S': userId},
            ':allowTime': {'S': allowTime}
        }
    )
    # print('query result: ', response['Items'])

    for item in response['Items']:
        text = item['body']['S']
        msg = item['msg']['S']
        type = item['type']['S']

        if type == 'text':
            memory_chain.chat_memory.add_user_message(text)
            if len(msg) > MSG_LENGTH:
                memory_chain.chat_memory.add_ai_message(msg[:MSG_LENGTH])                          
            else:
                memory_chain.chat_memory.add_ai_message(msg) 
                
            if len(msg) > MSG_LENGTH:
                memory_chat.save_context({"input": text}, {"output": msg[:MSG_LENGTH]})
            else:
                memory_chat.save_context({"input": text}, {"output": msg})
                
def getAllowTime():
    d = datetime.datetime.now() - datetime.timedelta(days = 2)
    timeStr = str(d)[0:19]
    print('allow time: ',timeStr)

    return timeStr

def isTyping(connectionId, requestId):    
    msg_proceeding = {
        'request_id': requestId,
        'msg': 'Proceeding...',
        'status': 'istyping'
    }
    #print('result: ', json.dumps(result))
    sendMessage(connectionId, msg_proceeding)

def readStreamMsg(connectionId, requestId, stream):
    msg = ""
    if stream:
        for event in stream:
            #print('event: ', event)
            msg = msg + event

            result = {
                'request_id': requestId,
                'msg': msg,
                'status': 'proceeding'
            }
            #print('result: ', json.dumps(result))
            sendMessage(connectionId, result)
    # print('msg: ', msg)
    return msg

_ROLE_MAP = {"human": "\n\nHuman: ", "ai": "\n\nAssistant: "}
def extract_chat_history_from_memory():
    chat_history = []
    chats = memory_chain.load_memory_variables({})    
    # print('chats: ', chats)

    for dialogue_turn in chats['chat_history']:
        role_prefix = _ROLE_MAP.get(dialogue_turn.type, f"{dialogue_turn.type}: ")
        history = f"{role_prefix[2:]}{dialogue_turn.content}"
        if len(history)>MSG_LENGTH:
            chat_history.append(history[:MSG_LENGTH])
        else:
            chat_history.append(history)

    return chat_history

def get_revised_question(llm, connectionId, requestId, query):    
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+')
    word_kor = pattern_hangul.search(str(query))
    print('word_kor: ', word_kor)

    if word_kor and word_kor != 'None':
        #condense_template = """{chat_history}

        #Human: 이전 대화와 다음의 <question>을 이용하여, 새로운 질문을 생성하여 질문만 전달합니다.

        #<question>            
        #{question}
        #</question>
            
        #Assistant: 새로운 질문:"""

        condense_template = """
        <history>
        {chat_history}
        </history>

        Human: <history>를 참조하여, 다음의 <question>의 뜻을 명확히 하는 새로운 질문을 한국어로 생성하세요. 새로운 질문은 원래 질문의 중요한 단어를 반드시 포함합니다.

        <question>            
        {question}
        </question>
            
        Assistant: 새로운 질문:"""
    else: 
        #condense_template = """{chat_history}    
        #Answer only with the new question.

        #Human: How would you ask the question considering the previous conversation: {question}

        #Assistant: Standalone question:"""


        #Given the following <history> and a follow up question, rephrase the follow up question to be a standalone question, in its original language. Answer only with the new question.


        condense_template = """
        <history>
        {chat_history}
        </history>
        Answer only with the new question.

        Human: using <history>, rephrase the follow up <question> to be a standalone question. The standalone question must have main words of the original question.
         
        <quesion>
        {question}
        </question>

        Assistant: Standalone question:"""

        #Given the following <history> and a follow up question, rephrase the follow up question to be a standalone question, in its original language. Answer only with the new question, in its original language. Answer only with the new question.

    print('condense_template: ', condense_template)

    print('start prompt!')

    condense_prompt_claude = PromptTemplate.from_template(condense_template)
        
    condense_prompt_chain = LLMChain(llm=llm, prompt=condense_prompt_claude)

    chat_history = extract_chat_history_from_memory()
    try:         
        revised_question = condense_prompt_chain.run({"chat_history": chat_history, "question": query})
        print('revised_question: '+revised_question)
        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                

        sendErrorMessage(connectionId, requestId, err_msg)        
        raise Exception ("Not able to request to LLM")
    
    if debugMessageMode=='true':
        # sendDebugMessage(connectionId, requestId, '[Debug]: '+revised_question)
        debug_msg_for_revised_question(llm, revised_question.replace("\n"," "), chat_history, connectionId, requestId)
    
    return revised_question.replace("\n"," ")

kendraRetriever = AmazonKendraRetriever(
    index_id=kendraIndex, 
    top_k=top_k, 
    region_name=kendra_region,
    attribute_filter = {
        "EqualsTo": {      
            "Key": "_language_code",
            "Value": {
                "StringValue": "ko"
            }
        },
    },
)

def retrieve_from_kendra(query, top_k):
    if kendra_method == 'kendra_retriever':
        relevant_docs = retrieve_from_kendra_using_kendra_retriever(query, top_k)
    else: 
        relevant_docs = retrieve_from_kendra_using_custom_retriever(query, top_k)
    
    return relevant_docs

def retrieve_from_kendra_using_kendra_retriever(query, top_k):
    print(f"query: {query} (kendra)")

    relevant_docs = []
    relevant_documents = kendraRetriever.get_relevant_documents(
        query=query,
        top_k=top_k,
    )
    #print('length of relevant_documents: ', len(relevant_documents))
    #print('relevant_documents: ', relevant_documents)

    rag_type = "kendra"
    api_type = "kendraRetriever"

    for i, document in enumerate(relevant_documents):
        #print('document.page_content:', document.page_content)
        #print('document.metadata:', document.metadata)
        print(f'## Document(retrieve_from_kendra_using_kendra_retriever) {i+1}: {document}')

        result_id = document.metadata['result_id']
        document_id = document.metadata['document_id']
        # source = document.metadata['source']
        title = document.metadata['title']
        excerpt = document.metadata['excerpt']

        uri = ""
        if "_source_uri" in document.metadata['document_attributes']:
            uri = document.metadata['document_attributes']['_source_uri']

        page = ""
        if "_excerpt_page_number" in document.metadata['document_attributes']:            
            page = document.metadata['document_attributes']['_excerpt_page_number']

        confidence = ""
        assessed_score = ""
            
        if page:
            doc_info = {
                "rag_type": rag_type,
                "api_type": api_type,
                "confidence": confidence,
                "metadata": {
                    #"type": query_result_type,
                    "document_id": document_id,
                    "source": uri,
                    "title": title,
                    "excerpt": excerpt,
                    "translated_excerpt": "",                    
                    "document_attributes": {
                        "_excerpt_page_number": page
                    }
                },
                #"query_id": query_id,
                #"feedback_token": feedback_token
                "assessed_score": assessed_score,
                "result_id": result_id
            }

        else: 
            doc_info = {
                "rag_type": rag_type,
                "api_type": api_type,
                "confidence": confidence,
                "metadata": {
                    #"type": query_result_type,
                    "document_id": document_id,
                    "source": uri,
                    "title": title,
                    "excerpt": excerpt,
                    "translated_excerpt": ""
                },
                #"query_id": query_id,
                #"feedback_token": feedback_token
                "assessed_score": assessed_score,
                "result_id": result_id
            }
            
        relevant_docs.append(doc_info)
    
    return relevant_docs    
    
def retrieve_from_kendra_using_custom_retriever(query, top_k):
    print(f"query: {query} (kendra)")

    index_id = kendraIndex    
    
    kendra_client = boto3.client(
        service_name='kendra', 
        region_name=kendra_region,
        config = Config(
            retries=dict(
                max_attempts=10
            )
        )
    )

    try:
        resp =  kendra_client.retrieve(
            IndexId = index_id,
            QueryText = query,
            PageSize = top_k,      
            AttributeFilter = {
                "EqualsTo": {      
                    "Key": "_language_code",
                    "Value": {
                        "StringValue": "ko"
                    }
                },
            },      
        )
        # print('retrieve resp:', json.dumps(resp))
        query_id = resp["QueryId"]

        if len(resp["ResultItems"]) >= 1:
            relevant_docs = []
            retrieve_docs = []
            for query_result in resp["ResultItems"]:
                #confidence = query_result["ScoreAttributes"]['ScoreConfidence']
                #if confidence == 'VERY_HIGH' or confidence == 'HIGH' or confidence == 'MEDIUM': only for "en"
                retrieve_docs.append(extract_relevant_doc_for_kendra(query_id=query_id, api_type="retrieve", query_result=query_result))
                # print('retrieve_docs: ', retrieve_docs)

            print('Looking for FAQ...')
            try:
                resp =  kendra_client.query(
                    IndexId = index_id,
                    QueryText = query,
                    PageSize = 4, # Maximum number of results returned for FAQ = 4 (default)
                    QueryResultTypeFilter = "QUESTION_ANSWER",  # 'QUESTION_ANSWER', 'ANSWER', "DOCUMENT"
                    AttributeFilter = {
                        "EqualsTo": {      
                            "Key": "_language_code",
                            "Value": {
                                "StringValue": "ko"
                            }
                        },
                    },      
                )
                print('query resp:', json.dumps(resp))
                query_id = resp["QueryId"]

                if len(resp["ResultItems"]) >= 1:
                    
                    for query_result in resp["ResultItems"]:
                        confidence = query_result["ScoreAttributes"]['ScoreConfidence']

                        #if confidence == 'VERY_HIGH' or confidence == 'HIGH' or confidence == 'MEDIUM': 
                        if confidence == 'VERY_HIGH' or confidence == 'HIGH': 
                            relevant_docs.append(extract_relevant_doc_for_kendra(query_id=query_id, api_type="query", query_result=query_result))

                            if len(relevant_docs) >= top_k:
                                break
                    # print('relevant_docs: ', relevant_docs)

                else: 
                    print('No result for FAQ')

            except Exception:
                err_msg = traceback.format_exc()
                print('error message: ', err_msg)
                raise Exception ("Not able to query from Kendra")

            for doc in retrieve_docs:                
                if len(relevant_docs) >= top_k:
                    break
                else:
                    relevant_docs.append(doc)
            
        else:  # fallback using query API
            print('No result for Retrieve API!')
            try:
                resp =  kendra_client.query(
                    IndexId = index_id,
                    QueryText = query,
                    PageSize = top_k,
                    #QueryResultTypeFilter = "DOCUMENT",  # 'QUESTION_ANSWER', 'ANSWER', "DOCUMENT"
                    AttributeFilter = {
                        "EqualsTo": {      
                            "Key": "_language_code",
                            "Value": {
                                "StringValue": "ko"
                            }
                        },
                    },      
                )
                print('query resp:', resp)
                query_id = resp["QueryId"]

                if len(resp["ResultItems"]) >= 1:
                    relevant_docs = []
                    for query_result in resp["ResultItems"]:
                        confidence = query_result["ScoreAttributes"]['ScoreConfidence']

                        if confidence == 'VERY_HIGH' or confidence == 'HIGH' or confidence == 'MEDIUM': 
                            relevant_docs.append(extract_relevant_doc_for_kendra(query_id=query_id, api_type="query", query_result=query_result))

                            if len(relevant_docs) >= top_k:
                                break
                    # print('relevant_docs: ', relevant_docs)

                else: 
                    print('No result for Query API. Finally, no relevant docs!')
                    relevant_docs = []

            except Exception:
                err_msg = traceback.format_exc()
                print('error message: ', err_msg)
                raise Exception ("Not able to query from Kendra")                

    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to retrieve from Kendra")     

    for i, rel_doc in enumerate(relevant_docs):
        print(f'## Document(retrieve_from_kendra_using_kendra_retriever) {i+1}: {json.dumps(rel_doc)}')  

    return relevant_docs

def priority_search(query, relevant_docs, bedrock_embeddings):
    excerpts = []
    for i, doc in enumerate(relevant_docs):
        # print('doc: ', doc)
        if doc['metadata']['translated_excerpt']:
            content = doc['metadata']['translated_excerpt']
        else:
            content = doc['metadata']['excerpt']
        
        excerpts.append(
            Document(
                page_content=content,
                metadata={
                    'name': doc['metadata']['title'],
                    'order':i,
                }
            )
        )  
    # print('excerpts: ', excerpts)

    embeddings = bedrock_embeddings
    vectorstore_confidence = FAISS.from_documents(
        excerpts,  # documents
        embeddings  # embeddings
    )            
    rel_documents = vectorstore_confidence.similarity_search_with_score(
        query=query,
        k=top_k
    )

    docs = []
    for i, document in enumerate(rel_documents):
        print(f'## Document(priority_search) {i+1}: {document}')

        order = document[0].metadata['order']
        name = document[0].metadata['name']
        assessed_score = document[1]
        print(f"{order} {name}: {assessed_score}")

        relevant_docs[order]['assessed_score'] = int(assessed_score)

        if assessed_score < 200:
            docs.append(relevant_docs[order])    
    # print('selected docs: ', docs)

    return docs

def extract_relevant_doc_for_kendra(query_id, api_type, query_result):
    rag_type = "kendra"
    if(api_type == 'retrieve'): # retrieve API
        excerpt = query_result["Content"]
        confidence = query_result["ScoreAttributes"]['ScoreConfidence']
        document_id = query_result["DocumentId"] 
        document_title = query_result["DocumentTitle"]
        
        document_uri = ""
        document_attributes = query_result["DocumentAttributes"]
        for attribute in document_attributes:
            if attribute["Key"] == "_source_uri":
                document_uri = str(attribute["Value"]["StringValue"])        
        if document_uri=="":  
            document_uri = query_result["DocumentURI"]

        doc_info = {
            "rag_type": rag_type,
            "api_type": api_type,
            "confidence": confidence,
            "metadata": {
                "document_id": document_id,
                "source": document_uri,
                "title": document_title,
                "excerpt": excerpt,
                "translated_excerpt": ""
            },
            "assessed_score": "",
        }
            
    else: # query API
        query_result_type = query_result["Type"]
        confidence = query_result["ScoreAttributes"]['ScoreConfidence']
        document_id = query_result["DocumentId"] 
        document_title = ""
        if "Text" in query_result["DocumentTitle"]:
            document_title = query_result["DocumentTitle"]["Text"]
        document_uri = query_result["DocumentURI"]
        feedback_token = query_result["FeedbackToken"] 

        page = ""
        document_attributes = query_result["DocumentAttributes"]
        for attribute in document_attributes:
            if attribute["Key"] == "_excerpt_page_number":
                page = str(attribute["Value"]["LongValue"])

        if query_result_type == "QUESTION_ANSWER":
            question_text = ""
            additional_attributes = query_result["AdditionalAttributes"]
            for attribute in additional_attributes:
                if attribute["Key"] == "QuestionText":
                    question_text = str(attribute["Value"]["TextWithHighlightsValue"]["Text"])
            answer = query_result["DocumentExcerpt"]["Text"]
            excerpt = f"{question_text} {answer}"
            excerpt = excerpt.replace("\n"," ") 
        else: 
            excerpt = query_result["DocumentExcerpt"]["Text"]

        if page:
            doc_info = {
                "rag_type": rag_type,
                "api_type": api_type,
                "confidence": confidence,
                "metadata": {
                    "type": query_result_type,
                    "document_id": document_id,
                    "source": document_uri,
                    "title": document_title,
                    "excerpt": excerpt,
                    "translated_excerpt": "",
                    "document_attributes": {
                        "_excerpt_page_number": page
                    }
                },
                "assessed_score": "",
                "query_id": query_id,
                "feedback_token": feedback_token
            }
        else: 
            doc_info = {
                "rag_type": rag_type,
                "api_type": api_type,
                "confidence": confidence,
                "metadata": {
                    "type": query_result_type,
                    "document_id": document_id,
                    "source": document_uri,
                    "title": document_title,
                    "excerpt": excerpt,
                    "translated_excerpt": "",
                },
                "assessed_score": "",
                "query_id": query_id,
                "feedback_token": feedback_token
            }
    return doc_info

def get_reference(docs, rag_method, rag_type, path, doc_prefix):
    if rag_method == 'RetrievalQA' or rag_method == 'ConversationalRetrievalChain':
        if rag_type == 'kendra':
            reference = "\n\nFrom\n"
            for i, doc in enumerate(docs):
                name = doc.metadata['title']     

                uri = ""
                if ("document_attributes" in doc.metadata) and ("_source_uri" in doc.metadata['document_attributes']):
                    uri = doc.metadata['document_attributes']['_source_uri']
                                    
                if ("document_attributes" in doc.metadata) and ("_excerpt_page_number" in doc.metadata['document_attributes']):
                    page = doc.metadata['document_attributes']['_excerpt_page_number']
                    reference = reference + f'{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>\n'
                else:
                    reference = reference + f'{i+1}. <a href={uri} target=_blank>{name}</a>\n'
        else:
            reference = "\n\nFrom\n"
            for i, doc in enumerate(docs):
                print(f'## Document(get_reference) {i+1}: {doc}')

                name = doc.metadata['name']
                page = doc.metadata['page']
                uri = doc.metadata['uri']

                reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>\n"

    elif rag_method == 'RetrievalPrompt':
        reference = "\n\nFrom\n"
        for i, doc in enumerate(docs):
            if doc['metadata']['translated_excerpt']:
                excerpt = str(doc['metadata']['excerpt']+'  [번역]'+doc['metadata']['translated_excerpt']).replace('"'," ") 
            else:
                excerpt = str(doc['metadata']['excerpt']).replace('"'," ")

            if doc['rag_type'] == 'kendra':                
                if doc['api_type'] == 'kendraRetriever': # provided by kendraRetriever from langchain
                    name = doc['metadata']['title']
                    uri = doc['metadata']['source']
                    reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
                elif doc['api_type'] == 'retrieve': # Retrieve. socre of confidence is only avaialbe for English
                    uri = doc['metadata']['source']
                    name = doc['metadata']['title']
                    reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
                else: # Query
                    confidence = doc['confidence']
                    if ("type" in doc['metadata']) and (doc['metadata']['type'] == "QUESTION_ANSWER"):
                        reference = reference + f"{i+1}. <a href=\"#\" onClick=\"alert(`{excerpt}`)\">FAQ ({confidence})</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
                    else:
                        uri = ""
                        if "title" in doc['metadata']:
                            #print('metadata: ', json.dumps(doc['metadata']))
                            name = doc['metadata']['title']
                            if name: 
                                uri = path+doc_prefix+parse.quote(name)

                        page = ""
                        if "document_attributes" in doc['metadata']:
                            if "_excerpt_page_number" in doc['metadata']['document_attributes']:
                                page = doc['metadata']['document_attributes']['_excerpt_page_number']
                                                
                        if page: 
                            reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name} ({confidence})</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
                        elif uri:
                            reference = reference + f"{i+1}. <a href={uri} target=_blank>{name} ({confidence})</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
            elif doc['rag_type'] == 'opensearch':
                print(f'## Document(get_reference) {i+1}: {doc}')
                
                page = ""
                if "document_attributes" in doc['metadata']:
                    if "_excerpt_page_number" in doc['metadata']['document_attributes']:
                        page = doc['metadata']['document_attributes']['_excerpt_page_number']
                uri = doc['metadata']['source']
                name = doc['metadata']['title']

                #print('opensearch page: ', page)

                if page:                
                    reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
                else:
                    reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
        
            elif doc['rag_type'] == 'faiss':
                print(f'## Document(get_reference) {i+1}: {doc}')
                
                page = ""
                if "document_attributes" in doc['metadata']:
                    if "_excerpt_page_number" in doc['metadata']['document_attributes']:
                        page = doc['metadata']['document_attributes']['_excerpt_page_number']
                uri = doc['metadata']['source']
                name = doc['metadata']['title']

                if page: 
                    reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
                elif uri:
                    reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
            
            elif doc['rag_type'] == 'search': # google search
                print(f'## Document(get_reference) {i+1}: {doc}')
                
                uri = doc['metadata']['source']
                name = doc['metadata']['title']
                reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
        
    return reference

def retrieve_from_vectorstore(query, top_k, rag_type):
    print(f"query: {query} ({rag_type})")

    relevant_docs = []
    if rag_type == 'faiss' and isReady:
        #relevant_documents = vectorstore_faiss.similarity_search(query)
        #query_embedding = bedrock_embeddings.embed_query(query)
        #print('query_embedding: ', query_embedding)
        #relevant_documents = vectorstore_faiss.similarity_search_by_vector(query_embedding)

        relevant_documents = vectorstore_faiss.similarity_search_with_score(
            query = query,
            k = top_k,
        )
        
        #relevant_documents1 = vectorstore_faiss.similarity_search_with_relevance_scores(query)
        #print('relevant_documents1: ',relevant_documents1)
        
        for i, document in enumerate(relevant_documents):
            #print('document.page_content:', document.page_content)
            #print('document.metadata:', document.metadata)
            print(f'## Document(retrieve_from_vectorstore) {i+1}: {document}')

            name = document[0].metadata['name']
            page = ""
            if "page" in document[0].metadata:
                page = document[0].metadata['page']
            uri = ""
            if "uri" in document[0].metadata:
                uri = document[0].metadata['uri']

            confidence = int(document[1])
            assessed_score = int(document[1])
            
            if page:
                doc_info = {
                    "rag_type": rag_type,
                    #"api_type": api_type,
                    "confidence": confidence,
                    "metadata": {
                        #"type": query_result_type,
                        #"document_id": document_id,
                        "source": uri,
                        "title": name,
                        "excerpt": document[0].page_content,
                        "translated_excerpt": "",
                        "document_attributes": {
                            "_excerpt_page_number": page
                        }
                    },
                    #"query_id": query_id,
                    #"feedback_token": feedback_token
                    "assessed_score": assessed_score,
                }

            else: 
                doc_info = {
                    "rag_type": rag_type,
                    #"api_type": api_type,
                    "confidence": confidence,
                    "metadata": {
                        #"type": query_result_type,
                        #"document_id": document_id,
                        "source": uri,
                        "title": name,
                        "excerpt": document[0].page_content,
                        "translated_excerpt": ""
                    },
                    #"query_id": query_id,
                    #"feedback_token": feedback_token
                    "assessed_score": assessed_score,
                }
            
            relevant_docs.append(doc_info)
            
    elif rag_type == 'opensearch':
        #relevant_documents = vectorstore_opensearch.similarity_search(
        #    query = query,
        #    k = top_k,
        #)

        relevant_documents = vectorstore_opensearch.similarity_search_with_score(
            query = query,
            k = top_k,
        )
        #print('(opensearch score) relevant_documents: ', relevant_documents)

        for i, document in enumerate(relevant_documents):
            #print('document.page_content:', document.page_content)
            #print('document.metadata:', document.metadata)
            print(f'## Document(retrieve_from_vectorstore) {i+1}: {document}')

            name = document[0].metadata['name']
            print('metadata: ', document[0].metadata)

            page = ""
            if "page" in document[0].metadata:
                page = document[0].metadata['page']
            uri = ""
            if "uri" in document[0].metadata:
                uri = document[0].metadata['uri']

            excerpt = document[0].page_content
            confidence = str(document[1])
            assessed_score = str(document[1])

            if page:
                print('page: ', page)
                doc_info = {
                    "rag_type": rag_type,
                    #"api_type": api_type,
                    "confidence": confidence,
                    "metadata": {
                        #"type": query_result_type,
                        #"document_id": document_id,
                        "source": uri,
                        "title": name,
                        "excerpt": excerpt,
                        "translated_excerpt": "",
                        "document_attributes": {
                            "_excerpt_page_number": page
                        }
                    },
                    #"query_id": query_id,
                    #"feedback_token": feedback_token
                    "assessed_score": assessed_score,
                }
            else:
                doc_info = {
                    "rag_type": rag_type,
                    #"api_type": api_type,
                    "confidence": confidence,
                    "metadata": {
                        #"type": query_result_type,
                        #"document_id": document_id,
                        "source": uri,
                        "title": name,
                        "excerpt": excerpt,
                        "translated_excerpt": ""
                    },
                    #"query_id": query_id,
                    #"feedback_token": feedback_token
                    "assessed_score": assessed_score,
                }
            relevant_docs.append(doc_info)

    return relevant_docs

from langchain.schema import BaseMessage
_ROLE_MAP = {"human": "\n\nHuman: ", "ai": "\n\nAssistant: "}
def _get_chat_history(chat_history):
    #print('_get_chat_history: ', chat_history)
    buffer = ""
    for dialogue_turn in chat_history:
        if isinstance(dialogue_turn, BaseMessage):
            role_prefix = _ROLE_MAP.get(dialogue_turn.type, f"{dialogue_turn.type}: ")
            buffer += f"\n{role_prefix}{dialogue_turn.content}"
        elif isinstance(dialogue_turn, tuple):
            human = "\n\nHuman: " + dialogue_turn[0]
            ai = "\n\nAssistant: " + dialogue_turn[1]
            buffer += "\n" + "\n".join([human, ai])
        else:
            raise ValueError(
                f"Unsupported chat history format: {type(dialogue_turn)}."
                f" Full chat history: {chat_history} "
            )
    #print('buffer: ', buffer)
    return buffer

def create_ConversationalRetrievalChain(llm, PROMPT, retriever):  
    condense_template = """\n\nHuman: Given the following <history> and a follow up question, rephrase the follow up question to be a standalone question, in its original language. Answer only with the new question.

    <history>
    {chat_history}
    </history>
    
    Follow Up Input: {question}
    
    Assistant: Standalone question:"""
    CONDENSE_QUESTION_PROMPT = PromptTemplate.from_template(condense_template)
        
    qa = ConversationalRetrievalChain.from_llm(
        llm=llm, 
        retriever=retriever,
        condense_question_prompt=CONDENSE_QUESTION_PROMPT, # chat history and new question
        combine_docs_chain_kwargs={'prompt': PROMPT},

        memory=memory_chain,
        get_chat_history=_get_chat_history,
        verbose=False, # for logging to stdout
        
        #max_tokens_limit=300,
        chain_type='stuff', # 'refine'
        rephrase_question=True,  # to pass the new generated question to the combine_docs_chain       
        return_source_documents=True, # retrieved source
        return_generated_question=False, # generated question
    )

    return qa

def retrieve_process_from_RAG(conn, query, top_k, rag_type):
    relevant_docs = []
    if rag_type == 'kendra':
        rel_docs = retrieve_from_kendra(query=query, top_k=top_k)      
        print('rel_docs (kendra): '+json.dumps(rel_docs))
    else:
        rel_docs = retrieve_from_vectorstore(query=query, top_k=top_k, rag_type=rag_type)
        print(f'rel_docs ({rag_type}): '+json.dumps(rel_docs))

    if(len(rel_docs)>=1):
        for doc in rel_docs:
            relevant_docs.append(doc)  
    
    conn.send(relevant_docs)
    conn.close()

def debug_msg_for_revised_question(llm, revised_question, chat_history, connectionId, requestId):
    global history_length, token_counter_history # debugMessageMode 

    history_context = ""
    token_counter_history = 0
    for history in chat_history:
        history_context = history_context + history

    if history_context:
        token_counter_history = llm.get_num_tokens(history_context)
        print('token_size of history: ', token_counter_history)

    history_length = len(history_context)

    sendDebugMessage(connectionId, requestId, f"새로운 질문: {revised_question}\n * 대화이력({str(history_length)}자, {token_counter_history} Tokens)을 활용하였습니다.")

def get_relevant_documents_using_parallel_processing(question, top_k):
    relevant_docs = []    

    processes = []
    parent_connections = []
    for rag in capabilities:
        parent_conn, child_conn = Pipe()
        parent_connections.append(parent_conn)
            
        process = Process(target=retrieve_process_from_RAG, args=(child_conn, question, top_k, rag))
        processes.append(process)

    for process in processes:
        process.start()
            
    for parent_conn in parent_connections:
        rel_docs = parent_conn.recv()

        if(len(rel_docs)>=1):
            for doc in rel_docs:
                relevant_docs.append(doc)    

    for process in processes:
        process.join()
    
    #print('relevant_docs: ', relevant_docs)
    return relevant_docs

def translate_process_from_relevent_doc(conn, llm, doc, bedrock_region):
    try: 
        translated_excerpt = traslation_to_korean(llm=llm, msg=doc['metadata']['excerpt'])
        print(f"translated_excerpt ({bedrock_region}): {translated_excerpt}")

        doc['metadata']['translated_excerpt'] = translated_excerpt
    
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)       
        raise Exception (f"Not able to translate: {doc}")   
    
    conn.send(doc)
    conn.close()

def translate_relevant_documents_using_parallel_processing(docs):
    selected_LLM = 0
    relevant_docs = []    
    processes = []
    parent_connections = []
    for doc in docs:
        parent_conn, child_conn = Pipe()
        parent_connections.append(parent_conn)
            
        llm = get_llm(profile_of_LLMs, selected_LLM)
        bedrock_region = profile_of_LLMs[selected_LLM]['bedrock_region']

        process = Process(target=translate_process_from_relevent_doc, args=(child_conn, llm, doc, bedrock_region))
        processes.append(process)

        selected_LLM = selected_LLM + 1
        if selected_LLM == len(profile_of_LLMs):
            selected_LLM = 0

    for process in processes:
        process.start()
            
    for parent_conn in parent_connections:
        doc = parent_conn.recv()
        relevant_docs.append(doc)    

    for process in processes:
        process.join()
    
    #print('relevant_docs: ', relevant_docs)
    return relevant_docs

def get_answer_using_RAG(llm, text, conv_type, connectionId, requestId, bedrock_embeddings, rag_type):
    global time_for_revise, time_for_rag, time_for_inference, time_for_priority_search, number_of_relevant_docs  # for debug
    time_for_revise = time_for_rag = time_for_inference = time_for_priority_search = number_of_relevant_docs = 0

    global time_for_rag_inference, time_for_rag_question_translation, time_for_rag_2nd_inference, time_for_rag_translation
    time_for_rag_inference = time_for_rag_question_translation = time_for_rag_2nd_inference = time_for_rag_translation = 0
    
    reference = ""
    if rag_type == 'all': # kendra, opensearch, faiss
        start_time_for_revise = time.time()

        revised_question = get_revised_question(llm, connectionId, requestId, text) # generate new prompt using chat history                    
        PROMPT = get_prompt_template(revised_question, conv_type, rag_type)
        # print('PROMPT: ', PROMPT)        

        end_time_for_revise = time.time()
        time_for_revise = end_time_for_revise - start_time_for_revise
        print('processing time for revised question: ', time_for_revise)

        relevant_docs = [] 
        if useParallelRAG == 'true':  # parallel processing
            print('start RAG for revised question')
            relevant_docs = get_relevant_documents_using_parallel_processing(question=revised_question, top_k=top_k)

            end_time_for_rag_inference = time.time()
            time_for_rag_inference = end_time_for_rag_inference - end_time_for_revise
            print('processing time for RAG (Inference): ', time_for_rag_inference)

            if allowDualSearch=='true' and isKorean(text)==True:
                print('start RAG for translated revised question')
                translated_revised_question = traslation_to_english(llm=llm, msg=revised_question)
                print('translated_revised_question: ', translated_revised_question)

                if debugMessageMode=='true':
                    sendDebugMessage(connectionId, requestId, f"새로운 질문: {revised_question}\n번역된 새로운 질문: {translated_revised_question}")

                end_time_for_rag_question_translation = time.time()
                time_for_rag_question_translation = end_time_for_rag_question_translation - end_time_for_rag_inference
                print('processing time for RAG (Question Translation): ', time_for_rag_question_translation)

                if allowDualSearchWithMulipleProcessing == True:
                    relevant_docs_using_translated_question = get_relevant_documents_using_parallel_processing(question=translated_revised_question, top_k=4)

                    end_time_for_rag_2nd_inference = time.time()
                    time_for_rag_2nd_inference = end_time_for_rag_2nd_inference - end_time_for_rag_question_translation
                    print('processing time for RAG (2nd Inference): ', time_for_rag_2nd_inference)
                    
                    docs_translation_required = []
                    if len(relevant_docs_using_translated_question)>=1:
                        for i, doc in enumerate(relevant_docs_using_translated_question):
                            if isKorean(doc)==False:
                                docs_translation_required.append(doc)
                            else:
                                print(f"original {i}: {doc}")
                                relevant_docs.append(doc)
                                           
                        translated_docs = translate_relevant_documents_using_parallel_processing(docs_translation_required)
                        for i, doc in enumerate(translated_docs):
                            print(f"#### {i} (ENG): {doc['metadata']['excerpt']}")
                            print(f"#### {i} (KOR): {doc['metadata']['translated_excerpt']}")
                            relevant_docs.append(doc)
                        
                        end_time_for_rag_translation = time.time()
                        time_for_rag_translation = end_time_for_rag_translation - end_time_for_rag_2nd_inference
                        print('processing time for RAG (translation): ', time_for_rag_translation)

                else:
                    relevant_docs_using_translated_question = []
                    for reg in capabilities:            
                        if reg == 'kendra':
                            rel_docs = retrieve_from_kendra(query=translated_revised_question, top_k=top_k)      
                            print('rel_docs (kendra): '+json.dumps(rel_docs))
                        else:
                            rel_docs = retrieve_from_vectorstore(query=translated_revised_question, top_k=top_k, rag_type=reg)
                            print(f'rel_docs ({reg}): '+json.dumps(rel_docs))
                    
                        if(len(rel_docs)>=1):
                            for doc in rel_docs:
                                relevant_docs_using_translated_question.append(doc)    

                    if len(relevant_docs_using_translated_question)>=1:
                        for i, doc in enumerate(relevant_docs_using_translated_question):
                            if isKorean(doc)==False:
                                translated_excerpt = traslation_to_korean(llm=llm, msg=doc['metadata']['excerpt'])
                                print(f"#### {i} (ENG): {doc['metadata']['excerpt']}")
                                print(f"#### {i} (KOR): {translated_excerpt}")

                                #doc['metadata']['excerpt'] = translated_excerpt
                                doc['metadata']['translated_excerpt'] = translated_excerpt
                                relevant_docs.append(doc)
                            else:
                                print(f"original {i}: {doc}")
                                relevant_docs.append(doc)
        else: # sequencial processing
            print('start the sequencial processing for multiple RAG')
            for reg in capabilities:            
                if reg == 'kendra':
                    rel_docs = retrieve_from_kendra(query=revised_question, top_k=top_k)      
                    print('rel_docs (kendra): '+json.dumps(rel_docs))
                else:
                    rel_docs = retrieve_from_vectorstore(query=revised_question, top_k=top_k, rag_type=reg)
                    print(f'rel_docs ({reg}): '+json.dumps(rel_docs))
                
                if(len(rel_docs)>=1):
                    for doc in rel_docs:
                        relevant_docs.append(doc)

        if debugMessageMode=='true':
            for i, doc in enumerate(relevant_docs):
                print(f"#### relevant_docs ({i}): {json.dumps(doc)}")

        end_time_for_rag = time.time()
        time_for_rag = end_time_for_rag - end_time_for_revise
        print('processing time for RAG: ', time_for_rag)

        selected_relevant_docs = []
        if len(relevant_docs)>=1:
            selected_relevant_docs = priority_search(revised_question, relevant_docs, bedrock_embeddings)
            print('selected_relevant_docs: ', json.dumps(selected_relevant_docs))

        if len(selected_relevant_docs)==0:
            print('No relevant document! So use google api')            
            api_key = google_api_key
            cse_id = google_cse_id 
            
            relevant_docs = []
            try: 
                service = build("customsearch", "v1", developerKey=api_key)
                result = service.cse().list(q=revised_question, cx=cse_id).execute()
                # print('google search result: ', result)

                if "items" in result:
                    for item in result['items']:
                        api_type = "google api"
                        excerpt = item['snippet']
                        uri = item['link']
                        title = item['title']
                        confidence = ""
                        assessed_score = ""
                        
                        doc_info = {
                            "rag_type": 'search',
                            "api_type": api_type,
                            "confidence": confidence,
                            "metadata": {
                                #"type": query_result_type,
                                # "document_id": document_id,
                                "source": uri,
                                "title": title,
                                "excerpt": excerpt,
                                "translated_excerpt": "",
                                #"document_attributes": {
                                #    "_excerpt_page_number": page
                                #}
                            },
                            #"query_id": query_id,
                            #"feedback_token": feedback_token
                            "assessed_score": assessed_score,
                            #"result_id": result_id
                        }
                        relevant_docs.append(doc_info)                
            except Exception:
                err_msg = traceback.format_exc()
                print('error message: ', err_msg)       

                sendErrorMessage(connectionId, requestId, err_msg)    
                raise Exception ("Not able to search using google api")   
            
            if len(relevant_docs)>=1:
                selected_relevant_docs = priority_search(revised_question, relevant_docs, bedrock_embeddings)
                print('selected_relevant_docs: ', json.dumps(selected_relevant_docs))
            # print('selected_relevant_docs (google): ', selected_relevant_docs)

        end_time_for_priority_search = time.time() 
        time_for_priority_search = end_time_for_priority_search - end_time_for_rag
        print('processing time for priority search: ', time_for_priority_search)
        number_of_relevant_docs = len(selected_relevant_docs)

        relevant_context = ""
        for document in selected_relevant_docs:
            if document['metadata']['translated_excerpt']:
                content = document['metadata']['translated_excerpt']
            else:
                content = document['metadata']['excerpt']
        
            relevant_context = relevant_context + content + "\n\n"
        print('relevant_context: ', relevant_context)

        try: 
            isTyping(connectionId, requestId)
            stream = llm(PROMPT.format(context=relevant_context, question=revised_question))
            msg = readStreamMsg(connectionId, requestId, stream)            
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)       
            sendErrorMessage(connectionId, requestId, err_msg)    
            raise Exception ("Not able to request to LLM")    

        if len(selected_relevant_docs)>=1 and enableReference=='true':
            reference = get_reference(selected_relevant_docs, rag_method, rag_type, path, doc_prefix)  

        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - end_time_for_priority_search
        print('processing time for inference: ', time_for_inference)
        
    else:        
        if rag_method == 'RetrievalQA': # RetrievalQA
            start_time_for_revise = time.time()

            revised_question = get_revised_question(llm, connectionId, requestId, text) # generate new prompt using chat history

            end_time_for_revise = time.time()
            time_for_revise = end_time_for_revise - start_time_for_revise
            print('processing time for revised question: ', time_for_revise)

            start_time_for_rag = time.time()
            PROMPT = get_prompt_template(revised_question, conv_type, rag_type)
            #print('PROMPT: ', PROMPT)

            if rag_type=='kendra':
                retriever = kendraRetriever
            elif rag_type=='opensearch':
                retriever = vectorstore_opensearch.as_retriever(
                    search_type="similarity", 
                    search_kwargs={
                        #"k": 3, 'score_threshold': 0.8
                        "k": top_k
                    }
                )
            elif rag_type=='faiss' and isReady:
                retriever = vectorstore_faiss.as_retriever(
                    search_type="similarity", 
                    search_kwargs={
                        #"k": 3, 'score_threshold': 0.8
                        "k": top_k
                    }
                )

            qa = RetrievalQA.from_chain_type(
                llm=llm,
                chain_type="stuff",
                retriever=retriever,
                return_source_documents=True,
                chain_type_kwargs={"prompt": PROMPT}
            )

            isTyping(connectionId, requestId)        
            result = qa({"query": revised_question})    
            print('result: ', result)

            msg = readStreamMsg(connectionId, requestId, result['result'])

            source_documents = result['source_documents']
            print('source_documents: ', source_documents)

            if len(source_documents)>=1 and enableReference=='true':
                reference = get_reference(source_documents, rag_method, rag_type, path, doc_prefix)    
            
            end_time_for_rag = time.time()
            time_for_rag = end_time_for_rag - end_time_for_revise
            print('processing time for RAG: ', time_for_rag)
            number_of_relevant_docs = len(source_documents)

        elif rag_method == 'ConversationalRetrievalChain': # ConversationalRetrievalChain
            start_time_for_rag = time.time()

            PROMPT = get_prompt_template(text, conv_type, rag_type)
            if rag_type == 'kendra':
                qa = create_ConversationalRetrievalChain(llm, PROMPT, retriever=kendraRetriever)            
            elif rag_type == 'opensearch': # opensearch
                vectorstoreRetriever = vectorstore_opensearch.as_retriever(
                    search_type="similarity", 
                    search_kwargs={
                        "k": 5
                    }
                )
                qa = create_ConversationalRetrievalChain(llm, PROMPT, retriever=vectorstoreRetriever)
            elif rag_type == 'faiss' and isReady: # faiss
                vectorstoreRetriever = vectorstore_faiss.as_retriever(
                    search_type="similarity", 
                    search_kwargs={
                        "k": 5
                    }
                )
                qa = create_ConversationalRetrievalChain(llm, PROMPT, retriever=vectorstoreRetriever)
            
            result = qa({"question": text})
            
            msg = result['answer']
            print('\nquestion: ', result['question'])    
            print('answer: ', result['answer'])    
            # print('chat_history: ', result['chat_history'])    
            print('source_documents: ', result['source_documents']) 

            if len(result['source_documents'])>=1 and enableReference=='true':
                reference = get_reference(result['source_documents'], rag_method, rag_type, path, doc_prefix)
            
            end_time_for_rag = time.time()
            time_for_rag = end_time_for_rag - start_time_for_rag
            print('processing time for RAG: ', time_for_rag)
            number_of_relevant_docs = len(result['source_documents'])
        
        elif rag_method == 'RetrievalPrompt': # RetrievalPrompt
            start_time_for_revise = time.time()

            revised_question = get_revised_question(llm, connectionId, requestId, text) # generate new prompt using chat history
            
            end_time_for_revise = time.time()
            time_for_revise = end_time_for_revise - start_time_for_revise
            print('processing time for revised question: ', time_for_revise)

            PROMPT = get_prompt_template(revised_question, conv_type, rag_type)
            #print('PROMPT: ', PROMPT)

            if rag_type == 'kendra':
                relevant_docs = retrieve_from_kendra(query=revised_question, top_k=top_k)
                if len(relevant_docs) >= 1:
                    relevant_docs = priority_search(revised_question, relevant_docs, bedrock_embeddings)
            else:
                relevant_docs = retrieve_from_vectorstore(query=revised_question, top_k=top_k, rag_type=rag_type)
            print('relevant_docs: ', json.dumps(relevant_docs))

            end_time_for_rag = time.time()
            time_for_rag = end_time_for_rag - end_time_for_revise
            print('processing time for RAG: ', time_for_rag)
            number_of_relevant_docs = len(relevant_docs)

            relevant_context = ""
            for document in relevant_docs:
                if document['metadata']['translated_excerpt']:
                    content = document['metadata']['translated_excerpt']
                else:
                    content = document['metadata']['excerpt']

                relevant_context = relevant_context + content + "\n\n"
            print('relevant_context: ', relevant_context)

            try: 
                isTyping(connectionId, requestId)
                stream = llm(PROMPT.format(context=relevant_context, question=revised_question))
                msg = readStreamMsg(connectionId, requestId, stream)
            except Exception:
                err_msg = traceback.format_exc()
                print('error message: ', err_msg)       

                sendErrorMessage(connectionId, requestId, err_msg)    
                raise Exception ("Not able to request to LLM")    

            if len(relevant_docs)>=1 and enableReference=='true':
                reference = get_reference(relevant_docs, rag_method, rag_type, path, doc_prefix)
            
            end_time_for_inference = time.time()
            time_for_inference = end_time_for_inference - end_time_for_rag
            print('processing time for inference: ', time_for_inference)

    global relevant_length, token_counter_relevant_docs
    
    if debugMessageMode=='true':   # extract chat history for debug
        chat_history_all = extract_chat_history_from_memory()
        # print('chat_history_all: ', chat_history_all)

        history_length = []
        for history in chat_history_all:
            history_length.append(len(history))
        print('chat_history length: ', history_length)

        relevant_length = len(relevant_context)
        token_counter_relevant_docs = llm.get_num_tokens(relevant_context)

    memory_chain.chat_memory.add_user_message(text)  # append new diaglog
    memory_chain.chat_memory.add_ai_message(msg)

    return msg, reference

def get_answer_using_ConversationChain(text, conversation, conv_type, connectionId, requestId, rag_type):
    conversation.prompt = get_prompt_template(text, conv_type, rag_type)
    try: 
        isTyping(connectionId, requestId)    
        stream = conversation.predict(input=text)
        #print('stream: ', stream)                    
        msg = readStreamMsg(connectionId, requestId, stream)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':   # extract chat history for debug
        chats = memory_chat.load_memory_variables({})
        chat_history_all = chats['history']
        # print('chat_history_all: ', chat_history_all)
        
        print('chat_history length: ', len(chat_history_all))

    return msg

def get_answer_from_PROMPT(llm, text, conv_type, connectionId, requestId):
    PROMPT = get_prompt_template(text, conv_type, "")
    #print('PROMPT: ', PROMPT)

    try: 
        isTyping(connectionId, requestId)
        stream = llm(PROMPT.format(input=text))
        msg = readStreamMsg(connectionId, requestId, stream)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")
    
    return msg

def create_metadata(bucket, key, meta_prefix, s3_prefix, uri, category, documentId):    
    title = key
    timestamp = int(time.time())

    metadata = {
        "Attributes": {
            "_category": category,
            "_source_uri": uri,
            "_version": str(timestamp),
            "_language_code": "ko"
        },
        "Title": title,
        "DocumentId": documentId,        
    }
    print('metadata: ', metadata)

    client = boto3.client('s3')
    try: 
        client.put_object(
            Body=json.dumps(metadata), 
            Bucket=bucket, 
            Key=meta_prefix+s3_prefix+'/'+key+'.metadata.json' 
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create meta file")

def get_text_speech(path, speech_prefix, bucket, msg):
    ext = "mp3"
    polly = boto3.client('polly')
    try:
        response = polly.start_speech_synthesis_task(
            Engine='neural',
            LanguageCode='ko-KR',
            OutputFormat=ext,
            OutputS3BucketName=bucket,
            OutputS3KeyPrefix=speech_prefix,
            Text=msg,
            TextType='text',
            VoiceId='Seoyeon'        
        )
        print('response: ', response)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create voice")
    
    object = '.'+response['SynthesisTask']['TaskId']+'.'+ext
    print('object: ', object)

    return path+speech_prefix+parse.quote(object)

def traslation_to_korean(llm, msg):
    PROMPT = """\n\nHuman: Here is an article, contained in <article> tags. Translate the article to Korean. Put it in <result> tags.
            
    <article>
    {input}
    </article>
                        
    Assistant:"""

    try: 
        translated_msg = llm(PROMPT.format(input=msg))
        #print('translated_msg: ', translated_msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to translate the message")
    
    msg = translated_msg[translated_msg.find('<result>')+9:len(translated_msg)-10]
    
    return msg.replace("\n"," ")

def traslation_to_english(llm, msg):
    PROMPT = """\n\nHuman: 다음의 <article>를 English로 번역하세요. 머리말은 건너뛰고 본론으로 바로 들어가주세요. 또한 결과는 <result> tag를 붙여주세요.

    <article>
    {input}
    </article>
                        
    Assistant:"""

    try: 
        translated_msg = llm(PROMPT.format(input=msg))
        #print('translated_msg: ', translated_msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to translate the message")
    
    return translated_msg[translated_msg.find('<result>')+9:len(translated_msg)-10]

def getResponse(connectionId, jsonBody):
    userId  = jsonBody['user_id']
    # print('userId: ', userId)
    requestId  = jsonBody['request_id']
    # print('requestId: ', requestId)
    requestTime  = jsonBody['request_time']
    # print('requestTime: ', requestTime)
    type  = jsonBody['type']
    # print('type: ', type)
    body = jsonBody['body']
    # print('body: ', body)
    conv_type = jsonBody['conv_type']  # conversation type
    print('Conversation Type: ', conv_type)

    rag_type = ""
    if 'rag_type' in jsonBody:
        if jsonBody['rag_type']:
            rag_type = jsonBody['rag_type']  # RAG type
            print('rag_type: ', rag_type)

    global vectorstore_opensearch, vectorstore_faiss, enableReference
    global map_chain, map_chat, memory_chat, memory_chain, isReady, debugMessageMode, selected_LLM, allowDualSearch

    # Multi-LLM
    profile = profile_of_LLMs[selected_LLM]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'selected_LLM: {selected_LLM}, bedrock_region: {bedrock_region}, modelId: {modelId}')
    # print('profile: ', profile)
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }            
        )
    )
    parameters = get_parameter(profile['model_type'], int(profile['maxOutputTokens']))
    print('parameters: ', parameters)

    # langchain for bedrock
    llm = Bedrock(
        model_id=modelId, 
        client=boto3_bedrock, 
        streaming=True,
        callbacks=[StreamingStdOutCallbackHandler()],
        model_kwargs=parameters)

    # embedding for RAG
    bedrock_embeddings = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = 'amazon.titan-embed-text-v1' 
    )    

    # create memory
    if userId in map_chat or userId in map_chain:  
        print('memory exist. reuse it!')        
        memory_chain = map_chain[userId]
        memory_chat = map_chat[userId]
        conversation = ConversationChain(llm=llm, verbose=False, memory=memory_chat)
        
    else: 
        print('memory does not exist. create new one!')
        memory_chain = ConversationBufferWindowMemory(memory_key="chat_history", output_key='answer', return_messages=True, k=10)
        map_chain[userId] = memory_chain
        
        memory_chat = ConversationBufferWindowMemory(human_prefix='Human', ai_prefix='Assistant', k=MSG_HISTORY_LENGTH)
        map_chat[userId] = memory_chat

        # memory_chain = ConversationBufferMemory(memory_key="chat_history", return_messages=True)
        #from langchain.memory import ConversationSummaryBufferMemory
        #memory_chat = ConversationSummaryBufferMemory(llm=llm, max_token_limit=1024,
        #    human_prefix='Human', ai_prefix='Assistant') #Maintains a summary of previous messages

        allowTime = getAllowTime()
        load_chat_history(userId, allowTime, conv_type, rag_type)
        conversation = ConversationChain(llm=llm, verbose=False, memory=memory_chat)

    # rag sources
    if conv_type == 'qa' and (rag_type == 'opensearch' or rag_type == 'all'):
        vectorstore_opensearch = OpenSearchVectorSearch(
            index_name = "rag-index-*", # all
            #index_name=f"rag-index-{userId}',
            is_aoss = False,
            ef_search = 1024, # 512(default)
            m=48,
            #engine="faiss",  # default: nmslib
            embedding_function = bedrock_embeddings,
            opensearch_url=opensearch_url,
            http_auth=(opensearch_account, opensearch_passwd), # http_auth=awsauth,
        )
    elif conv_type == 'qa' and rag_type == 'faiss':
        print('isReady = ', isReady)

    start = int(time.time())    

    msg = ""
    reference = ""
    speech_uri = ""
    if type == 'text' and body[:11] == 'list models':
        bedrock_client = boto3.client(
            service_name='bedrock',
            region_name=bedrock_region,
        )
        modelInfo = bedrock_client.list_foundation_models()    
        print('models: ', modelInfo)

        msg = f"The list of models: \n"
        lists = modelInfo['modelSummaries']
        
        for model in lists:
            msg += f"{model['modelId']}\n"
        
        msg += f"current model: {modelId}"
        print('model lists: ', msg)          

        sendResultMessage(connectionId, requestId, msg)  
    else:             
        if type == 'text':
            text = body
            print('query: ', text)

            querySize = len(text)
            textCount = len(text.split())
            print(f"query size: {querySize}, words: {textCount}")

            if text == 'enableReference':
                enableReference = 'true'
                msg  = "Referece is enabled"
            elif text == 'disableReference':
                enableReference = 'false'
                msg  = "Reference is disabled"
            elif text == 'useOpenSearch':
                rag_type = 'opensearch'
                msg  = "OpenSearch is selected for Knowledge Database"
            elif text == 'useFaiss':
                rag_type = 'faiss'
                msg  = "Faiss is selected for Knowledge Database"
            elif text == 'useKendra':
                rag_type = 'kendra'
                msg  = "Kendra is selected for Knowledge Database"
            elif text == 'enableDebug':
                debugMessageMode = 'true'
                msg  = "Debug messages will be delivered to the client."
            elif text == 'disableDebug':
                debugMessageMode = 'false'
                msg  = "Debug messages will not be delivered to the client."
            elif text == 'enableDualSearch':
                allowDualSearch = 'true'
                msg  = "Translated question is enabled"
            elif text == 'disableDualSearch':
                allowDualSearch = 'false'
                msg  = "Translated question is disabled"

            elif text == 'clearMemory':
                memory_chain.clear()
                map_chain[userId] = memory_chain
                memory_chat.clear()                
                map_chat[userId] = memory_chat
                conversation = ConversationChain(llm=llm, verbose=False, memory=memory_chat)
                    
                print('initiate the chat memory!')
                msg  = "The chat memory was intialized in this session."
            else:          
                if conv_type == 'qa':   # RAG
                    print(f'rag_type: {rag_type}, rag_method: {rag_method}')
                          
                    if rag_type == 'faiss' and isReady==False:                               
                        msg = get_answer_using_ConversationChain(text, conversation, conv_type, connectionId, requestId, rag_type)      

                        memory_chain.chat_memory.add_user_message(text)  # append new diaglog
                        memory_chain.chat_memory.add_ai_message(msg)
                    else: 
                        msg, reference = get_answer_using_RAG(llm, text, conv_type, connectionId, requestId, bedrock_embeddings, rag_type)     
                
                elif conv_type == 'normal' or conv_type == 'funny':      # normal
                    msg = get_answer_using_ConversationChain(text, conversation, conv_type, connectionId, requestId, rag_type)
                
                elif conv_type == 'none':   # no prompt
                    try: 
                        msg = llm(HUMAN_PROMPT+text+AI_PROMPT)
                    except Exception:
                        err_msg = traceback.format_exc()
                        print('error message: ', err_msg)

                        sendErrorMessage(connectionId, requestId, err_msg)    
                        raise Exception ("Not able to request to LLM")
                else: 
                    msg = get_answer_from_PROMPT(llm, text, conv_type, connectionId, requestId)

                # token counter
                if debugMessageMode=='true':
                    token_counter_input = llm.get_num_tokens(text)
                    token_counter_output = llm.get_num_tokens(msg)
                    print(f"token_counter: question: {token_counter_input}, answer: {token_counter_output}")
                
        elif type == 'document':
            isTyping(connectionId, requestId)
            
            object = body
            file_type = object[object.rfind('.')+1:len(object)]            
            print('file_type: ', file_type)

            if file_type == 'csv':
                docs = load_csv_document(path, doc_prefix, object)
                contexts = []
                for doc in docs:
                    contexts.append(doc.page_content)
                print('contexts: ', contexts)

                msg = get_summary(llm, contexts)

            elif file_type == 'pdf' or file_type == 'txt' or file_type == 'pptx' or file_type == 'docx':
                texts = load_document(file_type, object)

                docs = []
                for i in range(len(texts)):
                    docs.append(
                        Document(
                            page_content=texts[i],
                            metadata={
                                'name': object,
                                # 'page':i+1,
                                'uri': path+doc_prefix+parse.quote(object)
                            }
                        )
                    )
                print('docs[0]: ', docs[0])    
                print('docs size: ', len(docs))

                contexts = []
                for doc in docs:
                    contexts.append(doc.page_content)
                print('contexts: ', contexts)

                msg = get_summary(llm, contexts)
            else:
                msg = "uploaded file: "+object
                                
            if conv_type == 'qa':
                start_time = time.time()
                
                category = "upload"
                key = object
                documentId = category + "-" + key
                documentId = documentId.replace(' ', '_') # remove spaces
                documentId = documentId.replace(',', '_') # remove commas
                documentId = documentId.lower() # change to lowercase
                print('documentId: ', documentId)

                if useParallelUpload == 'false':
                    if rag_type == 'all':      
                        for type in capabilities:
                            print('rag_type: ', type)
                            if type=='kendra':      
                                print('upload to kendra: ', object)           
                                store_document_for_kendra(path, doc_prefix, object, documentId)  # store the object into kendra
                                                
                            else:
                                if file_type == 'pdf' or file_type == 'txt' or file_type == 'csv' or file_type == 'pptx' or file_type == 'docx':
                                    if type == 'faiss':
                                        if isReady == False:   
                                            embeddings = bedrock_embeddings
                                            vectorstore_faiss = FAISS.from_documents( # create vectorstore from a document
                                                    docs,  # documents
                                                    embeddings  # embeddings
                                            )
                                            isReady = True
                                        else:
                                            store_document_for_faiss(docs, vectorstore_faiss)

                                    elif type == 'opensearch':    
                                        store_document_for_opensearch(bedrock_embeddings, docs, userId, documentId)
                    else:
                        print('rag_type: ', rag_type)
                        if rag_type=='kendra':
                            print('upload to kendra: ', object)
                            store_document_for_kendra(path, doc_prefix, object, documentId)  # store the object into kendra
                                                
                        else:
                            if file_type == 'pdf' or file_type == 'txt' or file_type == 'csv' or file_type == 'pptx' or file_type == 'docx':
                                if rag_type == 'faiss':
                                    if isReady == False:
                                        embeddings = bedrock_embeddings
                                        vectorstore_faiss = FAISS.from_documents( # create vectorstore from a document
                                                docs,  # documents
                                                embeddings  # embeddings
                                        )
                                        isReady = True
                                    else:
                                        store_document_for_faiss(docs, vectorstore_faiss)
                                                                        
                                elif rag_type == 'opensearch':    
                                    store_document_for_opensearch(bedrock_embeddings, docs, userId, documentId)
                            
                else:                    
                    p1 = Process(target=store_document_for_kendra, args=(path, doc_prefix, object, documentId))
                    p1.start(); p1.join()
                    
                    if file_type == 'pdf' or file_type == 'txt' or file_type == 'csv' or file_type == 'pptx' or file_type == 'docx':
                        # opensearch
                        p2 = Process(target=store_document_for_opensearch, args=(bedrock_embeddings, docs, userId, documentId,))
                        p2.start(); p2.join()

                        # faiss
                        if isReady == False:   
                            embeddings = bedrock_embeddings
                            vectorstore_faiss = FAISS.from_documents( # create vectorstore from a document
                                docs,  # documents
                                embeddings  # embeddings
                            )
                            isReady = True
                        else: 
                            #p3 = Process(target=store_document_for_faiss, args=(docs, vectorstore_faiss))
                            #p3.start(); p3.join()
                            vectorstore_faiss.add_documents(docs)       

                create_metadata(bucket=s3_bucket, key=key, meta_prefix=meta_prefix, s3_prefix=s3_prefix, uri=path+doc_prefix+parse.quote(object), category=category, documentId=documentId)
                print('processing time: ', str(time.time() - start_time))
                        
        sendResultMessage(connectionId, requestId, msg+reference)
        # print('msg+reference: ', msg+reference)

        elapsed_time = time.time() - start
        print("total run time(sec): ", elapsed_time)

        if isKorean(msg)==False and (conv_type=='qa' or  conv_type == "normal"):
            translated_msg = traslation_to_korean(llm, msg)
            print('translated_msg: ', translated_msg)
            
            if speech_generation: # generate mp3 file
                speech_uri = get_text_speech(path=path, speech_prefix=speech_prefix, bucket=s3_bucket, msg=translated_msg)
                print('speech_uri: ', speech_uri)  
            msg = msg+'\n[한국어]\n'+translated_msg
        else:
            if speech_generation: # generate mp3 file
                speech_uri = get_text_speech(path=path, speech_prefix=speech_prefix, bucket=s3_bucket, msg=msg)
                print('speech_uri: ', speech_uri)  

        item = {    # save dialog
            'user_id': {'S':userId},
            'request_id': {'S':requestId},
            'request_time': {'S':requestTime},
            'type': {'S':type},
            'body': {'S':body},
            'msg': {'S':msg+reference}
        }
        client = boto3.client('dynamodb')
        try:
            resp =  client.put_item(TableName=callLogTableName, Item=item)
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)
            raise Exception ("Not able to write into dynamodb")        
        #print('resp, ', resp)

    if selected_LLM >= len(profile_of_LLMs)-1:
        selected_LLM = 0
    else:
        selected_LLM = selected_LLM + 1

    if speech_uri:
        speech = '\n' + f'<a href={speech_uri} target=_blank>{"[결과 읽어주기 (mp3)]"}</a>'                
        sendResultMessage(connectionId, requestId, msg+reference+speech)

        if conv_type=='qa' and debugMessageMode=='true' and reference:
            statusMsg = f"\n[통계]\nRegion: {bedrock_region}\nQuestion: {str(len(text))}자 / {token_counter_input}토큰\nAnswer: {str(len(msg))}자 / {token_counter_output}토큰\n"
            statusMsg = statusMsg + f"History: {str(history_length)}자 / {token_counter_history}토큰\n"
            statusMsg = statusMsg + f"RAG: {str(relevant_length)}자 / {token_counter_relevant_docs}토큰 ({number_of_relevant_docs})\n"

            statusMsg = statusMsg + f"Time(초): "            
            if time_for_revise != 0:
                statusMsg = statusMsg + f"{time_for_revise:.2f}(Revise), "
            if time_for_rag != 0:
                statusMsg = statusMsg + f"{time_for_rag:.2f}(RAG), "
            if time_for_priority_search != 0:
                statusMsg = statusMsg + f"{time_for_priority_search:.2f}(Priority) "
            if time_for_inference != 0:
                statusMsg = statusMsg + f"{time_for_inference:.2f}(Inference), "
            statusMsg = statusMsg + f"{elapsed_time:.2f}(전체)"
            
            if time_for_rag_inference != 0:
                statusMsg = statusMsg + f"\nRAG-Detail: {time_for_rag_inference:.2f}(Inference(KOR)), "
            if time_for_rag_question_translation != 0:
                statusMsg = statusMsg + f"{time_for_rag_question_translation:.2f}(Question(ENG)), "
            if time_for_rag_2nd_inference != 0:
                statusMsg = statusMsg + f"{time_for_rag_2nd_inference:.2f}(Inference(ENG)), "
            if time_for_rag_translation != 0:
                statusMsg = statusMsg + f"{time_for_rag_translation:.2f}(Doc Translation), "

            sendResultMessage(connectionId, requestId, msg+reference+speech+statusMsg)

    return msg, reference

def lambda_handler(event, context):
    # print('event: ', event)
    
    msg = ""
    if event['requestContext']: 
        connectionId = event['requestContext']['connectionId']        
        routeKey = event['requestContext']['routeKey']
        
        if routeKey == '$connect':
            print('connected!')
        elif routeKey == '$disconnect':
            print('disconnected!')
        else:
            body = event.get("body", "")
            #print("data[0:8]: ", body[0:8])
            if body[0:8] == "__ping__":
                # print("keep alive!")                
                sendMessage(connectionId, "__pong__")
            else:
                print('connectionId: ', connectionId)
                print('routeKey: ', routeKey)
        
                jsonBody = json.loads(body)
                print('request body: ', json.dumps(jsonBody))

                requestId  = jsonBody['request_id']
                try:
                    msg, reference = getResponse(connectionId, jsonBody)

                    print('msg+reference: ', msg+reference)
                                        
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)

                    sendErrorMessage(connectionId, requestId, err_msg)    
                    raise Exception ("Not able to send a message")

    return {
        'statusCode': 200
    }
