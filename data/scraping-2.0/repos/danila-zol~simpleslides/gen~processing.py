import openai


# user input

def get_input():
  
  q5_names_ = []
  q10_info = []
  q4_demographics = {}
  time = ["_", "400", "800", "1200"]
  
  q1 = input("Как называется ваша компания? ")
  q2 = input("Какие услуги вы предоставляете? ")
  q3 = input("На какой рынок ориентируется ваша деятельность? ")
  q4_1 = input("На людей какого ВОЗРАСТА ориентирована деятельность вашего бизнес проекта? ")
  q4_2 = input("На людей какого ПОЛА ориентирована деятельность вашего бизнес проекта? ")
  q4_3 = input("На людей какого РЕГИОНА(местности) проживания ориентирована деятельность вашего бизнес проекта? ")
  q4_4 =input("На людей каких ИНТЕРЕСОВ ориентирована деятельность вашего бизнес проекта? ") 
  q4_demographics = {"age":f'{q4_1}', "gender":f'{q4_2}', "location":f'{q4_3}', "interests":f'{q4_4}'}
  def add_name():
    q5_name = input("Назовите имя и фамилию одного из ваших сотрудников: ")
    q5_xper = input("Опишите его роль в команде: ")
    q5_names_.append({"name":q5_name, "expertise":q5_xper})
  add_name()
  more_names = input("Хотите добавить данные следующего сотрудника?(y/n) ").upper()
  while more_names == "Y":
    add_name()
    more_names = input("Хотите добавить данные следующего сотрудника?(y/n) ").upper()
  q6 = input("Чего ваш проект уже достиг на данном этапе? ")
  q7 = input("На что вы направите будующие инвестиции? ")
  q8 = input("На какой размер инвестиций вы рассчитываете? ")
  q9 = input("К какому сроку вы планируете оправдать инвестиции? ")
  def add_info():
    q10_name = input("Укажите название интернет ресурса/имя и фамилию ответчика ")
    q10_link = input("Введите ссылку или номер телефона ")
    q10_info.append({"name": f'{q10_name}', "link":f'{q10_link}'})
  add_contact = input("Хотите добавить контактную информацию?(y/n) ").upper()
  while add_contact == "Y":
    add_info()
    add_contact = input("Хотите добавить контактную информацию?(y/n) ").upper()
  q11_num = input("Сколько времени должен занять ваш питчдек? Введите циыру от одного до трех (1 - 2 минуты; 2 - 5 минут; 3 - 12 минут) ")
  try:
    q11 = time[q11_num]
  except:
    q11 = time[2]

  answers = (q1, q2, q3, q4_demographics, q5_names_, q6, q7, q8, q9, q10_info, q11)
  return answers


sample_input = {'company_name' : 'SimpleCoffe',
                'business_type' : 'A coffeeshop', 
                'target market': 'Saint-Petersburg, Russia',
                'target audience': {"age":"all ages", "gender":"all genders", "location":"Saint-Petersburg", "interests":"drinking tasty coffe"},
                'project team': [{"name":"Daria Ivanova",  "expertise":"Barista"}, {"name":"Dmitry Petrov", "expertise": "accountant"}],
                'The project have already achieved': "monthly proffit of 25000 rubles, 30000 visitors per month",
                'investment target': 'aquiring new equipment, making better product, doubbling our revenue',
                'investment amount': '10 million rubles', 
                'deadline on investment promises': '1 year',
                'contact information': "connect to us directly using the number 123456789",
                'length(symbols)': '200'} 

sample_input_rus = {'company_name' : "ПростоКофе",
                'business_type' : "кофейня", 
                'target market': "Санкт-Петербург, россия",
                'target audience': {"age":"всех возростов", "gender":"любого пола", "location":"Санкт-Петербург", "interests":"Разнообразных интересов"},
                'project team': [{"name":"Дарья Иванова",  "expertise":"Бариста"}, {"name":"Дмитрий Петров", "expertise": "Бухгалтер"}],
                'The project have already achieved': "20000 посетителей в месяц, прибыль более полумиллиона рублей",
                'investment target': 'aquiring new equipment, making better product, doubbling our revenue',
                'investment amount': '10 миллионов рублей', 
                'deadline on investment promises': '1 год',
                'contact information': "", 
                'length(symbols)': '200'}   

def formated_input():
  answers = get_input()
  user_input = {'company name' : f'{answers[0]}',
                  'business type' : f'{answers[1]}', 
                  'target market': f'{answers[2]}',
                  'target audience': answers[3],
                  'project team': answers[4],
                  'The project have already achieved': f"{answers[5]}",
                  'investment target': f'{answers[6]}',
                  'investment amount': f'{answers[7]}', 
                  'deadline on investment promises': f'{answers[8]}',
                  'information': f'{answers[9]}',
                  'length(symbols)': f'{answers[10]}'} 
  return user_input


def compress_dict(dict_):
  compressed = ''
  for key, data in dict_.items():
    compressed += f' {key}: {data},'
  return compressed

def get_txt_input(dict_):
  final_str = ''
  for key, data in dict_.items():
    if type(data) == list:
      temp_str = ''
      for subdict_ in data:
        temp_str += compress_dict(subdict_)
      final_str += f'{key}: {temp_str}\n'
    elif type(data) == dict:
      final_str += f'{key}: {compress_dict(data)}\n'
    else:
      final_str += f'{key}: {data}\n'

  print(f'user_data{final_str}')
  return final_str

#print(get_txt_input(sample_input_rus))
#print(get_txt_input(sample_input))

# requests



#completion = openai.ChatCompletion.create(
# model="gpt-3.5-turbo", 
# messages=[{"role": "user", "content": "What is the OpenAI mission?"}]
#)

#translation = openai.ChatCompletion.create(
#  model = "gpt-3.5-turbo", 
#  messages=[{"role": "user", "content": f"translate to english: {get_txt_input(sample_input_rus)}"}]
#)

import json

def get_json(user_info):
  print('user info accepted.')
  
  with open("gen/key.txt") as f:
    openai.api_key = f.read()

  with open("gen/questions.txt") as f:
    question = f.read()

  test_resp = openai.ChatCompletion.create(
    model = "gpt-3.5-turbo", 
    messages=[
      {"role": "system", "content": "You are a helpful assistant."},
      {"role": "user", "content": f"{question}\n {get_txt_input(user_info)}"}
      ]
  )

  contents = test_resp["choices"][0]["message"]["content"]

  print(contents)

  with open("gen/response_log.json", "w") as f:
    f.write(contents)

  
  response = ''
  
  with open("gen/response_log.json") as f:
    response = json.loads(f.read())

  return response

# translation

#import goslate

#gs = goslate.Goslate
#def translate(path):
#  with open(path) as f:
#    response = json.dumps(f.read())
#    response = gs.translate(response, 'en', 'rus' )
# return response

#print(translate("gen/response_log.json"))
# get slides

import collections
import collections.abc
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
import requests

def translate_text(text):
    url = f"https://translate.googleapis.com/translate_a/single?client=gtx&sl=en&tl=ru&dt=t&q={text}"
    response = requests.get(url)
    data = response.json()
    return data[0][0][0] + "."




def get_title_slides(user_info):
  
  global prs, slide_layout

  prs = Presentation()
  slide_layout = prs.slide_layouts[6]
  slide_layout2 = prs.slide_layouts[2]

  print(f'user info: {user_info}')
  response = get_json(user_info)
  print(f'api response: {response}')

  #with open("response_log.json") as f:
  # response = json.loads(f.read())

  #title slide

  first_slide_title = user_info["company_name"]

  prs.slide_height = Inches(9)    
  prs.slide_width = Inches(16)      
  first_slide = prs.slides.add_slide(slide_layout2)
  company_name = first_slide.shapes.title
  pitch_lable = first_slide.placeholders[1]

  company_name.text = first_slide_title
  pitch_lable.text = "pitch deck 2023"

  
  #presentation body
  
  for element in response["slides"]:
    title_text = translate_text(element["title"])
    discr = element["content"].values()

    #translation
    translated = []
    #print(discr)
    for element in discr:
      if type(element) == list:
        temp_str = ''
        for dic in element:
          temp_str += f"{translate_text(': '.join(dic.values()))}\n"
        translated.append(temp_str)
      elif element not in ['\n', ' ', ':', '.', ',']:
        translated.append(translate_text(element))
    discr = ''.join(translated)
    #print(discr)
    count = 0
    temp = ''
    string = ''
    for symbol in discr:
      if count == 0 and symbol == ' ':
        symbol = ''   
      if count >= 42 or symbol == ':':
        temp += symbol
        if (symbol in [' ', ':']):
          string += f'{temp}\n'
          temp = ''
          count = 0
        else:
          count += 1
      else:
        temp += symbol
        count += 1
    string += temp
    
    discr = string.replace('_', ' ').replace('.,', '.').replace('..', '. ')
    
    # split_list = discr.split(':')
    # final_list = []
    # for element in split_list:
    #   element += ':'
    #   final_list += element.split('.')
    # translated = []
    # for element in final_list:
    #   if element not in ['\n', ' ', ':', '.', ',']:
    #     translated.append(translate_text(element))
    #     print(translated)

    # discr = ''.join(translated).replace(':.', ': ').replace('.:', '').replace('.', '')


    size1 = 60
    if len(title_text) > 20:
      size1 = 50


    prs.slide_height = Inches(9)    
    prs.slide_width = Inches(16)      
    slide = prs.slides.add_slide(slide_layout)
    

    # width = Inches(prs.slide_width * 0.8)
    # height = Inches(prs.slide_height * 0.8)
    left = Inches(1) 
    top = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, prs.slide_width, prs.slide_height)
    txBox.word_wrap = True
    title_box = txBox.text_frame
    title_box.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_box.vertical_anchor = MSO_ANCHOR.TOP
    title_par = title_box.add_paragraph()
    title_par.word_wrap = True
    title_par.alignment = PP_ALIGN.LEFT
    title_par.text = title_text
 
    title_par.font.bold = True
    title_par.font.size = Pt(size1) 
    
    discr_par = title_box.add_paragraph()
    discr_par.alignment = PP_ALIGN.LEFT
    discr_par.text = discr
    discr_par.font.size = Pt(20) 


  return prs 


