'''
This LangChain Agent creates presentations from different contents.
'''

#%% ---------------------------------------------  IMPORTS  ----------------------------------------------------------#
import streamlit as st
import nbformat
from credentials import OPENAI_API_KEY
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.chat_models import ChatOpenAI
import tempfile
import os
from langchain.vectorstores import FAISS
from langchain.document_loaders import TextLoader
from main import rec_streamlit, speak_answer, get_transcript_whisper
import time
import en_core_web_sm
import spacy_streamlit
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.chains import RetrievalQA
from pptx import Presentation
import openai
from pptx.util import Inches


current_directory = os.path.dirname(os.path.abspath(__file__))
template_path = os.path.join(current_directory, "../templates/template.pptx")
output_path = os.path.join(current_directory, "../pres/output.pptx")

#%% ----------------------------------------  LANGCHAIN PRELOADS -----------------------------------------------------#
embeddings = OpenAIEmbeddings(disallowed_special=(), openai_api_key=OPENAI_API_KEY)
llm_doc = ChatOpenAI(openai_api_key=OPENAI_API_KEY, request_timeout=120)


# --------------------  SETTINGS  -------------------- #
st.set_page_config(page_title="Home", layout="wide")
st.markdown("""<style>.reportview-container .main .block-container {max-width: 95%;}</style>""", unsafe_allow_html=True)

# --------------------- HOME PAGE -------------------- #
st.title("PRESENTATION AGENT 📊")
st.write("""Use the power of LLMs with LangChain and OpenAI to scan through your Notebooks and create a Powerpoint presentation. 
        Find information and insight's with lightning speed. 🚀 Create new content with the support of state of the art language models and 
        and voice command your way through your documents. 🎙️""")


#%% ---------------------------------------  PREPROCESS DOCUMENTS ----------------------------------------------------#
def process_and_load_files(files, embeddings):
    """FAISS vector store for the documents"""
    loaded_files = []
    num_docs_processed = 0
    for file in files:
        try:
            # Save the uploaded file to a temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file.name.split('.')[-1]}") as temp_file:
                temp_file.write(file.getvalue())
                temp_file_path = temp_file.name

            if file_extension == ".ipynb":
                # Load up the notebook and extract the text
                notebook = nbformat.read(temp_file_path, as_version=4)
                loaded_files.extend([cell.source for cell in notebook.cells if cell.cell_type == "markdown"])
                num_docs_processed += len(loaded_files)
                print(f"Loaded {file.name}")

            else:
                # Load up the file as a doc and split
                loader = TextLoader(temp_file_path, encoding="utf-8")
                docs = loader.load_and_split()
                loaded_files.extend(docs)
                num_docs_processed += len(docs)
                print(f"Loaded {file.name}")

            # Delete the temporary file
            os.remove(temp_file_path)
        except Exception as e:
            print(f"Error loading {file.name}: {e}")
            pass

    if loaded_files:
        return FAISS.from_documents(loaded_files, embeddings), loaded_files, num_docs_processed
    else:
        return None, [], 0


#%% ---------------------------------------  CREATE PRESENTATION -----------------------------------------------------#
systemprompt = """You are a master at creating short paragraphs for presentations and you either summarize larger
                content that you get, or elaborate on a short topic that is given in a keyword."""


def run_GPT4(systemprompt, prompt, temperature):
    """Run GPT4 with the prompt and return the response"""
    openai.api_key = OPENAI_API_KEY
    completion = openai.ChatCompletion.create(
        model="gpt-4",
        temperature=temperature,
        messages=[
            {"role": "system", "content": systemprompt},
            {"role": "user", "content": prompt},
        ]
    )
    answer = completion.choices[0].message.content

    return answer


#%% ---------------------------------------  CREATE PRESENTATION -----------------------------------------------------#
def create_presentation_nb(notebook, template_path):
    prs = Presentation(template_path)
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[4]

    # Set the presentation title
    for cell in notebook.cells:
        if cell.cell_type == "markdown" and cell.source.startswith("# "):
            title_slide = prs.slides.add_slide(title_slide_layout)
            title = title_slide.shapes.title
            title.text = cell.source[2:].strip()
            break

    # Create slides for each second-level header
    for cell in notebook.cells:
        if cell.cell_type == "markdown" and cell.source.startswith("## "):
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            lines = cell.source.split('\n')
            title.text = lines[0][3:].strip()

            # Add content to the slide
            content = ""
            for line in lines[1:]:
                content += line.strip() + "\n"

            for subcell in notebook.cells[notebook.cells.index(cell) + 1:]:
                if subcell.cell_type == "markdown" and subcell.source.startswith("### "):
                    content += subcell.source[4:].strip() + "\n\n"
                elif subcell.cell_type == "markdown" and subcell.source.startswith("## "):
                    break
                elif subcell.cell_type == "markdown":
                    lines = subcell.source.split('\n')
                    for line in lines:
                        if not line.startswith("### "):
                            content += line.strip() + "\n"

            # Pass the content to GPT-4 for re-creation
            gpt4_content = run_GPT4(systemprompt, content, temperature=0.5)

            # Find the existing text shape on the slide and insert the GPT-4 generated content
            for shape in slide.shapes:
                if shape.has_text_frame and not shape.text.startswith("Click to edit"):
                    text_frame = shape.text_frame
                    text_frame.text = gpt4_content
                    break

    return prs




# ----------------- SIDE BAR SETTINGS ---------------- #
st.sidebar.subheader("Settings:")
tts_enabled = st.sidebar.checkbox("Enable Text-to-Speech", value=False)
ner_enabled = st.sidebar.checkbox("Enable NER in Response", value=False)

# ------------------ FILE UPLOADER ------------------- #
st.sidebar.subheader("File Uploader:")
uploaded_files = st.sidebar.file_uploader("Choose files", type=["csv", "html", "css", "py", "pdf", "ipynb"],
                                          accept_multiple_files=True)
st.sidebar.metric("Number of files uploaded", len(uploaded_files))
st.sidebar.color_picker("Pick a color for the answer space", "#C14531")

# Initialize docsearch as None
docsearch = None

# ------------------- FILE HANDLER ------------------- #
if uploaded_files:

    file_index = st.sidebar.selectbox("Select a file to display", options=[f.name for f in uploaded_files])
    selected_file = uploaded_files[[f.name for f in uploaded_files].index(file_index)]
    file_extension = selected_file.name.split(".")[-1]

    if file_extension in ["html", "css", "py"]:
        try:
            file_content = selected_file.getvalue().decode("utf-8")

            # --- Display the file content as code---
            with st.expander("Document Expander (Press button on the right to fold or unfold)", expanded=True):
                st.subheader("Uploaded Document:")
                st.code(file_content, language=file_extension)

        except Exception as e:
            st.write(f"Error reading {file_extension.upper()} file:", e)

    elif file_extension == "ipynb":
        try:
            nb_content = nbformat.read(selected_file, as_version=4)
            nb_filtered = [cell for cell in nb_content["cells"] if cell["cell_type"] in ["code", "markdown"]]
            nb_cell_content = [cell["source"] for cell in nb_filtered]


            # --- Display the file content as code---
            with st.expander("Document Expander (Press button on the right to fold or unfold)", expanded=True):
                st.subheader("Uploaded Document:")
                for cell in nb_filtered:
                    if cell["cell_type"] == "code":
                        st.code(cell["source"], language="python")
                    elif cell["cell_type"] == "markdown":
                        st.markdown(cell["source"])

                create_presentation_button = st.button("Create Presentation")

                # --------------------- PRESENTATION --------------------- #
                # Create presentation with the selected template
                if create_presentation_button:
                    prs = create_presentation_nb(nb_content, template_path)
                    prs.save(output_path)
                    st.success("Presentation created successfully!")

        except Exception as e:
            st.write(f"Error reading {file_extension.upper()} file:", e)


        # Filter HTML and Python files
        files = [f for f in uploaded_files if f.name.split(".")[-1] in ["html", "py", "ipynb", "css"]]


        # Process and load the HTML and Python files into the FAISS index
        docsearch, successfully_loaded_files, num_docs_processed = process_and_load_files(files, embeddings)
        st.write(f"You have {num_docs_processed} documents processed.")


# --------------------- USER INPUT --------------------- #
user_input = st.text_area("")
# If record button is pressed, rec_streamlit records and the output is saved
audio_bytes = rec_streamlit()


# ------------------- TRANSCRIPTION -------------------- #
if audio_bytes or user_input:

    if audio_bytes:
        try:
            with open("audio.wav", "wb") as file:
                file.write(audio_bytes)
        except Exception as e:
            st.write("Error recording audio:", e)
        transcript = get_transcript_whisper("audio.wav")
    else:
        transcript = user_input

    st.write("**Recognized:**")
    st.write(transcript)

    if any(word in transcript for word in ["abort recording"]):
        st.write("... Script stopped by user")
        exit()

    # ----------------------- ANSWER ----------------------- #
    with st.spinner("Fetching answer ..."):
        time.sleep(6)

    qa = RetrievalQA.from_chain_type(llm=llm_doc, chain_type="stuff", retriever=docsearch.as_retriever())
    answer = qa.run(transcript)
    st.write(answer)
    speak_answer(answer, tts_enabled)
    st.success("**Interaction finished**")

