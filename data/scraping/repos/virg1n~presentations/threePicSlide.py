import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openai
import config
from PIL import Image, ImageEnhance
import requests
from io import BytesIO
import numpy as np
from PIL import Image, ImageDraw
import os
import random
import cv2
import time
from googletrans import Translator
from functions import image1024, image256, imgToCircle, forBeautifullText, saveImg1024, saveImg256, question, getImgFromGoogle, forBeautifullTextBySimbols, saveImg512, cropPicture, cropPicture8x9, saveImg512WithDirectory, generateImage


# With Image Generation
def threePicSlideSDFirst(slide, topic, timecode, directory, textForSlide, topicForName):
    txBox = slide.shapes.add_textbox(Pt(0), Pt(50), width = Pt(1920), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(5), Pt(250), width = Pt(1920), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=64)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)

    slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png', Pt(120), Pt(530), width=Pt(493), height=Pt(493))
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png", Pt(713), Pt(530), width=Pt(493), height=Pt(493))
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\thirdFor2Pictures{topicForName}{timecode}.png", Pt(1306), Pt(530), width=Pt(493), height=Pt(493))


def threePicSlideSDSecond(slide, topic, timecode, directory, textForSlide, topicForName):
    txBox = slide.shapes.add_textbox(Pt(500), Pt(50), width = Pt(920), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(505), Pt(250), width = Pt(910), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=30)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)
    slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png', Pt(50), Pt(50), width=Pt(450), height=Pt(450))
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png", Pt(50), Pt(580), width=Pt(450), height=Pt(450))
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\thirdFor2Pictures{topicForName}{timecode}.png", Pt(1420), Pt(50), width=Pt(450), height=Pt(450))


def threePicSlideSDThird(slide, topic, timecode, directory, textForSlide, topicForName):
    txBox = slide.shapes.add_textbox(Pt(500), Pt(50), width = Pt(920), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(505), Pt(250), width = Pt(910), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=30)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)

    slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png', Pt(50), Pt(50), width=Pt(450), height=Pt(450))
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png", Pt(1420), Pt(580), width=Pt(450), height=Pt(450))
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\thirdFor2Pictures{topicForName}{timecode}.png", Pt(1420), Pt(50), width=Pt(450), height=Pt(450))


# With Images from google
def threePicSlideGoogleFirst(slide, topic, directory, textForSlide):
    txBox = slide.shapes.add_textbox(Pt(0), Pt(50), width = Pt(1920), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(5), Pt(250), width = Pt(1920), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=64)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg', Pt(120), Pt(530), width=Pt(493), height=Pt(493))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png', Pt(120), Pt(530), width=Pt(493), height=Pt(493))

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000003.jpg', Pt(713), Pt(530), width=Pt(493), height=Pt(493))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000003.png', Pt(713), Pt(530), width=Pt(493), height=Pt(493))

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.jpg', Pt(1306), Pt(530), width=Pt(493), height=Pt(493))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.png', Pt(1306), Pt(530), width=Pt(493), height=Pt(493))


def threePicSlideGoogleSecond(slide, topic, directory, textForSlide):
    txBox = slide.shapes.add_textbox(Pt(500), Pt(50), width = Pt(920), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(505), Pt(250), width = Pt(910), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=30)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg', Pt(50), Pt(50), width=Pt(450), height=Pt(450))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png', Pt(50), Pt(50), width=Pt(450), height=Pt(450))

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.jpg', Pt(50), Pt(580), width=Pt(450), height=Pt(450))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.png', Pt(50), Pt(580), width=Pt(450), height=Pt(450))

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000003.jpg', Pt(1420), Pt(50), width=Pt(450), height=Pt(450))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000003.png', Pt(1420), Pt(50), width=Pt(450), height=Pt(450))


def threePicSlideGoogleThird(slide, topic, directory, textForSlide):
    txBox = slide.shapes.add_textbox(Pt(500), Pt(50), width = Pt(920), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(505), Pt(250), width = Pt(910), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=30)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg', Pt(50), Pt(50), width=Pt(450), height=Pt(450))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png', Pt(50), Pt(50), width=Pt(450), height=Pt(450))

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.jpg', Pt(1420), Pt(580), width=Pt(450), height=Pt(450))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.png', Pt(1420), Pt(580), width=Pt(450), height=Pt(450))

    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000003.jpg', Pt(1420), Pt(50), width=Pt(450), height=Pt(450))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000003.png', Pt(1420), Pt(50), width=Pt(450), height=Pt(450))








