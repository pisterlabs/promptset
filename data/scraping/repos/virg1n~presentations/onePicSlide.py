import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openai
import config
from PIL import Image, ImageEnhance
import requests
from io import BytesIO
import numpy as np
from PIL import Image, ImageDraw
import os
import random
import cv2
import time
from googletrans import Translator
from functions import image1024, image256, imgToCircle, forBeautifullText, saveImg1024, saveImg256, question, getImgFromGoogle, forBeautifullTextBySimbols, saveImg512, cropPicture, cropPicture8x9, saveImg512WithDirectory, generateImage
from functions import darken_hex_color, hex_to_rgb, make_color_more_white, createReactangle, createDelay, createCircle
from pptx.enum.shapes import MSO_SHAPE


# With Image Generation
def onePicSlideSDFirst(slide, topic, timecode, directory, PromptForPictures, textForSlide ):
    
    txBox = slide.shapes.add_textbox(Pt(0), Pt(30), width = Pt(1920), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(135)
    generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures)
    # saveImg512WithDirectory(PromptForPictures, f"{str(topic[3:])}{str(timecode)}", dir=directory)
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(85), Pt(250), Pt(700), Pt(700))


    txBox = slide.shapes.add_textbox(Pt(810), Pt(350), width = Pt(1110), height = Pt(500))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=30)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(45)


def onePicSlideSDSecond(slide, topic, timecode, directory, PromptForPictures, textForSlide ):
    txBox = slide.shapes.add_textbox(Pt(0), Pt(30), width = Pt(1920), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(135)
    generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures)
    # saveImg512WithDirectory(PromptForPictures, f"{str(topic[3:])}{str(timecode)}", dir=directory)
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(1100), Pt(250), Pt(700), Pt(700))

    txBox = slide.shapes.add_textbox(Pt(20), Pt(350), width = Pt(1090), height = Pt(500))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=30)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(45)


def onePicSlideSDThird(slide, topic, timecode, directory, PromptForPictures, textForSlide ):
    txBox = slide.shapes.add_textbox(Pt(0), Pt(30), width = Pt(1920), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(135)
    generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures)
    # saveImg512WithDirectory(PromptForPictures, f"{str(topic[3:])}{str(timecode)}", dir=directory)
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(1100), Pt(250), Pt(700), Pt(700))

    txBox = slide.shapes.add_textbox(Pt(20), Pt(350), width = Pt(1090), height = Pt(500))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=30)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(45)


def onePicSlideSDFirstDocs(slide, topic, timecode, directory, PromptForPictures, textForSlide, BGcolor):
    shapes = slide.shapes
    createReactangle(left=0, top=0, height=1080, width=840, color=hex_to_rgb(BGcolor), shapes=shapes)


    generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures, negative_prompt="borders")
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(840), Pt(0), Pt(1080), Pt(1080))

    createReactangle(left=120, top=250, height=720, width=890, color=hex_to_rgb("#FFFFFF"), shapes=shapes)


    txBox = slide.shapes.add_textbox(Pt(120), Pt(270), width = Pt(890), height = Pt(720))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox = slide.shapes.add_textbox(Pt(120), Pt(470), width = Pt(890), height = Pt(720))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=33)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)


def onePicSlideSDSecondDocs(slide, topic, timecode, directory, PromptForPictures, textForSlide, BGcolor):
    shapes = slide.shapes
    createReactangle(left=0, top=0, height=1080, width=500, color=hex_to_rgb("#FFFFFF"), shapes=shapes)
    createReactangle(left=500, top=0, height=1080, width=1520, color=hex_to_rgb(BGcolor), shapes=shapes)


    generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures, negative_prompt="borders")
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(120), Pt(200), Pt(780), Pt(780))



    txBox = slide.shapes.add_textbox(Pt(950), Pt(270), width = Pt(970), height = Pt(150))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.size = Pt(100)


    txBox = slide.shapes.add_textbox(Pt(950), Pt(470), width = Pt(970), height = Pt(720))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=33)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.size = Pt(40)

# NEED TO CHANGE WIDTH AND HEIGHT
def onePicSlideSDThirdSquareDocs(slide, topic, timecode, directory, PromptForPictures, textForSlide, BGcolor):
    shapes = slide.shapes
    createReactangle(left=0, top=0, height=1080, width=1920, color=hex_to_rgb(BGcolor), shapes=shapes)
    createReactangle(left=100, top=100, height=880, width=700, color=hex_to_rgb("#FFFFFF"), shapes=shapes)


    generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures, negative_prompt="borders")
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(820), Pt(100), Pt(880), Pt(880))



    txBox = slide.shapes.add_textbox(Pt(100), Pt(200), width = Pt(700), height = Pt(150))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    try:
        p.text = forBeautifullTextBySimbols(string=topic, simbols=28)[3:]
    except:
        p.text = topic[3:]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox = slide.shapes.add_textbox(Pt(100), Pt(380), width = Pt(700), height = Pt(720))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=28)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)

# NEED TO CHANGE WIDTH AND HEIGHT
# def onePicSlideSDThirdWMoreHDocs(slide, topic, timecode, directory, PromptForPictures, textForSlide, BGcolor):
#     shapes = slide.shapes
#     createReactangle(left=0, top=0, height=1080, width=1920, color=hex_to_rgb(BGcolor), shapes=shapes)
#     createReactangle(left=100, top=100, height=880, width=700, color=hex_to_rgb("#FFFFFF"), shapes=shapes)


#     generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures, negative_prompt="borders", w=512, h=512)
#     slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(820), Pt(100), Pt(880), Pt(595))



#     txBox = slide.shapes.add_textbox(Pt(100), Pt(200), width = Pt(700), height = Pt(150))
#     tf = txBox.text_frame
#     p = tf.add_paragraph()
#     try:
#         p.text = forBeautifullTextBySimbols(string=topic, simbols=28)[3:]
#     except:
#         p.text = topic[3:]
#     p.alignment = PP_ALIGN.CENTER
#     p.font.name = 'Comic Sans'
#     p.font.size = Pt(100)


#     txBox = slide.shapes.add_textbox(Pt(100), Pt(380), width = Pt(700), height = Pt(720))
#     tf = txBox.text_frame
#     p = tf.add_paragraph()
#     p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=28)
#     p.alignment = PP_ALIGN.CENTER
#     p.font.name = 'Comic Sans'
#     p.font.size = Pt(40)


# 4
# {
#   Text:"asdasd",
    # Text-Margin : "2px"
# }
def onePicSlideSDFourthDocs(slide, topic, timecode, directory, PromptForPictures, textForSlide, BGcolor):

    generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures, negative_prompt="borders")
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(840), Pt(0), Pt(1080), Pt(1080))

    shapes = slide.shapes
    createDelay(left=-380, top=-220, radius=1500, color=hex_to_rgb("#FFFFFF"), shapes=shapes)



    txBox = slide.shapes.add_textbox(Pt(50), Pt(130), width = Pt(900), height = Pt(150))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    try:
        p.text = forBeautifullTextBySimbols(string=topic, simbols=16)[3:]
    except:
        p.text = topic[3:]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(120)


    txBox = slide.shapes.add_textbox(Pt(50), Pt(400), width = Pt(900), height = Pt(720))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=28)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(48)

# 5
def onePicSlideSDFifthDocs(slide, topic, timecode, directory, PromptForPictures, textForSlide, BGcolor):

    generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures, negative_prompt="borders, many pictures, text, letters")
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(0), Pt(0), Pt(1080), Pt(1080))

    shapes = slide.shapes
    createReactangle(left=1080, top=0, width=840, height=1080, color=hex_to_rgb("#FFFFFF"), shapes=shapes)
    createReactangle(left=800, top=300, width=1080, height=700, color=hex_to_rgb("#F1F1F1"), shapes=shapes)



    txBox = slide.shapes.add_textbox(Pt(800), Pt(305), width = Pt(1040), height = Pt(150))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    try:
        p.text = forBeautifullTextBySimbols(string=topic, simbols=18)[3:]
    except:
        p.text = topic[3:]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(110)


    txBox = slide.shapes.add_textbox(Pt(800), Pt(500), width = Pt(1040), height = Pt(720))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=33)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(44)


# NEED TO CHANGE WIDTH AND HEIGHT
def onePicSlideSDEighthDocs(slide, topic, timecode, directory, PromptForPictures, textForSlide, BGcolor):

    shapes = slide.shapes
    createReactangle(left=0, top=0, height=1100, width=1100, color=hex_to_rgb("#FFFFFF"), shapes=shapes)
    createReactangle(left=80, top=80, height=940, width=940, color=hex_to_rgb('#313138'), shapes=shapes)


    generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures+" nearly #313138 colors, main content on left side", negative_prompt="borders")
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(1100), Pt(0), Pt(1080), Pt(1080))



    txBox = slide.shapes.add_textbox(Pt(100), Pt(270), width = Pt(900), height = Pt(150))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    try:
        p.text = forBeautifullTextBySimbols(string=topic, simbols=18)[3:]
    except:
        p.text = topic[3:]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.size = Pt(100)


    txBox = slide.shapes.add_textbox(Pt(100), Pt(470), width = Pt(900), height = Pt(720))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=33)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.size = Pt(40)


# 10
def onePicSlideSDTenthDocs(slide, topic, timecode, directory, PromptForPictures, textForSlide, BGcolor):

    shapes = slide.shapes
    # createReactangle(left=80, top=80, height=940, width=940, color=hex_to_rgb('#313138'), shapes=shapes)


    generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures+", main content on left side", negative_prompt="borders")
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(0), Pt(0), Pt(1080), Pt(1080))

    createCircle(left=800, top=-200, radius=1500, color=hex_to_rgb("#FFFFFF"), shapes=shapes)


    txBox = slide.shapes.add_textbox(Pt(900), Pt(190), width = Pt(900), height = Pt(150))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    try:
        p.text = forBeautifullTextBySimbols(string=topic, simbols=18)[3:]
    except:
        p.text = topic[3:]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(140)


    txBox = slide.shapes.add_textbox(Pt(900), Pt(470), width = Pt(900), height = Pt(720))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=33)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(50)

# 11
def onePicSlideSDEleventhDocs(slide, topic, timecode, directory, PromptForPictures, textForSlide, BGcolor):

    shapes = slide.shapes

    
    createReactangle(left=0, top=0, height=1080, width=1920, color=hex_to_rgb(make_color_more_white(BGcolor, 70)), shapes=shapes)
    createReactangle(left=70, top=100, height=680, width=1780, color=hex_to_rgb('#D3D2D1'), shapes=shapes)
    createReactangle(left=70, top=780, height=200, width=1780, color=hex_to_rgb('#EBE9EB'), shapes=shapes)
    createReactangle(left=140, top=60, height=600, width=600, color=hex_to_rgb('FFFFFF'), shapes=shapes)

    generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures, negative_prompt="borders")
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(165), Pt(85), Pt(550), Pt(550))

    # createReactangle(left=800, top=0, height=200, width=50, color=hex_to_rgb(BGcolor), shapes=shapes)
    # createCircle(left=800, top=-200, radius=1500, color=hex_to_rgb("#FFFFFF"), shapes=shapes)


    txBox = slide.shapes.add_textbox(Pt(830), Pt(110), width = Pt(900), height = Pt(150))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    try:
        p.text = forBeautifullTextBySimbols(string=topic, simbols=18)[3:]
    except:
        p.text = topic[3:]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(130)


    txBox = slide.shapes.add_textbox(Pt(830), Pt(340), width = Pt(900), height = Pt(720))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=33)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(44)


# 12
# def onePicSlideSDTwelvethDocs(slide, topic, timecode, directory, PromptForPictures, textForSlide, BGcolor):

#     shapes = slide.shapes

#     createReactangle(left=0, top=0, height=1080, width=1920, color=hex_to_rgb('#f5f5f5'), shapes=shapes)
#     createReactangle(left=25, top=25, height=1055, width=650, color=hex_to_rgb(make_color_more_white(BGcolor, 90)), shapes=shapes)

#     generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures, negative_prompt="borders")
#     slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(725), Pt(265), Pt(550), Pt(550))


#     txBox = slide.shapes.add_textbox(Pt(830), Pt(110), width = Pt(900), height = Pt(150))
#     tf = txBox.text_frame
#     p = tf.add_paragraph()
#     try:
#         p.text = forBeautifullTextBySimbols(string=topic, simbols=18)[3:]
#     except:
#         p.text = topic[3:]
#     p.alignment = PP_ALIGN.CENTER
#     p.font.name = 'Comic Sans'
#     p.font.size = Pt(130)


#     txBox = slide.shapes.add_textbox(Pt(830), Pt(340), width = Pt(900), height = Pt(720))
#     tf = txBox.text_frame
#     p = tf.add_paragraph()
#     p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=33)
#     p.alignment = PP_ALIGN.CENTER
#     p.font.name = 'Comic Sans'
#     p.font.size = Pt(44)



def onePicSlideSDThirteenthDocs(slide, topic, timecode, directory, PromptForPictures, textForSlide, BGcolor):
    shapes = slide.shapes
    # createReactangle(left=0, top=0, height=1080, width=1920, color=hex_to_rgb(make_color_more_white(BGcolor, 70)), shapes=shapes)
    createReactangle(left=165, top=800, height=50, width=700, color=hex_to_rgb(make_color_more_white(BGcolor,80)), shapes=shapes)
    # createReactangle(left=70, top=780, height=200, width=1780, color=hex_to_rgb('#EBE9EB'), shapes=shapes)
    # createReactangle(left=140, top=60, height=600, width=600, color=hex_to_rgb('FFFFFF'), shapes=shapes)

    generateImage(f"{str(topic[3:])}{str(timecode)}", dir=directory, prompt=PromptForPictures, negative_prompt="borders")
    slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png", Pt(100), Pt(100), Pt(700), Pt(700))

    # createReactangle(left=800, top=0, height=200, width=50, color=hex_to_rgb(BGcolor), shapes=shapes)
    # createCircle(left=800, top=-200, radius=1500, color=hex_to_rgb("#FFFFFF"), shapes=shapes)


    txBox = slide.shapes.add_textbox(Pt(830), Pt(110), width = Pt(900), height = Pt(150))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    try:
        p.text = forBeautifullTextBySimbols(string=topic, simbols=18)[3:]
    except:
        p.text = topic[3:]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(130)


    txBox = slide.shapes.add_textbox(Pt(830), Pt(340), width = Pt(900), height = Pt(720))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=33)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(44)













# With Images from google
def onePicSlideGoogleWlowerHFirst(slide, topic, directory, textForSlide):
    txBox = slide.shapes.add_textbox(Pt(0), Pt(120), width = Pt(1120), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(5), Pt(370), width = Pt(1110), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=40)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)
    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg', Pt(1110), Pt(0), width=Pt(800), height=Pt(1080))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png', Pt(1110), Pt(0), width=Pt(800), height=Pt(1080))


def onePicSlideGoogleWlowerHSecond(slide, topic, directory, textForSlide):
    txBox = slide.shapes.add_textbox(Pt(800), Pt(120), width = Pt(1100), height = Pt(100)) #Pt(1100), Pt(0), width=Pt(820), height=Pt(1080)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(800), Pt(370), width = Pt(1080), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=33)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)
    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg', Pt(0), Pt(0), width=Pt(800), height=Pt(1080))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png', Pt(0), Pt(0), width=Pt(800), height=Pt(1080))


def onePicSlideGoogleWequalHFirst(slide, topic, directory, textForSlide):
    txBox = slide.shapes.add_textbox(Pt(0), Pt(120), width = Pt(750), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(5), Pt(370), width = Pt(740), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=27)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)
    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg', Pt(750), Pt(0), width=Pt(1170), height=Pt(1080))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png', Pt(750), Pt(0), width=Pt(1170), height=Pt(1080))


def onePicSlideGoogleWequalHSecond(slide, topic, directory, textForSlide):
    txBox = slide.shapes.add_textbox(Pt(1170), Pt(120), width = Pt(750), height = Pt(100)) #Pt(750), Pt(0), width=Pt(1170), height=Pt(1080)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(100)


    txBox1 = slide.shapes.add_textbox(Pt(1175), Pt(370), width = Pt(740), height = Pt(500))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=27)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)
    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg', Pt(0), Pt(0), width=Pt(1170), height=Pt(1080))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png', Pt(0), Pt(0), width=Pt(1170), height=Pt(1080))


def onePicSlideGoogleWmoreHFirst(slide, topic, directory, textForSlide):
    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg', Pt(0), Pt(0), width=Pt(1920), height=Pt(1080))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png', Pt(0), Pt(0), width=Pt(1920), height=Pt(1080))
    rectangle = slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\AI\\rectangle2.png", Pt(453), Pt(610), width=Pt(1066), height=Pt(450))
    txBox = slide.shapes.add_textbox(Pt(453), Pt(614), width = Pt(1066), height = Pt(450))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(65)

    txBox1 = slide.shapes.add_textbox(Pt(453), Pt(690), width = Pt(1066), height = Pt(450))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=37)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)


def onePicSlideGoogleWmoreHSecond(slide, topic, directory, textForSlide):
    try:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg', Pt(453), Pt(460), width=Pt(1066), height=Pt(600))
    except:
        slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png', Pt(453), Pt(460), width=Pt(1066), height=Pt(600))
    txBox = slide.shapes.add_textbox(Pt(0), Pt(5), width = Pt(1900), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic.split(' ')[-1]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(60)


    txBox1 = slide.shapes.add_textbox(Pt(5), Pt(107), width = Pt(1900), height = Pt(460))
    tf1 = txBox1.text_frame
    p = tf1.add_paragraph()
    p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=60)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(40)

