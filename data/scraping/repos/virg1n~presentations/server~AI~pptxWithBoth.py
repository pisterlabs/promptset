import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openai
import config
from PIL import Image, ImageEnhance
import requests
from io import BytesIO
import PyPDF2 
import re
import tiktoken
import numpy as np
from PIL import Image, ImageDraw
import os
import random
import cv2
import time
from googletrans import Translator
from functions import image1024, image256, imgToCircle, forBeautifullText, saveImg1024, saveImg256, question, getImgFromGoogle, forBeautifullTextBySimbols, saveImg512, cropPicture, cropPicture8x9, saveImg512WithDirectory, generateImage
from onePicSlide import onePicSlideSDFirst, onePicSlideSDSecond, onePicSlideGoogleWlowerHFirst, onePicSlideGoogleWlowerHSecond, onePicSlideGoogleWequalHFirst, onePicSlideGoogleWequalHSecond, onePicSlideGoogleWmoreHFirst, onePicSlideGoogleWmoreHSecond
from twoPicSlide import twoPicSlideSDFirst, twoPicSlideSDSecond, twoPicSlideSDThird, twoPicSlideGoogleFirstThreeFirst, twoPicSlideGoogleFirstThreeSecond, twoPicSlideGoogleFirstThreeThird, twoPicSlideGoogleSecondOneFirst
from threePicSlide import threePicSlideSDFirst, threePicSlideSDSecond, threePicSlideSDThird, threePicSlideGoogleFirst, threePicSlideGoogleSecond, threePicSlideGoogleThird


# openai.api_key = "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"
# openai.api_key = "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"
prs = Presentation()
prs.slide_height = Pt(1080)
prs.slide_width = Pt(1920)

translator = Translator()

def firstSlide(slide, topic, usersname):
    top = Pt(400)
    width = Pt(1920)
    left = Pt(0)
    height = Pt(150)
    txBox = slide.shapes.add_textbox(left, top, width, height)# (left, top, width, height)

    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = topic
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans' #Arial Rounded MT Bold
    p.font.size = Pt(130)
    p = tf.add_paragraph()
    p.text = "Prepared by " + usersname
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60)
    p.font.name = 'Comic Sans'
    p.font.color.rgb = RGBColor(0, 0, 0)
    print("1st Slide done")


def generatePlan(slide, themes, language):
    txBox = slide.shapes.add_textbox(Pt(20), Pt(35), width = Pt(1920), height = Pt(100))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    if language != 'en':
        p.text = translator.translate("Plan", dest=language).text
    else:
        p.text = "Plan"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Comic Sans'
    p.font.size = Pt(120)
    txBox1 = slide.shapes.add_textbox(Pt(20), Pt(280), width = Pt(1920), height = Pt(150))
    tf1 = txBox1.text_frame
    if language == "en":
        for i in range(0, len(themes)):
            p = tf1.add_paragraph()
            p.text = str(themes[i][0])
            p.font.size = Pt(90)
    else:
        for i in range(0, len(themes)):
            p = tf1.add_paragraph()
            p.text = translator.translate(str(themes[i][0]), dest=language).text
            p.font.size = Pt(90)
    print("2nd Slide done")

def generateDeffaultSlide(slide, quantityOfPictures, methodOfPictures, PromptForPictures, textForSlide):
    if quantityOfPictures == 1:
        generate1PicturesSlide(slide=slide, quantityOfPictures=1, methodOfPictures=methodOfPictures, PromptForPictures=PromptForPictures, textForSlide=textForSlide, varitation=1)
    if quantityOfPictures == 2:
        generate2PicturesSlide(slide=slide, quantityOfPictures=2, methodOfPictures=methodOfPictures, PromptForPictures=PromptForPictures, textForSlide=textForSlide)
    if quantityOfPictures == 3:
        generate3PicturesSlide(slide=slide, quantityOfPictures=3, methodOfPictures=methodOfPictures, PromptForPictures=PromptForPictures, textForSlide=textForSlide)
    if quantityOfPictures >= 4:
        generate4PicturesSlide(slide=slide, quantityOfPictures=4, methodOfPictures=methodOfPictures, PromptForPictures=PromptForPictures, textForSlide=textForSlide)


def getImage(directory, nameOfImage):
    try:
        image = cv2.imread(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{nameOfImage}.jpg')
        return image
    except:
        image = cv2.imread(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{nameOfImage}.png')
        return image


def delete3ImagesWithTimecode(usersname, timecode):
    try:
        os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000001.jpg")
    except:
        os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000001.png")

    try:
        os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000002.jpg")
    except:
        os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000002.png")

    try:
        os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000003.jpg")
    except:
        os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000003.png")


def delete3Images(directory):
    try:
        os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg")
    except:
        os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.png")

    try:
        os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.jpg")
    except:
        os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000002.png")

    try:
        os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000003.jpg")
    except:
        os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000003.png")



def generate1PicturesSlide(slide, methodOfPictures, PromptForPictures, textForSlide, topic, varitation, usersname, timecode, directory, prompt2):
    print(methodOfPictures)
    if methodOfPictures == "dalle":
        if random.randint(1,2) == 1:
            onePicSlideSDFirst(slide=slide, topic=topic, timecode=timecode, directory=directory, PromptForPictures=PromptForPictures, textForSlide=textForSlide)
        else:
            onePicSlideSDSecond(slide=slide, topic=topic, timecode=timecode, directory=directory, PromptForPictures=PromptForPictures, textForSlide=textForSlide)    

        # Delete Images
        os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\{str(topic[3:])}{str(timecode)}.png")
    
    else:
        getImgFromGoogle(prompt=PromptForPictures, dir=directory)
        print(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\000001.jpg')
        h, w, c = getImage(directory=directory, nameOfImage="000001").shape

        # Rare
        if w/h < 1:
            if random.randint(1,2) == 1:
                onePicSlideGoogleWlowerHFirst(slide=slide, topic=topic, directory=directory, textForSlide=textForSlide)
            else:
                onePicSlideGoogleWlowerHSecond(slide=slide, topic=topic, directory=directory, textForSlide=textForSlide)

        # Main
        elif (w/h - 1.083 <= 0.35):
            if random.randint(1,2) == 1:
                onePicSlideGoogleWequalHFirst(slide=slide, topic=topic, directory=directory, textForSlide=textForSlide)
            else:
                onePicSlideGoogleWequalHSecond(slide=slide, topic=topic, directory=directory, textForSlide=textForSlide)

        # Rare
        else:
            if random.randint(1,3) != 3:
                onePicSlideGoogleWmoreHFirst(slide=slide, topic=topic, directory=directory, textForSlide=textForSlide)
            else:
                onePicSlideGoogleWmoreHSecond(slide=slide, topic=topic, directory=directory, textForSlide=textForSlide)

        # Delete Images
        try:
            os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000001.jpg")
        except:
            os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000001.png")


def generate2PicturesSlide(slide, methodOfPictures, PromptForPictures, textForSlide, topic, varitation, usersname, timecode, directory, prompt2):
    if methodOfPictures == "dalle":
        k = random.randint(1,3)
        if k == 1:
            twoPicSlideSDFirst(slide=slide, topic=topic, timecode=timecode, directory=directory, PromptForPictures=PromptForPictures,prompt2=prompt2, textForSlide=textForSlide)
            topicForName = topic.replace(' ', '').replace('.', '')
            # topicForName = topic.replace(' ', '').replace('.', '')
            # generateImage(f"firstFor2Pictures{topicForName}{timecode}", dir=directory, prompt=PromptForPictures)
            # generateImage(f"secondFor2Pictures{topicForName}{timecode}", dir=directory, prompt=prompt2)


            # # saveImg512WithDirectory(PromptForPictures, f"firstFor2Pictures{topicForName}{timecode}", dir=directory)
            # # saveImg512WithDirectory(prompt2, f"secondFor2Pictures{topicForName}{timecode}", dir=directory)
            # txBox = slide.shapes.add_textbox(Pt(0), Pt(50), width = Pt(1920), height = Pt(100))
            # tf = txBox.text_frame
            # p = tf.add_paragraph()
            # p.text = topic.split(' ')[-1]
            # p.alignment = PP_ALIGN.CENTER
            # p.font.name = 'Comic Sans'
            # p.font.size = Pt(100)


            # txBox1 = slide.shapes.add_textbox(Pt(5), Pt(250), width = Pt(1920), height = Pt(500))
            # tf1 = txBox1.text_frame
            # p = tf1.add_paragraph()
            # p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=65)
            # p.alignment = PP_ALIGN.CENTER
            # p.font.name = 'Comic Sans'
            # p.font.size = Pt(40)

            # slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png', Pt(200), Pt(530), width=Pt(512), height=Pt(512))
            # slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png', Pt(1008), Pt(530), width=Pt(512), height=Pt(512))
        elif k == 2:
            twoPicSlideSDSecond(slide=slide, topic=topic, timecode=timecode, directory=directory, PromptForPictures=PromptForPictures,prompt2=prompt2, textForSlide=textForSlide)
            topicForName = topic.replace(' ', '').replace('.', '')
            # topicForName = topic.replace(' ', '').replace('.', '')
            # generateImage(f"firstFor2Pictures{topicForName}{timecode}", dir=directory, prompt=PromptForPictures)
            # generateImage(f"secondFor2Pictures{topicForName}{timecode}", dir=directory, prompt=prompt2)

            # # saveImg512WithDirectory(PromptForPictures, f"firstFor2Pictures{topicForName}{timecode}", dir=directory)
            # # saveImg512WithDirectory(prompt2, f"secondFor2Pictures{topicForName}{timecode}", dir=directory)

            # txBox = slide.shapes.add_textbox(Pt(0), Pt(200), width = Pt(1290), height = Pt(100))
            # tf = txBox.text_frame
            # p = tf.add_paragraph()
            # p.text = topic.split(' ')[-1]
            # p.alignment = PP_ALIGN.CENTER
            # p.font.name = 'Comic Sans'
            # p.font.size = Pt(100)


            # txBox1 = slide.shapes.add_textbox(Pt(10), Pt(400), width = Pt(1270), height = Pt(500))
            # tf1 = txBox1.text_frame
            # p = tf1.add_paragraph()
            # p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=40)
            # p.alignment = PP_ALIGN.CENTER
            # p.font.name = 'Comic Sans'
            # p.font.size = Pt(50)

            # slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png', Pt(1290), Pt(50), width=Pt(475), height=Pt(475))
            # slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png', Pt(1290), Pt(545), width=Pt(475), height=Pt(475))
        else:
            twoPicSlideSDThird(slide=slide, topic=topic, timecode=timecode, directory=directory, PromptForPictures=PromptForPictures,prompt2=prompt2, textForSlide=textForSlide)
            topicForName = topic.replace(' ', '').replace('.', '')
            # topicForName = topic.replace(' ', '').replace('.', '')
            # generateImage(f"firstFor2Pictures{topicForName}{timecode}", dir=directory, prompt=PromptForPictures)
            # generateImage(f"secondFor2Pictures{topicForName}{timecode}", dir=directory, prompt=prompt2)


            # # saveImg512WithDirectory(PromptForPictures, f"firstFor2Pictures{topicForName}{timecode}", dir=directory)
            # # saveImg512WithDirectory(prompt2, f"secondFor2Pictures{topicForName}{timecode}", dir=directory)
            # txBox = slide.shapes.add_textbox(Pt(630), Pt(200), width = Pt(1290), height = Pt(100))
            # tf = txBox.text_frame
            # p = tf.add_paragraph()
            # p.text = topic.split(' ')[-1]
            # p.alignment = PP_ALIGN.CENTER
            # p.font.name = 'Comic Sans'
            # p.font.size = Pt(100)


            # txBox1 = slide.shapes.add_textbox(Pt(640), Pt(400), width = Pt(1270), height = Pt(500))
            # tf1 = txBox1.text_frame
            # p = tf1.add_paragraph()
            # p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=40)
            # p.alignment = PP_ALIGN.CENTER
            # p.font.name = 'Comic Sans'
            # p.font.size = Pt(50)

            # slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png', Pt(155), Pt(50), width=Pt(475), height=Pt(475))
            # slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png', Pt(155), Pt(545), width=Pt(475), height=Pt(475))
        os.remove(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png')
        os.remove(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png')
    else:
        k = random.randint(1,4)
        if k != 4:
            getImgFromGoogle(prompt=PromptForPictures, dir=directory)
            cropPicture(path=f"{directory}\\000001")

            getImgFromGoogle(prompt=prompt2, dir=directory)
            cropPicture(path=f"{directory}\\000002")

            if k == 1:
                twoPicSlideGoogleFirstThreeFirst(slide=slide, topic=topic, textForSlide=textForSlide, directory=directory)
            
            elif k == 2:
                twoPicSlideGoogleFirstThreeSecond(slide=slide, topic=topic, textForSlide=textForSlide, directory=directory)

            elif k == 3:
                twoPicSlideGoogleFirstThreeThird(slide=slide, topic=topic, textForSlide=textForSlide, directory=directory)

        else:
            getImgFromGoogle(prompt=PromptForPictures, dir=directory)
            cropPicture8x9(path=f"{directory}\\000001")

            getImgFromGoogle(prompt=prompt2, dir=directory)
            cropPicture8x9(path=f"{directory}\\000002")

            twoPicSlideGoogleSecondOneFirst(slide=slide, topic=topic, timecode=timecode, directory=directory, PromptForPictures=PromptForPictures,prompt2=prompt2, textForSlide=textForSlide)
        
        delete3ImagesWithTimecode(usersname=usersname, timecode=timecode)
        
        # try:
        #     os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000001.jpg")
        # except:
        #     os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000001.png")

        # try:
        #     os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000002.jpg")
        # except:
        #     os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000002.png")

        # try:
        #     os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000003.jpg")
        # except:
        #     os.remove(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{str(usersname.replace(' ', ''))}{str(timecode)}\\000003.png")


def generate3PicturesSlide(slide, methodOfPictures, PromptForPictures, textForSlide, topic, varitation, usersname, timecode, directory, prompt2):
    if methodOfPictures == "dalle":
        k = random.randint(1,3)
        topicForName = topic.replace(' ', '').replace('.', '')

        generateImage(f"firstFor2Pictures{topicForName}{timecode}", dir=directory, prompt=PromptForPictures)
        generateImage(f"secondFor2Pictures{topicForName}{timecode}", dir=directory, prompt=prompt2)
        generateImage(f"thirdFor2Pictures{topicForName}{timecode}", dir=directory, prompt=prompt2)

        # saveImg512WithDirectory(PromptForPictures, f"firstFor2Pictures{topicForName}{timecode}", dir=directory)
        # saveImg512WithDirectory(prompt2, f"secondFor2Pictures{topicForName}{timecode}", dir=directory)
        # saveImg512WithDirectory(prompt2, f"thirdFor2Pictures{topicForName}{timecode}", dir=directory)
        if k == 1:
            threePicSlideSDFirst(slide=slide, topic=topic, timecode=timecode, directory=directory, textForSlide=textForSlide, topicForName=topicForName)

            # txBox = slide.shapes.add_textbox(Pt(0), Pt(50), width = Pt(1920), height = Pt(100))
            # tf = txBox.text_frame
            # p = tf.add_paragraph()
            # p.text = topic
            # p.alignment = PP_ALIGN.CENTER
            # p.font.name = 'Comic Sans'
            # p.font.size = Pt(100)


            # txBox1 = slide.shapes.add_textbox(Pt(5), Pt(250), width = Pt(1920), height = Pt(500))
            # tf1 = txBox1.text_frame
            # p = tf1.add_paragraph()
            # p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=64)
            # p.alignment = PP_ALIGN.CENTER
            # p.font.name = 'Comic Sans'
            # p.font.size = Pt(40)

            # slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png', Pt(120), Pt(530), width=Pt(493), height=Pt(493))
            # slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png", Pt(713), Pt(530), width=Pt(493), height=Pt(493))
            # slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\thirdFor2Pictures{topicForName}{timecode}.png", Pt(1306), Pt(530), width=Pt(493), height=Pt(493))
        elif k == 2:
            threePicSlideSDSecond(slide=slide, topic=topic, timecode=timecode, directory=directory, textForSlide=textForSlide, topicForName=topicForName)

            # txBox = slide.shapes.add_textbox(Pt(500), Pt(50), width = Pt(920), height = Pt(100))
            # tf = txBox.text_frame
            # p = tf.add_paragraph()
            # p.text = topic.split(' ')[-1]
            # p.alignment = PP_ALIGN.CENTER
            # p.font.name = 'Comic Sans'
            # p.font.size = Pt(100)


            # txBox1 = slide.shapes.add_textbox(Pt(505), Pt(250), width = Pt(910), height = Pt(500))
            # tf1 = txBox1.text_frame
            # p = tf1.add_paragraph()
            # p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=30)
            # p.alignment = PP_ALIGN.CENTER
            # p.font.name = 'Comic Sans'
            # p.font.size = Pt(40)
            # slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png', Pt(50), Pt(50), width=Pt(450), height=Pt(450))
            # slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png", Pt(50), Pt(580), width=Pt(450), height=Pt(450))
            # slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\thirdFor2Pictures{topicForName}{timecode}.png", Pt(1420), Pt(50), width=Pt(450), height=Pt(450))

        else:
            threePicSlideSDThird(slide=slide, topic=topic, timecode=timecode, directory=directory, textForSlide=textForSlide, topicForName=topicForName)

            # txBox = slide.shapes.add_textbox(Pt(500), Pt(50), width = Pt(920), height = Pt(100))
            # tf = txBox.text_frame
            # p = tf.add_paragraph()
            # p.text = topic.split(' ')[-1]
            # p.alignment = PP_ALIGN.CENTER
            # p.font.name = 'Comic Sans'
            # p.font.size = Pt(100)


            # txBox1 = slide.shapes.add_textbox(Pt(505), Pt(250), width = Pt(910), height = Pt(500))
            # tf1 = txBox1.text_frame
            # p = tf1.add_paragraph()
            # p.text = forBeautifullTextBySimbols(string=textForSlide, simbols=30)
            # p.alignment = PP_ALIGN.CENTER
            # p.font.name = 'Comic Sans'
            # p.font.size = Pt(40)

            # slide.shapes.add_picture(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png', Pt(50), Pt(50), width=Pt(450), height=Pt(450))
            # slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png", Pt(1420), Pt(580), width=Pt(450), height=Pt(450))
            # slide.shapes.add_picture(f"C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\thirdFor2Pictures{topicForName}{timecode}.png", Pt(1420), Pt(50), width=Pt(450), height=Pt(450))
        os.remove(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\firstFor2Pictures{topicForName}{timecode}.png')
        os.remove(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\secondFor2Pictures{topicForName}{timecode}.png')
        os.remove(f'C:\\Users\\Bogdan\\java\\pyth\\t2p\\presentations\\server\\imgs\\{directory}\\thirdFor2Pictures{topicForName}{timecode}.png')
    else:
        getImgFromGoogle(prompt=PromptForPictures, dir=directory, max_num=2)
        cropPicture(path=f"{directory}\\000001")
        cropPicture(path=f"{directory}\\000002")

        getImgFromGoogle(prompt=prompt2, dir=directory)
        cropPicture(path=f"{directory}\\000003")

        k = random.randint(1, 3)
        if k == 1:
            threePicSlideGoogleFirst(slide=slide, topic=topic, directory=directory, textForSlide=textForSlide)

        elif k == 2:
            threePicSlideGoogleSecond(slide=slide, topic=topic, directory=directory, textForSlide=textForSlide)

        else:
            threePicSlideGoogleThird(slide=slide, topic=topic, directory=directory, textForSlide=textForSlide)

        delete3Images(directory=directory)


def generate4PicturesSlide(slide, methodOfPictures, PromptForPictures, textForSlide, topic, varitation, usersname, timecode, directory):
    pass



def GeneratePresentationWith2TypesOfPictures(topic, numOfSlides, font, usersname, plan=True, interestFact=True, timecode=1):
    print("start presentation")
    themes = []
    finalThemesFromPrompt = []
    language = translator.detect(topic).lang
    # if language != "en":
    #     themesFromPrompt = question(f"write to me like above a {numOfSlides-2} topics for {numOfSlides-2} slides about {translator.translate(topic).text}, without introducing and conclusion slide with 2-5 words. say about every slide what is better search pictures in network or generate with dalle. after topic write only 1 word: 'dalle' and write better prompt for falle for this picture devided by '-' if dalle is better. Write 'network' and write better querry for this picture devided by '-' if network is better to find picture for this slide. say about every slide how much picture is needed for this slide with only one integer digit devided by ' - ' in interval from 0 to 3. Then for every slide generate text for this slide about particular topic in 50 words devided by ' - '. ").replace('\n', ' - ').replace('"', '').split(' - ')
    # else:
    #     themesFromPrompt = question(f"write to me like above a {numOfSlides-2} topics for {numOfSlides-2} slides about {topic}, without introducing and conclusion slide with 2-5 words. say about every slide what is better search pictures in network or generate with dalle. after topic write only 1 word: 'dalle' and write better prompt for falle for this picture devided by '-' if dalle is better. Write 'network' and write better querry for this picture devided by '-' if network is better to find picture for this slide. say about every slide how much picture is needed for this slide with only one integer digit devided by ' - ' in interval from 0 to 3. Then for every slide generate text for this slide about particular topic in 50 words devided by ' - '. ").replace('\n', ' - ').replace('"', '').split(' - ')
    # for i in range(len(themesFromPrompt)):
    #     if themesFromPrompt[i] != '':
    #         finalThemesFromPrompt.append(themesFromPrompt[i])
    # # print(finalThemesFromPrompt)
    # # finalThemesFromPrompt = ['Colors', 'network', 'Rainbow colors', '2', 'An image showing the colors of the rainbow in order, from red to violet', 'The colors of the rainbow are a beautiful natural phenomenon. This image displays the colors in order, from red to violet, and can be used to teach children about the colors of the rainbow. ', 'Formation', 'dalle', 'Rainbow formation', '1', 'A picture of a rainbow forming in the sky, with clear distinction of the colors', 'Rainbows are formed when sunlight is refracted and reflected by water droplets in the air. This picture captures the moment of a rainbow forming in the sky, with clear distinction of the colors. ', 'Mythology', 'network', 'Rainbow mythology', '3', 'An image of a painting or artwork depicting a rainbow in mythology, such as the Norse Bifrost or the Greek Iris', 'Rainbows have been a part of mythology and folklore for centuries. This image shows a painting or artwork depicting a rainbow in mythology, such as the Norse Bifrost or the Greek Iris. ', 'Double Rainbow', 'dalle', 'Double rainbow', '2', 'A picture of a double rainbow, with both rainbows clearly visible', 'Double rainbows are a rare and beautiful sight. This picture captures the moment of a double rainbow, with both rainbows clearly visible and the colors vibrant and distinct. ', 'Scientific Explanation', 'network', 'Rainbow diagram', '1', 'A diagram or illustration explaining the scientific process behind the formation of a rainbow', 'Rainbows are a result of the refraction and reflection of light through water droplets in the air. This diagram or illustration explains the scientific process behind the formation of a rainbow, including the angles of refraction and reflection.']
    # for i in range(0, len(finalThemesFromPrompt), 6):
    #     themes.append(finalThemesFromPrompt[i:i+6])
    # print(themes)
    # [['1. Formation', 'dalle', '"coal formation"', '3', 'Coal is formed from the remains of ancient plants that were buried and subjected to high pressure and heat over millions of years. The process of coal formation is called coalification.'], ['2. Types', 'network', '"types of coal"', '4', 'There are four main types of coal: anthracite, bituminous, subbituminous, and lignite. Each type has different properties and uses, with anthracite being the highest quality and lignite being the lowest.'], ['3. Mining', 'dalle', '"coal mining"', '2', "Coal mining involves extracting coal from the earth's surface or underground. It can be a dangerous and environmentally damaging process, but it is also a vital source of energy for many countries."], ['4. Uses', 'network', '"coal uses"', '3', 'Coal is primarily used for electricity generation, but it is also used in steel production, cement manufacturing, and other industrial processes. It has been a key source of energy for centuries.'], ['5. Environmental impact', 'dalle', '"coal and the environment"', '4', 'Coal mining and burning have significant environmental impacts, including air and water pollution, habitat destruction, and greenhouse gas emissions. Efforts are being made to reduce these impacts and transition to cleaner energy sources.']]
    themes = ['1. Formation of Coal', 'dalle', '2. Types of Coal', 'network', '3. Uses of Coal', 'dalle', '4. Environmental Impact of Coal', 'network', '5. Future of Coal', 'dalle']
    slides = []
    blank_slide_layout = prs.slide_layouts[6]

    # FIRST SLIDE
    slides.append(prs.slides.add_slide(blank_slide_layout)) #slides[0] 13.333 in x 7.5 in
    firstSlide(slide=slides[0], topic=topic, usersname=usersname)

    # PLAN
    if(plan == True):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        generatePlan(slide=slides[1], themes=themes, language=language)
    

    # Main Slides
    # slides.append(prs.slides.add_slide(blank_slide_layout))
    # translator.translate('Hello', dest=translator.detect('привет мир').lang).text
    
    for i in range(len(themes)):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        if themes[i][3] == 3:
            generate3PicturesSlide(slide=slides[2 + i], methodOfPictures=themes[i][1], PromptForPictures=themes[i][2],
                                    textForSlide=translator.translate(themes[i][5], dest=language).text, topic=translator.translate(themes[i][0], dest=language).text, varitation=1, usersname=usersname, timecode=timecode, 
                                    directory = str(usersname).replace(' ', '') + str(timecode),
                                    prompt2=themes[i][4])
        elif themes[i][3] == 2:
            generate2PicturesSlide(slide=slides[2 + i], methodOfPictures=themes[i][1], PromptForPictures=themes[i][2],
                                    textForSlide=translator.translate(themes[i][5], dest=language).text, topic=translator.translate(themes[i][0], dest=language).text, varitation=1, usersname=usersname, timecode=timecode, 
                                    directory = str(usersname).replace(' ', '') + str(timecode),
                                    prompt2=themes[i][4])
        else:
            generate1PicturesSlide(slide=slides[2 + i], methodOfPictures=themes[i][1], PromptForPictures=themes[i][2],
                                    textForSlide=translator.translate(themes[i][5], dest=language).text, topic=translator.translate(themes[i][0], dest=language).text, varitation=1, usersname=usersname, timecode=timecode, 
                                    directory = str(usersname).replace(' ', '') + str(timecode),
                                    prompt2=themes[i][4])
        print(f"{i + 3}'s slide Done by {themes[i][1]}")


    prs.save(f'{usersname+topic}.pptx')
    slides = []


# GeneratePresentationWith2TypesOfPictures(topic="Лягушки", numOfSlides=7, font="DoraDura", usersname="Suvernev Bogdan", timecode=1)