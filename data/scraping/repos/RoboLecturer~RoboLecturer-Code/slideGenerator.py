#  script for lecture script and slide generation from a recieved table of contents

#pip install python-pptx
#pip install openai 
import openai
import re
from pkg.chat import chat_model
from pptx import Presentation
from pptx.util import Inches

# topic1 = input("Select first topic:")
# topic2 = input("Select second topic:")
# topic3 = input("Select third topic:")
# topics = [topic1, topic2, topic3]

# pres_title = input("What is the title for your presentation: ")
pres_title = "Some Science Topics"

topics = ["Astronomy", "Photosynthesis", "Plant cells", "Animal Cells"]


# generate text 
def generate_text(topic):

    query = f"write a 7 sentence paragraph about {topic} in the style of an excited lecturer and then summarize this paragraph into 5 bullet points with -"

    # get response to query using selelcted model
    response = chat_model.getResponse(query)
    
    return response


def prep_text(GPT_output):
    # Regular expression to extract bullet points
    bullet_point_regex = r'•\s(.*?)\n'

    # Split the text into a paragraph and bullet points
    parts = GPT_output.split("\n\n-")
    paragraph = parts[0]
    bullet_points = parts[1].split("\n-")
    return paragraph, bullet_points

# add title for slide
def add_title(slide, title):
    title_frame = slide.shapes.title
    title_frame.text = title

#prepare bullet points for slide 
def add_bullet_points(slide, bullet_points):
    bullet_point_frame = slide.shapes.placeholders[1].text_frame
    for point in bullet_points:
        p = bullet_point_frame.add_paragraph()
        p.text = point
        p.level = 1

# add title and bullet points to slide will incorporate images later 

def add_text_and_image_to_slide(prs, slide_num, title, bullet_points, notes=""):
    slide_layout = prs.slide_layouts[slide_num]
    slide = prs.slides.add_slide(slide_layout)
    add_title(slide, title)
    add_bullet_points(slide, bullet_points)

    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = notes



# put presentation together 
def generate_presentation(content):
    prs = Presentation()
    add_text_and_image_to_slide(prs,0,pres_title,"")
    for i in content:
        script_and_bullets = generate_text(i)
        para, bullets = prep_text(script_and_bullets)
        add_text_and_image_to_slide(prs,1,i,bullets,para)

    prs.save("grade_school_sci_pres_test10.pptx")


generate_presentation(topics)