import openai
import sys
import os
import json
from pptx import Presentation
import streamlit as st
import io
from PIL import Image as image

# Configure our page.
st.set_page_config(
    page_title="Lean Canvas GPT",
    page_icon="💭",
    layout="centered",
    initial_sidebar_state="expanded",
)
st.title("Lean Canvas GPT")
st.markdown("⚡️ _Create a Lean Canvas presentation for your business idea in seconds_")
st.markdown("The [lean business canvas](https://en.wikipedia.org/wiki/Business_Model_Canvas) is a visual tool for quickly outlining and evaluating key aspects of a business idea, including the problem, solution, value proposition, customer segments, and financials.")
image = image.open('demo_image.png')
st.image(image, use_column_width=True)

# Customize the footer
hide_footer_style = """
<style>
#MainMenu {visibility: hidden;}
footer {visibility: visible;}
footer:after {
    content:' | Made with ❤️ by @upster';
</style>
"""
st.markdown(hide_footer_style, unsafe_allow_html=True)


# Get the GPT-3.5-turbo API key from the OPENAI_API_KEY environment variable
api_key = os.environ.get('OPENAI_API_KEY')

if not api_key:
    st.error("Please set the OPENAI_API_KEY environment variable.")
    sys.exit(1)

# Set up the OpenAI API client
openai.api_key = api_key

# Function to generate text using GPT-3.5-turbo
def generate_text(prompt):
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=2048,
        temperature=0.7,
        n=1,
        stop=None
    )
    return response.choices[0].text.strip()

# Collect the idea_description from the user via streamlit. Wait for user to hit enter.
idea_description = st.text_input("Enter your idea description here:")

def render_page(idea_description):
    # Generate the contents of the Lean Canvas for the idea in JSON format
    lean_canvas_prompt = f"Create a lean canvas for this business idea: {idea_description}\n\nYour response must be in JSON format like: {{ \"Problem\": [\"foo\", \"bar\"] }}\n\nTry to add 3 bullet points in each category.\n\nYour response:\n\nAI-RESPONSE:"
    lean_canvas_text = generate_text(lean_canvas_prompt)

    # Parse the JSON response to get the Lean Canvas data
    lean_canvas = json.loads(lean_canvas_text)

    # Generate the fancy name and tagline for the idea in JSON format
    name_and_tagline_prompt = f"Generate a fancy name and tag line for this idea: {idea_description}\n\nYour response must be in JSON format like: {{ \"name\": \"bla\", \"tagline\": \"foo\" }}\n\nYour response:\n\nAI-RESPONSE:"
    name_and_tagline_text = generate_text(name_and_tagline_prompt)

    # Parse the JSON response to get the fancy name and tagline
    name_and_tagline = json.loads(name_and_tagline_text)
    fancy_name = name_and_tagline["name"]
    tagline = name_and_tagline["tagline"]

    # Create a new PowerPoint presentation
    presentation = Presentation()

    # Add the first slide with the fancy name and tagline
    slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(slide_layout)
    title_placeholder = slide.placeholders[0]
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = fancy_name
    subtitle_placeholder.text = tagline

    # Loop through the Lean Canvas and create a slide for each section
    for section, content in lean_canvas.items():
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title_placeholder = slide.placeholders[0]
        content_placeholder = slide.placeholders[1]
        title_placeholder.text = section
        content_placeholder.text = '\n'.join(content)

    # Save the presentation to an in-memory binary stream
    pptx_io = io.BytesIO()
    presentation.save(pptx_io)

    # Reset the stream's position to the beginning
    pptx_io.seek(0)

    # Create a download button for the generated presentation
    st.download_button(
        label="Download Presentation",
        data=pptx_io,
        file_name=f"{fancy_name}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    st.success(f'PowerPoint presentation "{fancy_name}.pptx" has been created successfully.')
    

if st.button("✨ Generate Presentation"):
    try:
        render_page(idea_description)
    except Exception as e:
        st.error(f"Something went wrong: {e}")
