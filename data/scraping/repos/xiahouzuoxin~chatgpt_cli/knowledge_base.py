import os
import warnings
import hashlib
import concurrent
import multiprocessing
import openai
import tiktoken
from tenacity import retry, wait_random_exponential, stop_after_attempt
from tqdm import tqdm
import mimetypes
import chromadb
import docx2txt
import pptx
import csv
from PyPDF2 import PdfReader

cpu_count = multiprocessing.cpu_count()

# Global variables
tokenizer = tiktoken.get_encoding(
    "cl100k_base"
)  # The encoding scheme to use for tokenization

class Embedding:
    def __init__(self, source='sentence_transformer', model='all-MiniLM-L6-v2') -> None:
        self.param_source = source
        self.param_model = model
        if source == 'sentence_transformer':
            from sentence_transformers import SentenceTransformer
            self.model = SentenceTransformer(model)
            self._encode = lambda text: list(self.model.encode(text))
        elif source == 'openai':
            # "text-embedding-ada-002"
            self.model = model
            self._encode = self._get_gpt_embedding

    @retry(wait=wait_random_exponential(min=1, max=20), stop=stop_after_attempt(6))
    def _get_gpt_embedding(self, text: str) -> list[float]:
        return openai.Embedding.create(input=[text], model=self.model)["data"][0]["embedding"]
    
    def encode(self, text: str) -> list[float]:
        if not isinstance(text, str) or not text.strip():
            raise ValueError(f"Embedding input text must be a non-empty string: {text}")
        result = self._encode(text)
        if self.param_source == 'sentence_transformer':
            # convert np.float to float
            result = [v.item() for v in result]
        return result

    def get_embedding(self, texts: list):
        return [self.encode(text) for text in texts]
    
    def get_embedding_parallel(self, texts: str, num_workers=cpu_count):
        with concurrent.futures.ThreadPoolExecutor(max_workers=num_workers) as executor:
            results = list(
                tqdm(
                    executor.map(
                        lambda text: self.encode(text), 
                        texts
                    ), 
                    total=len(texts)
                )
            )
        return results
    
    def cal_gpt_embedding_price(df, text_col='text', base_price_1k=0.0004):
        from transformers import GPT2TokenizerFast
        tokenizer = GPT2TokenizerFast.from_pretrained("gpt2")
        df['n_tokens'] = df[text_col].apply(lambda x: len(tokenizer.encode(x)))
        total_tokens = df['n_tokens'].sum()
        return base_price_1k * (total_tokens / 1000), total_tokens

class VectorRetrieval():
    def __init__(self, workspace) -> None:
        if not os.path.exists(workspace):
            os.mkdir(workspace)

        self.chroma_client = chromadb.PersistentClient(path=workspace)
        self.collection = self.chroma_client.get_or_create_collection(name='knowledge_base')

        # embedding model
        # Reference: https://www.sbert.net/docs/pretrained_models.html
        self.emb_model_config = {
            'source': 'sentence_transformer',
            'model': 'all-MiniLM-L6-v2',
            'embedding_dim': 384,
            'max_seq_len': 256
        }
        self.emb_model = Embedding(source=self.emb_model_config['source'], model=self.emb_model_config['model'])

    def add_index_for_texts(self, texts: list, num_workers=cpu_count):
        texts = list(set(texts))
        num_workers = min(num_workers, len(texts))
        text_embs = self.emb_model.get_embedding_parallel(texts, num_workers=num_workers)
        ids = [hashlib.sha1(_text.encode("utf-8")).hexdigest() for _text in texts]
        self.collection.add(documents=texts, embeddings=text_embs, ids=ids)

    def add_index_for_docs(self, path: str, num_workers=5):
        docs = []
        for root, dirs, files in os.walk(path):
            for file in files:
                docs.append( os.path.join(root, file) )
        if len(docs) == 0:
            return 0
        num_workers = min(num_workers, len(docs))
        with concurrent.futures.ThreadPoolExecutor(max_workers=num_workers) as executor:
            results = list(
                tqdm(
                    executor.map(
                        lambda doc: self.extract_chunk_texts_from_file(doc), 
                        docs
                    ), 
                    total=len(docs)
                )
            )
        texts = [_text for _texts in results for _text in _texts]
        self.add_index_for_texts(texts)
        return len(texts)

    # Extract text from a file based on its mimetype
    def extract_chunk_texts_from_file(self, file):
        """Return the text content of a file."""
        file_mimetype = mimetypes.guess_type(file)[0]
        if file_mimetype == "application/pdf":
            # Extract text from pdf using PyPDF2
            reader = PdfReader(file)
            extracted_text = " ".join([page.extract_text() for page in reader.pages])
        elif file_mimetype == "text/plain" or file_mimetype == "text/markdown":
            # Read text from plain text file
            extracted_text = file.read().decode("utf-8")
            file.close()
        elif file_mimetype == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            # Extract text from docx using docx2txt
            extracted_text = docx2txt.process(file)
        elif file_mimetype == "text/csv":
            # Extract text from csv using csv module
            extracted_text = ""
            decoded_buffer = (line.decode("utf-8") for line in file)
            reader = csv.reader(decoded_buffer)
            for row in reader:
                extracted_text += " ".join(row) + "\n"
        elif file_mimetype == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            # Extract text from pptx using python-pptx
            extracted_text = ""
            presentation = pptx.Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                extracted_text += run.text + " "
                        extracted_text += "\n"
        else:
            # Unsupported file type
            extracted_text = None
            warnings.warn("Unsupported file type of {}: {}".format(file, file_mimetype))
        if extracted_text is None:
            return []

        text_chunks = self.get_text_chunks(extracted_text, chunk_size=int(self.emb_model_config['max_seq_len'] * 0.8))
        return text_chunks
    
    def get_text_chunks(self, text, chunk_size, min_chunk_size_chars=350, min_chunk_length_to_embed=5):
        '''
        Split a text into chunks of ~chunk_size length, based on punctuation and newline boundaries.
        Reference: https://github.com/openai/chatgpt-retrieval-plugin/blob/main/services/chunks.py
        '''
        if not text or text.isspace():
            return []
        
        # Tokenize the text
        tokens = tokenizer.encode(text, disallowed_special=())

        chunks = []
        while tokens:
            chunk = tokens[:chunk_size]
            chunk_text = tokenizer.decode(chunk)
            if not chunk_text or chunk_text.isspace():
                tokens = tokens[len(chunk):]
                continue

            # Find the last period or punctuation mark in the chunk
            last_punctuation = max(
                chunk_text.rfind("."),
                chunk_text.rfind("?"),
                chunk_text.rfind("!"),
                chunk_text.rfind("\n"),
            )

            # If there is a punctuation mark, and the last punctuation index is before MIN_CHUNK_SIZE_CHARS
            if last_punctuation != -1 and last_punctuation > min_chunk_size_chars:
                # Truncate the chunk text at the punctuation mark
                chunk_text = chunk_text[: last_punctuation + 1]

            # Remove any newline characters and strip any leading or trailing whitespace
            chunk_text_to_append = chunk_text.replace("\n", " ").strip()

            if len(chunk_text_to_append) > min_chunk_length_to_embed:
                # Append the chunk text to the list of chunks
                chunks.append(chunk_text_to_append)

            # Remove the tokens corresponding to the chunk text from the remaining tokens
            tokens = tokens[len(tokenizer.encode(chunk_text, disallowed_special=())) :]

        # Handle the remaining tokens
        if tokens:
            remaining_text = tokenizer.decode(tokens).replace("\n", " ").strip()
            if len(remaining_text) > min_chunk_length_to_embed:
                chunks.append(remaining_text)

        return chunks
    
    def query(self, texts: str, limit=3):
        query_embs = self.emb_model.get_embedding_parallel(texts, num_workers=min(len(texts), cpu_count))
        return self.collection.query(query_embeddings=query_embs, n_results=limit)

if __name__ == '__main__':
    vector_retrieval = VectorRetrieval('./knowledge/vector_db/')
    vector_retrieval.add_index_for_docs(path='./knowledge/docs')
    print('Generate knowledge base down.')
    
    print( vector_retrieval.query(['How are you'], limit=1) )