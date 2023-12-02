from django.shortcuts import render, HttpResponse
from django.http import JsonResponse

import os
import sys
from langchain.indexes import VectorstoreIndexCreator
from langchain.vectorstores import Chroma
from langchain.embeddings import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.llms import OpenAI
from langchain.chains import RetrievalQA
from langchain.document_loaders import DirectoryLoader
import chromadb
from chromadb.utils import embedding_functions
import io
import uuid
import datetime


# Import Loaders
from langchain.document_loaders import PyPDFLoader
from langchain.document_loaders import Docx2txtLoader
from langchain.document_loaders import TextLoader
from langchain.document_loaders.csv_loader import CSVLoader
from langchain.document_loaders import UnstructuredPowerPointLoader

# Import File Loaders
import PyPDF2
from docx import Document
from pptx import Presentation
import os
import csv

os.environ["OPENAI_API_KEY"] = ""


# Sets up environ variables for openai
os.environ["TOKENIZERS_PARALLELISM"] = "false"
os.environ["OPENAI_API_KEY"] = "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"

# Creates a new embedding function (OpenAI Embedding Function)
openai_ef = embedding_functions.OpenAIEmbeddingFunction(
    model_name="text-embedding-ada-002",
    api_key="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",
)

# Gets the chroma client
chroma_client = chromadb.HttpClient(host="localhost", port=8000)


# Generates a random ID
def uid():
    return str(uuid.uuid4())


# Deletes a collection in the databse
def delete_collection(collection):
    chroma_client.delete_collection(name=collection)


# Creates a collection in the databse
def create_collection(collection):
    chroma_client.create_collection(name=collection, embedding_function=openai_ef)


# Gets a collection in the databse
def get_collection(collection_name):
    return chroma_client.get_collection(name=collection_name)


def query_to_file(text):
    pass


# Uploads a file to the collection
def upload_file(file_name, collection):
    text = ""
    ext = get_extension(file_name)
    if ext == ".docx":
        text = load_docx(file_name)
    elif ext == ".pdf":
        text = load_pdf(file_name)
    elif ext == ".csv":
        text = load_csv(file_name)
    elif ext == ".pptx":
        text = load_ppt(file_name)
    elif ext == ".txt":
        text = load_txt(file_name)

    collection.add(documents=[text], metadatas=[{"source": file_name}], ids=[uid()])


# Does a similarity search between the prompt and the data in the collection
def query_chroma(collection, prompt):
    query = collection.query(
        query_texts=[
            prompt,
        ],
        n_results=1,
    )

    return query


def get_all_collections():
    return chroma_client.list_collections()


# results = query_chroma(get_collection("test"), "What is the most common rating?")

# a = results["documents"][0]


# with open("temp.csv", "w") as f:
#     f.write(a[0] + "\n")

# a = pd.read_csv("temp.csv")

# agent = create_pandas_dataframe_agent(
#     OpenAI(temperature=0, max_tokens=100), a, verbose=True
# )

# print(agent.run("What is the most common opening?"))

# query = sys.argv[1]
# print(query)
# loader = TextLoader("temp.csv")
# index = VectorstoreIndexCreator().from_loaders([loader])
# print(index.query(query))


def load_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page in range(len(pdf_reader.pages)):
        text += pdf_reader.pages[page].extract_text()
    return text


def load_docx(file):
    doc = Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text
    return text


def load_csv(file):
    csv_data = file.read().decode("utf-8")
    csv_reader = csv.reader(csv_data.splitlines())
    text = ""
    for row in csv_reader:
        text += ", ".join(row) + "\n"  # Convert CSV rows to a string
    return text


def load_ppt(file):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text


def load_txt(file):
    with open(file, "r") as f:
        return f.read()


def get_extension(file):
    return os.path.splitext(file)[1]


# Create your views here.
def index(request):
    return render(request, "index.html")


def store_files(files):
    current_datetime = datetime.datetime.now()
    directory = os.path.join("static", "original_files")
    for file in files:
        file_name = current_datetime.__str__() + file.name
        file_name = file_name.strip(" ")
        file_path = os.path.join(directory, file_name)
        # Save the file to the directory.
        with open(file_path, "wb") as f:
            f.write(file.read())

        upload_file(file_path, get_collection("test"))


def upload(request):
    loaders = []
    if request.method == "POST" and request.FILES:
        uploaded_files = request.FILES.getlist("files")
        store_files(uploaded_files)

        return HttpResponse("Uploaded Files")

    elif request.method == "POST":
        print("Hello")
        data = request.POST.get("data")
        query = query_chroma(get_collection("test"), data)
        print(query["documents"][0][0])
        with open("static/context/context.txt", "w") as f:
            f.write(query["documents"][0][0] + "\n")

        loader = TextLoader("static/context/context.txt")
        index = VectorstoreIndexCreator().from_loaders([loader])
        result = index.query(data)

        return JsonResponse({"result": result})

        # for loader in loaders:
        #     index.from_loaders([loader])
        #     print(1)

        # print(index)
        # index.

        # print(index)

        # additional_data = request.POST.get("additional_data", "")
        # print(additional_data)
        # print(uploaded_files[0].name)
        # print(load_docx(uploaded_files[0]))
        # loader = PyPDFLoader(uploaded_files[0])
        # pages = loader.load_and_split()
