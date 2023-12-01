import openai
from pptx import Presentation

from decker.ErrorHandling import ErrorHandling
from decker.Validation import Validation
from decker.Logging import Logging
from decker.UnitTesting import TestComponent


class PitchDeckGenerator:
    def __init__(self):
        self.presentation = Presentation()
        self.error_handling = ErrorHandling()
        self.validation = Validation()
        self.logging = Logging()
        self.test_component = TestComponent()

    def generate_content(self, prompt):
        response = openai.Completion.create(engine="text-davinci-003", prompt=prompt, max_tokens=100)
        return response.choices[0].text.strip()

    def create_slide(self, title, content):
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title_placeholder.text = title
        content_placeholder.text = content

    def create_pitch_deck(self, brand, xcom_links, crunchbase_links):
        slides_content = [
            ("Introduction", f"Generate an introduction for a pitch deck for {brand}.", xcom_links),
            ("Team Overview", f"Generate a team overview for a pitch deck for {brand}.", xcom_links),
            ("Problem Statement", f"Generate a problem statement for a pitch deck for {brand}.", xcom_links),
            ("Solution", f"Generate a description of the solution for a pitch deck for {brand}.", xcom_links),
            ("Market Size", f"Generate a description of the market size for a pitch deck for {brand}.", xcom_links),
            ("Business Model", f"Generate a description of the business model for a pitch deck for {brand}.", xcom_links),
            ("Marketing Strategy", f"Generate a description of the marketing strategy for a pitch deck for {brand}.", xcom_links),
            ("Financial Projections", f"Generate financial projections for a pitch deck for {brand}.", crunchbase_links),
            ("Conclusion", f"Generate a conclusion for a pitch deck for {brand}.", xcom_links)
        ]

        for title, prompt, links in slides_content:
            content = self.generate_content(prompt)
            self.create_slide(title, content)

        self.presentation.save(f'{brand}_pitch_deck.pptx')