# Sidekick365 is a Python module that provides functions to create a PowerPoint presentation from a Word document.
# Steve Goodman 2023/09/27
# The module uses the following libraries:
import os
import sys
import requests
import json
import base64
import pptx
import docx
import openai
import random
import pywinauto


# function wrapping Word Doc conversion to AI generated PowerPoint with images
def GeneratePowerPointFromWord(wordfile, powerpointfile, customphrase=""):
    markdown = ExtractDoc(wordfile)
    summary = GenPresentation(markdown, customphrase)
    CreatePPTX(summary, powerpointfile)
    # count lines in summary with # in them and store as slidecount
    slidecount = 0
    for line in summary.splitlines():
        if line.startswith("#"):
            slidecount += 1
    return slidecount

# function to open PowerPoint and apply the designer
def OpenPowerPointAndApplyDesigner(powerpointfile,slidecount):
    # get full path of file provided 
    filepath = os.path.abspath(powerpointfile)
    # use pywinauto to start powerpoint with the powerpointfile
    app = pywinauto.application.Application(backend="uia").start(f"C:\\Program Files\\Microsoft Office\\root\\Office16\\POWERPNT.EXE \"{filepath}\"")
    # wait for the powerpoint window to be ready
    window = app.window(title_re=".*PowerPoint")
    window.wait('ready', timeout=30)
    # find the thumbnails window
    # find the designer button
    designer=window.Designer
    # click the designer button
    designer.click_input()
    # wait for the designerListBox window to be ready
    window.DesignerListBox.wait('ready', timeout=30)
    # find the designerListBox
    designerListBox = window.DesignerListBox
    # apply to the first slide
    designerListBox.ListItem1.click_input()
    # iterate through thumbnails_list
    # iterate through slidecount
    for i in range(2,slidecount+1):
        # click the main window titlebar
        window.click_input()
        # send page down to move to the next slide
        pywinauto.keyboard.send_keys("{PGDN}")
        # find the first listitem in the designerListBox and click it
        designerListBox.ListItem1.click_input()
    # move to the first slide
    window.click_input()
    pywinauto.keyboard.send_keys("{HOME}")
    return True  

# create a function openai_summarise with text as input and return summary
def GenPresentation(doctext,customphrase=""):

    openai.api_type = "azure"
    openai.api_key = os.environ["OPENAI_API_KEY"]
    openai.api_base = os.environ["OPENAI_API_BASE"]
    openai.api_version = "2023-07-01-preview"
    # create the prompt
    conversation= [{"role":"system","content":"You are an AI assistant that will expect a markdown format document to be provided and respond using markdown format that captures the points of the document and will be suitable for creating a slide presentation for a professional audience. The first heading will be a suitable title based on the content. The subsequent headings will have between 3 and 10 bullet points, then the phrase \"image prompt:\" followed by a suitable dall-e prompt, then a suitable color for the slide background. Optionally the user may add a customization phrase to guide the output. This will be a line prefixed by SideKick365_custom:"},{"role":"user","content":"#The Long Walk of Bernard Jenkins\n##Chapter One: The Departure\nBernard Jenkins had always dreamed of seeing the world. He was born in a small town in Queensland, Australia, where he lived with his parents and his younger sister. He loved reading books and magazines about different countries and cultures, and he imagined himself as an explorer, a traveler, a adventurer. He wanted to go beyond the horizon, to discover new places and people, to have amazing experiences and stories to tell.\nBut his family was poor, and they could not afford to travel. His father worked as a mechanic, and his mother was a seamstress. They barely had enough money to pay the bills and buy food. They lived in a modest house, with no TV, no phone, no car. Bernard had to walk to school every day, and he often wore the same clothes for weeks. He felt trapped in his small world, and he longed for something more.\nWhen he was 18, he decided to make his dream come true. He had saved up some money from doing odd jobs around the town, and he bought a backpack, a tent, a sleeping bag, a map, and a compass. He told his parents and his sister that he was going to walk to the UK, the land of his ancestors, and that he would be back in three years. They thought he was crazy, and they tried to stop him. They said it was too dangerous, too expensive, too impossible. They said he would never make it, that he would die on the way, that he would regret it for the rest of his life. They begged him to stay, to find a job, to settle down, to be realistic.\nBut Bernard was determined. He had a vision, a goal, a passion. He wanted to prove himself, to challenge himself, to fulfill himself. He wanted to live his life on his own terms, to follow his heart, to chase his dream. He hugged his family goodbye, and he set off on his journey. He walked out of his town, out of his state, out of his country. He walked across deserts, mountains, jungles, rivers. He walked through heat, cold, rain, snow. He walked alone, with no one to help him, to guide him, to support him. He walked with courage, with perseverance, with faith. He walked for three years, until he reached the UK.\n##Chapter Two: The Arrival\nBernard Jenkins arrived in the UK on a cold and cloudy day in November. He had walked over 20,000 kilometers, crossing 18 countries, facing countless dangers, hardships, and obstacles. He had seen wonders and horrors, beauty and ugliness, kindness and cruelty. He had met people of all kinds, some friendly, some hostile, some indifferent. He had learned languages, customs, histories, and traditions. He had grown physically, mentally, emotionally, and spiritually. He had changed, he had evolved, he had transformed.\nHe was no longer the naive and restless boy who had left Australia. He was now a mature and confident man who had seen the world. He had a wealth of knowledge, a depth of wisdom, a breadth of perspective. He had a sense of accomplishment, a feeling of gratitude, a vision of purpose. He had a story to tell, a story to share, a story to inspire.\nHe walked into London, the capital of the UK, and he felt a surge of excitement and curiosity. He wanted to explore the city, to see the sights, to experience the culture. He wanted to find a place to stay, to get a job, to make a living. He wanted to start a new chapter in his life, to create a new identity, to build a new future.\nHe walked into a pub, and he asked for a pint of beer. He sat at the bar, and he looked around. He saw people of different ages, races, genders, and backgrounds. He saw people talking, laughing, arguing, flirting. He saw people living, working, loving, dreaming. He felt a connection, a belonging, a home.\nHe started a conversation with the bartender, a friendly and talkative man named Tom. He told him his name, his origin, his story. Tom was amazed, impressed, and curious. He asked him questions, he gave him compliments, he offered him advice. He told him about the city, the country, the people. He told him about the opportunities, the challenges, the possibilities. He told him about his own life, his own dreams, his own struggles.\nTom liked Bernard, and he wanted to help him. He told him that he knew a place where he could stay, a friend who had a spare room. He told him that he knew a job where he could work, a company that was hiring. He told him that he knew a business where he could invest, a startup that was growing. He told him that he knew a woman who he could date, a friend who was single. He told him that he knew a lot of things, and that he was willing to share them with him.\nBernard was grateful, touched, and intrigued. He accepted Tom's offer, he thanked him for his kindness, he agreed to his plan. He finished his beer, he paid his bill, he followed his lead. He walked out of the pub, and he walked into his new life.\n"},{"role":"assistant","content":"# The Long Walk of Bernard Jenkins - His Story\n# The Long Walk of Bernard Jenkins\n* Bernard Jenkins had always dreamed of seeing the world.\n* He was born in a small town in Queensland, Australia, where he lived with his parents and his younger sister.\n* He loved reading books and magazines about different countries and cultures, and he imagined himself as an explorer, a traveler, a adventurer.\n* He wanted to go beyond the horizon, to discover new places and people, to have amazing experiences and stories to tell.\n* But his family was poor, and they could not afford to travel.\n* His father worked as a mechanic, and his mother was a seamstress.\n* They barely had enough money to pay the bills and buy food.\n* They lived in a modest house, with no TV, no phone, no car.\nimage prompt: a traveller looking beyond the horizon\n# A Journey to the UK\n* When he was 18, he decided to make his dream come true.\n* He saved up money from doing odd jobs around the town\n* He bought a backpack, a tent, a sleeping bag, a map, and a compass\n* He told his parents and his sister that he was going to walk to the UK\n* He would be back in three years\n* His family thought he was crazy and tried to stop him\n* They said it was too dangerous, too expensive, too impossible\n* They said he would never make it, that he would die on the way, that he would regret it for the rest of his life\n* They begged him to stay, to find a job, to settle down, to be realistic\nimage prompt: a man arriving in the UK looking towards London\n# Following His Dreams: Bernard's Journey\n* Bernard was determined to prove himself and follow his heart\n* He set off on a journey to the UK, walking through deserts, mountains, jungles, and rivers\n* He faced heat, cold, rain, and snow, and walked alone with no help or guidance\n* He persevered with courage and faith, and after three years, he reached the UK\nimage prompt: a man working hard, vintage look\n# Bernard Jenkins' Journey of Transformation\n* Bernard Jenkins embarked on a 20,000 kilometer journey across 18 countries\nHe faced countless dangers, hardships, and obstacles\n* He encountered people of all kinds, and learned languages, customs, histories, and traditions\n* He emerged a mature and confident man, with a wealth of knowledge, a depth of wisdom, and a breadth of perspective\n* He had a sense of accomplishment, a feeling of gratitude, and a vision of purpose\n* He had a story to tell, a story to share, and a story to inspire\nimage prompt: an inspiring photo of a man looking ahead with a variety of people from different backgrounds alongside him\n# Exploring London and Finding a Home\n* He walked into London, the capital of the UK, and felt a surge of excitement and curiosity\n* He wanted to explore the city, to see the sights, to experience the culture\n* He wanted to find a place to stay, to get a job, to make a living\n* He wanted to start a new chapter in his life, to create a new identity, to build a new future\nHe walked into a pub and asked for a pint of beer\n* He saw people of different ages, races, genders, and backgrounds\n* He saw people talking, laughing, arguing, flirting\n* He saw people living, working, loving, dreaming\nHe started a conversation with the bartender, a friendly and talkative man named Tom\n* Tom asked him questions, gave him compliments, and offered him advice\n* He told him about the city, the country, the people\nimage prompt: a world traveller with a backpack, talking to a person in a deep discussion, in a pub\n# Tom Helps Bernard Start a New Life\n* Tom liked Bernard and wanted to help him\n* He offered him a place to stay\n* He offered him a job\n* Bernard accepted Tom's offer and thanked him for his kindness\nHe followed Tom's lead and walked into his new life\nimage prompt: two people shaking hands and smiling with an inspirational business background"}]
    # customize the prompt if a customization_phrase was provided
    if customphrase != "":
        doctext = doctext + "\nSideKick365_custom:" + customphrase
    conversation.append({"role": "user", "content": doctext})

    # create the chat completion
    response = openai.ChatCompletion.create(
        engine="gpt-4-32k",
        messages = conversation,
        temperature=0.5,
        max_tokens=3500,
        top_p=0.95,
        frequency_penalty=0,
        presence_penalty=0,
        stop=None)
    # use doctext to create a query using ChatCompletion and store the response as summary
    summary = response['choices'][0]['message']['content']
    # return summary
    return summary

# create function create_image with prompt as input and return the filename to the image
def GenImage(imageprompt):
    # set parameters
    openai.api_type = "azure"
    openai.api_key = os.environ["OPENAI_API_KEY"]
    openai.api_base = os.environ["OPENAI_API_BASE"]
    openai.api_version = "2023-06-01-preview"

    response = openai.Image.create(
        prompt=imageprompt,
        size='1024x1024',
        n=1
    )
    image_url = response["data"][0]["url"]
    # get image from url
    image = requests.get(image_url)
    # write image to random temporary file
    # get filesystem temporary folder
    tmpdir = os.environ['TEMP']
    # set a filename to tmp/tmp_ plus the prompt with spaces and commas replaced by underscores and with the .png extension
    filename = tmpdir + "/tmp_" + imageprompt.replace(" ", "_").replace(",", "_") + ".png"
    # write image to file
    with open(filename, "wb") as file:
        file.write(image.content)
    # return base64 image
    return filename

# function to create a pptx from a markdown file that uses the first heading as the title slide and the subsequent headings as slide titles. 
# Each subsequent slide has a title section and two columns.
# Each subsequent heading has bullet points that should be inserted, a color indicated, and a description of an image to be generated using openai_dalle(prompt) and inserted on the slide from the base64 image returned. The bullet points and image should be inserted into the two columns, with the image in the right column.
# The pptx is saved to the filename provided.
def CreatePPTX(summary, filename):
    # process the markdown format summary into a dictionary with the first heading as the title, and then each subsequent heading with dictionary entries for the bullet points, color, and image prompt
    # create a dictionary to hold the summary
    summary_dict = {}
    # split the summary into lines
    summary_lines = summary.splitlines()
    # set a slide number variable
    slidenumber = 0
    # loop through the lines in summary_lines
    for line in summary_lines:
        # if the line starts with a # then it is a heading
        if line.startswith("#"):
            # increment the slide number
            slidenumber += 1
            # remove the # from the line
            line = line.replace("#", "")
            # trim the line
            line = line.strip()
            # set the heading as the key in the dictionary
            heading = line
            # set the slidenumber as a dictionary
            summary_dict[slidenumber] = {}
            # set the heading as a string
            summary_dict[slidenumber]["heading"] = heading
            # set the bullet points as a list
            summary_dict[slidenumber]["bullet_points"] = []
            # set the color as a string
            summary_dict[slidenumber]["color"] = ""
            # set the image prompt as a string
            summary_dict[slidenumber]["image_prompt"] = ""
            # set the base64 encoded image as a string
            summary_dict[slidenumber]["image"] = ""
        # if the line starts with a * then it is a bullet point
        elif line.startswith("*"):
            # remove the * from the line
            line = line.replace("*", "")
            # trim the line
            line = line.strip()
            # add the bullet point to the list
            summary_dict[slidenumber]["bullet_points"].append(line)
        # if the line starts with a image prompt then it is a image prompt
        elif line.startswith("image prompt:"):
            # remove image prompt: from the line
            line = line.replace("image prompt:", "")
            # trim the line
            line = line.strip()
            # set the image prompt
            summary_dict[slidenumber]["image_prompt"] = "a photo of " + line + ", photoreastic, corporate stock photo"
            # create the image
            imagefilename = GenImage(line)
            # set image to the image filename
            summary_dict[slidenumber]["image"] = imagefilename

    # create a new PowerPoint document
    document = pptx.Presentation()
    # use summary_dict to create the slides
    # start with summary_dict[1] as the title slide
    # create the title slide
    slide = document.slides.add_slide(document.slide_layouts[0])
    # set the title
    title = slide.shapes.title
    # set the title text
    title.text = summary_dict[1]["heading"]
    # loop through the remaining slides in summary_dict from 1 to the end
    for slidenumber in range(2, len(summary_dict)):
        # randomly choose between layout 1 and layout 3 but not 2
        layout = random.choice([1, 3])
        if layout == 2:
            layout = 3
        
        # add a new slide with a title and two columns
        slide = document.slides.add_slide(document.slide_layouts[layout])
        # set the title
        title = slide.shapes.title
        # set the title text
        title.text = summary_dict[slidenumber]["heading"]
        # add the bullet points to the left column
        # set the left column
        left_column = slide.shapes.placeholders[1]
        # set the left column text frame
        left_column_text_frame = left_column.text_frame
        
        # loop through the bullet points and add
        for bullet_point in summary_dict[slidenumber]["bullet_points"]:
            left_column_text_frame.add_paragraph().text = bullet_point
        # if count of summary_dict[slidenumber]["bullet_points"] is more than 3 then loop through each paragraph in left_column_text_frame and set the font size to 14pt
        if len(summary_dict[slidenumber]["bullet_points"]) > 3:
            for paragraph in left_column_text_frame.paragraphs:
                paragraph.font.size = pptx.util.Pt(14)


        if layout == 3:
            # set the right column
            right_column = slide.shapes.placeholders[2]
            imagefilename = summary_dict[slidenumber]["image"]
            # add a the image as right_column_picture the same dimensions and location as the right column 
            right_column_picture = slide.shapes.add_picture(imagefilename, right_column.left, right_column.top, right_column.width, right_column.height)
            # remove the right_column placeholder
            right_column.element.getparent().remove(right_column.element)


        

    # add a blank slide to the end
    slide = document.slides.add_slide(document.slide_layouts[6])

    # save the PowerPoint document
    document.save(filename)

# function to extract a word document using python-docx and return it in markdown text
def ExtractDoc(filename):
    # open document
    document = docx.Document(filename)
    # long string to hold document text
    doc_text = ""
    # loop through document paragraphs
    for para in document.paragraphs:
        # format headings as md headings
        if para.style.name == "Heading 1":
            doc_text += "# " + para.text
            doc_text += "\n"
        elif para.style.name == "Heading 2":
            doc_text += "## " + para.text
            doc_text += "\n"
        elif para.style.name == "Heading 3":
            doc_text += "### " + para.text
            doc_text += "\n"
        elif para.style.name == "Heading 4":
            doc_text += "#### " + para.text
            doc_text += "\n"
        elif para.style.name == "Heading 5":
            doc_text += "##### " + para.text
            doc_text += "\n"
        elif para.style.name == "Heading 6":
            doc_text += "###### " + para.text
            doc_text += "\n"
        # format bold text as md bold
        elif para.style.name == "Strong":
            doc_text += "**" + para.text + "**"
            doc_text += "\n"
        # format italic text as md italic
        elif para.style.name == "Emphasis":
            doc_text += "*" + para.text + "*"
            doc_text += "\n"
        # format hyperlinks as md hyperlinks
        elif para.style.name == "Hyperlink":
            doc_text += "[" + para.text + "](" + para.hyperlink.address + ")"
            doc_text += "\n"
        # format bullet points as md bullet points
        elif para.style.name == "List Bullet":
            doc_text += "* " + para.text
            doc_text += "\n"
        # format numbered lists as md numbered lists
        elif para.style.name == "List Number":
            doc_text += "1. " + para.text
            doc_text += "\n"
        # format block quotes as md block quotes
        elif para.style.name == "Quote":
            doc_text += "> " + para.text
            doc_text += "\n"
        # format code as md code
        elif para.style.name == "Code":
            doc_text += "`" + para.text + "`"
            doc_text += "\n"
        # format title as md title
        elif para.style.name == "Title":
            doc_text += "# " + para.text
            doc_text += "\n"
        # add paragraph text to doc_text
        doc_text += para.text
        # add new line
        doc_text += "\n"
    # return doc_text
    return doc_text
    